#!/usr/bin/env python3
"""
Append diagnostic slides to existing report PPTX.
Adds slides about: Round1 floater analysis, Z-layer analysis,
filtering stages, ellipsoidal plateau design, anchor generation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

PPTX = "/home/wosas/Desktop/Incremental_mapping_test/gs_floaterLab/context/ppt/report_20260705.pptx"
IMGS = Path("/home/wosas/Desktop/Incremental_mapping_test/gs_floaterLab/context/ppt/imgs")

C_BG     = RGBColor(0x0F, 0x17, 0x2A)
C_ACCENT = RGBColor(0x38, 0xBD, 0xF8)
C_ACCENT2= RGBColor(0xFB, 0xBF, 0x24)
C_GREEN  = RGBColor(0x34, 0xD3, 0x99)
C_RED    = RGBColor(0xF8, 0x71, 0x71)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
C_PANEL  = RGBColor(0x1E, 0x2D, 0x45)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation(PPTX)

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BG

def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=13, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color

def title_bar(slide, title, sub=None):
    add_rect(slide, 0, 0, W, Inches(0.07), C_ACCENT)
    txt(slide, title, Inches(0.5), Inches(0.18), Inches(12), Inches(0.7),
        size=30, bold=True)
    if sub:
        txt(slide, sub, Inches(0.5), Inches(0.88), Inches(12), Inches(0.4),
            size=15, color=C_ACCENT)

def img(slide, path, l, t, w, h=None):
    if not Path(path).exists():
        return
    if h:
        slide.shapes.add_picture(str(path), l, t, w, h)
    else:
        slide.shapes.add_picture(str(path), l, t, w)

def bullet_box(slide, items, l, t, w, h, title=None, tc=C_ACCENT):
    add_rect(slide, l, t, w, h, C_PANEL)
    add_rect(slide, l, t, Inches(0.04), h, tc)
    y = t + Inches(0.15)
    if title:
        txt(slide, title, l+Inches(0.12), y, w-Inches(0.2), Inches(0.33),
            size=13, bold=True, color=tc)
        y += Inches(0.36)
    for item in items:
        txt(slide, f"• {item}", l+Inches(0.12), y, w-Inches(0.2), Inches(0.33),
            size=12, color=C_WHITE)
        y += Inches(0.32)

def table(slide, headers, rows, l, t, w, hdr=C_ACCENT):
    nc = len(headers); cw = w/nc; rh = Inches(0.38)
    for ci, h in enumerate(headers):
        add_rect(slide, l+cw*ci, t, cw, rh, hdr)
        txt(slide, h, l+cw*ci+Inches(0.05), t+Inches(0.05),
            cw-Inches(0.1), rh-Inches(0.05),
            size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        bg = C_PANEL if ri%2==0 else C_BG
        for ci, cell in enumerate(row):
            add_rect(slide, l+cw*ci, t+rh*(ri+1), cw, rh, bg)
            tc = C_RED if (cell.startswith("-") and "dB" in cell) \
                 else C_GREEN if (cell.startswith("+") and "dB" in cell) \
                 else C_WHITE
            txt(slide, cell,
                l+cw*ci+Inches(0.05), t+rh*(ri+1)+Inches(0.05),
                cw-Inches(0.1), rh-Inches(0.05),
                size=11, color=tc, align=PP_ALIGN.CENTER)
    return t + rh*(len(rows)+1)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide A — Round 1: Floater 실태 분석
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs); fill_bg(sl)
title_bar(sl, "Round 1 진단: Floater 실태 분석 (exp08 기준)",
          "exp08 30k checkpoint 323,864 Gaussians 전수 분석")

img(sl, IMGS/"round1_overview.png",
    Inches(0.3), Inches(1.35), Inches(5.5))

bullet_box(sl, [
    "총 323,864 Gaussians 분석",
    "Low-opacity (<0.1): 39.0% = 126,360개",
    "NN 거리 중앙값: 0.025m (대부분 표면 근처)",
    ">1m 거리: 0.1% (227개) — 진짜 floater",
    "depth_residual<-0.2m (빈 공간 침범): 72개 (0.0%)",
    "→ opacity로만 floater 찾으면 오탐 많음",
    "→ 실제 floater(빈 공간)는 극소수",
], Inches(6.0), Inches(1.35), Inches(4.2), Inches(3.5),
   title="Round 1 수치 요약")

bullet_box(sl, [
    "Low-opacity AND depth<-0.2m 동시: 30개 (진짜 floater)",
    "depth<-0.2m만: 42개 (opacity 필터로 못 잡는 floater!)",
    "Low-opacity만: 126,330개 (표면 위 정상 Gaussian)",
    "→ opacity 기반 pruning은 표면 Gaussian 과다 제거 위험",
], Inches(10.4), Inches(1.35), Inches(2.7), Inches(3.5),
   title="핵심 결론", tc=C_ACCENT2)

img(sl, IMGS/"round1c_z_outlier.png",
    Inches(0.3), Inches(5.0), Inches(12.7), Inches(2.2))


# ═══════════════════════════════════════════════════════════════════════════════
# Slide B — Z-Layer 분석 (MPS vs ORB)
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs); fill_bg(sl)
title_bar(sl, "Z-Layer 분석: MPS vs ORB-SLAM 앵커 분포",
          "높이(Z)별 포인트 밀도 & coverage — Pop2 구간 식별")

txt(sl, "MPS Z-Layer (semidense, 626k pts)",
    Inches(0.3), Inches(1.3), Inches(6.3), Inches(0.35),
    size=14, bold=True, color=C_ACCENT)
img(sl, IMGS/"mps_zlayer.png",
    Inches(0.3), Inches(1.7), Inches(6.3), Inches(5.5))

txt(sl, "ORB-SLAM Z-Layer (sparse, 6,492 filtered pts)",
    Inches(6.8), Inches(1.3), Inches(6.3), Inches(0.35),
    size=14, bold=True, color=C_ACCENT2)
img(sl, IMGS/"orb_zlayer.png",
    Inches(6.8), Inches(1.7), Inches(6.3), Inches(5.5))

add_rect(sl, 0, H-Inches(0.45), W, Inches(0.45), C_PANEL)
txt(sl,
    "핵심: Z > +2m 구간(Pop2 floater 위험 구간)에서 ORB 앵커 거의 0개 → plateau가 작동 안 함. "
    "MPS semidense는 해당 구간에 28,860pts 존재하나 dist_std ~0.4~1m (위치 불확실).",
    Inches(0.3), H-Inches(0.42), Inches(12.7), Inches(0.38),
    size=12, color=C_ACCENT2)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide C — 앵커 필터링 4단계
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs); fill_bg(sl)
title_bar(sl, "앵커 필터링 Pipeline (4단계)",
          "raw ORB points → 6,492 high-quality anchors")

img(sl, IMGS/"filter_stage0.png",
    Inches(0.3), Inches(1.35), Inches(6.5), Inches(5.8))

table(sl,
    ["단계", "기준", "제거", "남은 pts"],
    [
        ["Stage 0 (raw)", "필터 없음", "—", "7,182"],
        ["Stage 1 (Z-bound)", "카메라 bound 내", "11", "7,171"],
        ["Stage 2 (obs≥3)", "track length ≥ 3", "67", "7,104"],
        ["Stage 3 (kNN iso.)", "k=5, 3×median", "612", "6,492 ✓"],
    ],
    Inches(6.9), Inches(1.35), Inches(6.1),
)

bullet_box(sl, [
    "Z-bound: 카메라 trajectory box 바깥 제거",
    "obs≥3: track이 짧은 불안정 point 제거",
    "kNN isolation: 주변 5개 평균보다 3배 이상",
    "  멀리 떨어진 고립 outlier 제거",
    "최종 6,492개: median kNN 거리 0.144m",
    "  → tau = clip(0.4~0.9 × 0.144, ...) 계산 기준",
], Inches(6.9), Inches(3.2), Inches(6.1), Inches(2.9),
   title="각 단계 설명")

img(sl, IMGS/"filter_stage3.png",
    Inches(6.9), Inches(5.4), Inches(6.1), Inches(1.8)) if (IMGS/"filter_stage3.png").exists() else None


# ═══════════════════════════════════════════════════════════════════════════════
# Slide D — Ellipsoidal Plateau 설계 (v2)
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs); fill_bg(sl)
title_bar(sl, "Ellipsoidal Plateau 설계 시각화 (v2)",
          "kNN PCA로 surface normal 추정 → 법선 tight / 접선 loose")

img(sl, IMGS/"ellipsoid_v2.png",
    Inches(0.3), Inches(1.35), Inches(7.0), Inches(5.9))

bullet_box(sl, [
    "각 앵커에서 kNN-5 이웃의 PCA 분해",
    "최소 고유벡터 = surface normal u_n",
    "나머지 두 벡터 = tangent u_t1, u_t2",
], Inches(7.6), Inches(1.35), Inches(5.4), Inches(1.8),
   title="Normal 추정 방법")

bullet_box(sl, [
    "τ_n = clip(0.4 × h_j, 0.03, 0.30m)",
    "τ_t = clip(0.9 × h_j, 0.03, 0.60m)",
    "h_j = kNN-5 거리 (spacing 추정)",
    "→ 법선 방향은 tight (0.054m 중앙값)",
    "→ 접선 방향은 loose (0.121m 중앙값)",
    "→ Gaussian이 표면 방향으로는 자유롭게 이동",
    "→ 표면 밖으로는 강하게 제한",
], Inches(7.6), Inches(3.3), Inches(5.4), Inches(3.0),
   title="Tau 파라미터")

txt(sl, "XY coverage (layer4): Spherical=11.9%  →  Ellipsoidal=14.0% (+2.1%p)",
    Inches(7.6), Inches(6.45), Inches(5.4), Inches(0.35),
    size=13, bold=True, color=C_GREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide E — Virtual Anchor 생성 (DepthPro v4)
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs); fill_bg(sl)
title_bar(sl, "Virtual Anchor 생성: DepthPro v4",
          "Monocular depth → 3D 투영 → Poisson-disk spacing (D_target=0.5m)")

img(sl, IMGS/"depthpro_completed.png",
    Inches(0.3), Inches(1.35), Inches(7.5), Inches(5.9))

bullet_box(sl, [
    "57장 MPS keyframe 이미지에 DepthPro 추론",
    "픽셀별 depth → 카메라 pose로 3D 투영",
    "stride=16px 스캔, min_hits=2 (multi-view 검증)",
    "Poisson-disk: 0.5m 간격 이하 중복 제거",
    "voxel 크기 0.30m, D_target=0.50m",
], Inches(8.0), Inches(1.35), Inches(5.0), Inches(2.8),
   title="생성 Pipeline")

table(sl,
    ["앵커 소스", "점 수", "비고"],
    [
        ["DepthPro virtual (v4)", "7,338 pts", "★ exp18/19/20 사용"],
        ["Metric3D virtual (v4)", "9,110 pts", "exp17 사용, 성능 최악"],
        ["Depth-Anything-V2 (v4)", "~8k pts", "미사용"],
        ["ORB-SLAM filtered", "6,492 pts", "exp15/16 사용"],
    ],
    Inches(8.0), Inches(4.3), Inches(5.0),
)

bullet_box(sl, [
    "DepthPro: 하부(L1) coverage 0% (depth 불안정)",
    "L2-L7 일관성 높음 → 중간층 coverage 양호",
    "MPS keyframe pose 기준 생성",
    "  → MPS 학습에 좌표계 일치 (exp19/20)",
    "  → ORB 학습에는 좌표계 불일치 (exp18은 사실상 무효)",
], Inches(8.0), Inches(5.65), Inches(5.0), Inches(1.55),
   title="품질 특성", tc=C_ACCENT2)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide F — Round3 & Round5 개입 결과
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs); fill_bg(sl)
title_bar(sl, "Round 3 & Round 5 개입 실험 요약",
          "Perturbation 분석 / camera-bound init filtering 효과")

txt(sl, "Round 3: Perturbation 분석",
    Inches(0.3), Inches(1.3), Inches(6.3), Inches(0.35),
    size=14, bold=True, color=C_ACCENT)
img(sl, IMGS/"round3_perturbation.png",
    Inches(0.3), Inches(1.7), Inches(6.3), Inches(3.0))

bullet_box(sl, [
    "opacity / scale / position을 perturbation해서",
    "  각 Gaussian이 렌더링에 얼마나 기여하는지 측정",
    "대부분 Gaussian은 영향 없음 → safe to prune",
    "일부 고영향 Gaussian이 전체 품질 좌우",
    "→ gradient-free floater 식별 아이디어 기원",
], Inches(0.3), Inches(4.8), Inches(6.3), Inches(2.4),
   title="Round 3 목적 및 결론")

txt(sl, "Round 5: Init Filtering 효과",
    Inches(6.8), Inches(1.3), Inches(6.3), Inches(0.35),
    size=14, bold=True, color=C_ACCENT2)
img(sl, IMGS/"round5b_intervention.png",
    Inches(6.8), Inches(1.7), Inches(6.3), Inches(3.0))

table(sl,
    ["항목", "before", "after"],
    [
        ["제거 pts", "—", "46,276 (7.38%)"],
        ["Z-outlier @500iter", "46,264개", "508개 (-99%)"],
        ["Z-outlier @30k", "1,474개", "385개"],
        ["|Z| max @30k", "42.71m", "4.85m"],
        ["PSNR @30k", "33.012", "32.855 (-0.16dB)"],
    ],
    Inches(6.8), Inches(4.8), Inches(6.3),
)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide G — exp13 비교 & 전체 PSNR 흐름
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs); fill_bg(sl)
title_bar(sl, "전체 실험 PSNR 흐름 (Round 1~6)",
          "exp08 baseline 33.012 dB 기준")

img(sl, IMGS/"exp13_comparison.png",
    Inches(0.3), Inches(1.35), Inches(6.5), Inches(5.5))

table(sl,
    ["실험", "PSNR @30k", "vs exp08", "비고"],
    [
        ["exp08 (MPS baseline)", "33.012", "—", "현재 best"],
        ["exp12 (sparse depth)", "32.587", "-0.43 dB", "미채택"],
        ["exp13 (init filter)", "32.855", "-0.16 dB", "Pop1 해결"],
        ["exp_orb_baseline", "29.023", "-3.99 dB", "ORB init 한계"],
        ["exp15 (spher, ORB)", "27.908", "-5.10 dB", "좌표계 불일치"],
        ["exp16 (ellip, ORB)", "28.924", "-4.09 dB", "좌표계 불일치"],
        ["exp18 (DP, ORB)", "28.934", "-4.08 dB", "좌표계 불일치"],
        ["exp19 (DP, MPS)", "32.753", "-0.26 dB", "첫 올바른 실험"],
        ["exp20 (scheduled)", "TBD", "TBD", "진행 중"],
    ],
    Inches(6.8), Inches(1.35), Inches(6.3),
)

add_rect(sl, 0, H-Inches(0.5), W, Inches(0.5), C_PANEL)
txt(sl,
    "ORB init 실험(exp15-18)은 좌표계 불일치로 plateau loss 무효. "
    "MPS init(exp19)이 첫 유효 실험: -0.26 dB. "
    "exp20 (lambda schedule)에서 개선 기대.",
    Inches(0.3), H-Inches(0.47), Inches(12.7), Inches(0.43),
    size=12, color=C_ACCENT2)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(PPTX)
print(f"Appended 7 diagnostic slides → {PPTX}")
print(f"Total slides: {len(prs.slides)}")
