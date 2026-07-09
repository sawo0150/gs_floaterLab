#!/usr/bin/env python3
"""
Round 2 슬라이드 추가 (exp20/21 결과 + exp22-24 계획)
기존 report_20260705.pptx에 슬라이드 5개 추가
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

PPTX_PATH = "/home/wosas/Desktop/Incremental_mapping_test/gs_floaterLab/context/ppt/report_20260705.pptx"

# ── Colors ────────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x0F, 0x17, 0x2A)
C_ACCENT  = RGBColor(0x38, 0xBD, 0xF8)
C_ACCENT2 = RGBColor(0xFB, 0xBF, 0x24)
C_GREEN   = RGBColor(0x34, 0xD3, 0x99)
C_RED     = RGBColor(0xF8, 0x71, 0x71)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY    = RGBColor(0x94, 0xA3, 0xB8)
C_PANEL   = RGBColor(0x1E, 0x2D, 0x45)
C_BORDER  = RGBColor(0x38, 0x4F, 0x6E)
C_ORANGE  = RGBColor(0xF9, 0x73, 0x16)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation(PPTX_PATH)

def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def fill_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def slide_title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(0.07), C_ACCENT)
    add_text(slide, title,
             Inches(0.5), Inches(0.18), Inches(12), Inches(0.7),
             size=32, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.88), Inches(12), Inches(0.4),
                 size=16, color=C_ACCENT)

def add_bullet_box(slide, items, l, t, w, h, title=None, title_color=C_ACCENT):
    add_rect(slide, l, t, w, h, C_PANEL)
    add_rect(slide, l, t, Inches(0.04), h, title_color)
    y = t + Inches(0.15)
    if title:
        add_text(slide, title, l + Inches(0.12), y, w - Inches(0.2), Inches(0.35),
                 size=14, bold=True, color=title_color)
        y += Inches(0.38)
    for item in items:
        add_text(slide, f"• {item}", l + Inches(0.12), y,
                 w - Inches(0.2), Inches(0.35), size=13, color=C_WHITE)
        y += Inches(0.34)

def add_table(slide, headers, rows, l, t, w,
              hdr_color=C_ACCENT, odd_color=C_PANEL, even_color=C_BG):
    n_cols = len(headers)
    col_w  = w / n_cols
    row_h  = Inches(0.40)
    for ci, hdr in enumerate(headers):
        add_rect(slide, l + col_w * ci, t, col_w, row_h, hdr_color)
        add_text(slide, hdr,
                 l + col_w * ci + Inches(0.05), t + Inches(0.05),
                 col_w - Inches(0.1), row_h - Inches(0.05),
                 size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        bg = odd_color if ri % 2 == 0 else even_color
        for ci, cell in enumerate(row):
            add_rect(slide, l + col_w * ci, t + row_h * (ri + 1), col_w, row_h, bg)
            txt_color = C_WHITE
            if cell.startswith("-") and "dB" in cell:
                txt_color = C_RED
            elif cell.startswith("+") and "dB" in cell:
                txt_color = C_GREEN
            add_text(slide, cell,
                     l + col_w * ci + Inches(0.05),
                     t + row_h * (ri + 1) + Inches(0.05),
                     col_w - Inches(0.1), row_h - Inches(0.05),
                     size=12, color=txt_color, align=PP_ALIGN.CENTER)
    return t + row_h * (len(rows) + 1)

def psnr_tile(slide, label, psnr_7k, psnr_30k, delta, x, y, w=Inches(3.8)):
    color = C_RED if delta.startswith("-") else C_GREEN
    add_rect(slide, x, y, w, Inches(1.8), C_PANEL)
    add_rect(slide, x, y, Inches(0.04), Inches(1.8), color)
    add_text(slide, label,  x + Inches(0.12), y + Inches(0.12), w - Inches(0.2), Inches(0.32),
             size=13, bold=True, color=C_GRAY)
    add_text(slide, f"@7k:   {psnr_7k} dB",  x + Inches(0.12), y + Inches(0.50), w - Inches(0.2), Inches(0.30),
             size=13, color=C_WHITE)
    add_text(slide, f"@30k: {psnr_30k} dB", x + Inches(0.12), y + Inches(0.85), w - Inches(0.2), Inches(0.30),
             size=20, bold=True, color=C_WHITE)
    add_text(slide, delta, x + Inches(0.12), y + Inches(1.3), w - Inches(0.2), Inches(0.35),
             size=16, bold=True, color=color)


# ═══════════════════════════════════════════════════════════════════════════════
# 섹션 구분 슬라이드 — Round 2
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, W, Inches(0.08), C_ACCENT)
add_rect(sl, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT)
add_text(sl, "Round 2",
         Inches(1), Inches(2.2), Inches(11), Inches(1.0),
         size=56, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_text(sl, "Opacity Weighting · Lambda 강화 · 실패 분석",
         Inches(1), Inches(3.3), Inches(11), Inches(0.6),
         size=24, color=C_GRAY, align=PP_ALIGN.CENTER)
add_text(sl, "exp20 / exp21  |  2026-07-05",
         Inches(1), Inches(4.1), Inches(11), Inches(0.5),
         size=18, color=C_ACCENT2, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide: exp20 결과 분석
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "exp20: Lambda Schedule 결과 — 실패",
                "start_iter=1000 → λ=0.1이 densification과 충돌, Gaussian 수 급감")

# PSNR tiles
psnr_tile(sl, "exp08 (baseline)", "28.3505", "33.012", "+0.00 dB",  Inches(0.4),  Inches(1.4))
psnr_tile(sl, "exp19 (λ=0.01 fixed)", "28.0899", "32.7533", "-0.26 dB", Inches(4.4), Inches(1.4))
psnr_tile(sl, "exp20 (λ schedule)", "27.1255", "31.6934", "-1.32 dB", Inches(8.4), Inches(1.4))

# Failure analysis
add_bullet_box(sl, [
    "start_iter=1000에서 λ=0.10으로 즉시 시작",
    "densification 구간(iter 1000~7000)과 정면 충돌",
    "Gaussian 수: 323k → 218k (-32%) — 대규모 소멸",
    "강한 plateau gradient가 densification의 clone/split을 무력화",
    "photometric coverage 축소 → PSNR 급락",
], Inches(0.4), Inches(3.5), Inches(5.8), Inches(3.3),
   title="실패 원인", title_color=C_RED)

# Fix idea
add_bullet_box(sl, [
    "densification은 iter 7000까지 진행",
    "→ start_iter를 7000 이후로 설정해야 안전",
    "그래야 Gaussian 구조 확립 후 plateau loss 적용",
    "exp21 설계: start_iter=7000, λ=0.10, opacity_weight=True",
], Inches(6.4), Inches(3.5), Inches(6.5), Inches(2.5),
   title="교훈 & exp21 설계", title_color=C_ACCENT)

# Gaussian count evidence
add_rect(sl, Inches(6.4), Inches(6.1), Inches(6.5), Inches(0.8), C_PANEL)
add_rect(sl, Inches(6.4), Inches(6.1), Inches(0.04), Inches(0.8), C_ACCENT2)
add_text(sl, "Gaussian 소멸 증거:  323,864 → 218,053 (-32%)  (exp20 vs baseline@7k)",
         Inches(6.55), Inches(6.2), Inches(6.2), Inches(0.6),
         size=13, color=C_ACCENT2)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide: exp21 결과 분석 — Opacity Weighting
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "exp21: Opacity Weighting 결과 — 예상 외 악화",
                "start_iter=7000, λ=0.10, opacity_weight=True → PSNR -2.24 dB (최악)")

psnr_tile(sl, "exp08 (baseline)", "28.3505", "33.012", "+0.00 dB",  Inches(0.4),  Inches(1.4))
psnr_tile(sl, "exp19 (λ=0.01)", "28.0899", "32.7533", "-0.26 dB", Inches(4.4), Inches(1.4))
psnr_tile(sl, "exp21 (opacity weight)", "28.1377", "30.7702", "-2.24 dB", Inches(8.4), Inches(1.4))

# What opacity_weight does
add_bullet_box(sl, [
    "Loss = opacity × max(D-1,0)²",
    "high-opacity floater → gradient on xyz (move) + opacity (→0)",
    "예상: floater가 투명해져 densification pruning에 걸림",
    "실제: 표면 Gaussian도 plateau 밖이면 opacity 감소 가능",
], Inches(0.4), Inches(3.5), Inches(5.8), Inches(2.7),
   title="Opacity Weighting 원리")

# Why it failed
add_bullet_box(sl, [
    "exp19 coverage: 90%가 plateau 밖 (D_aniso > 1)",
    "즉, 전체 Gaussian의 90%에 opacity 감소 gradient",
    "표면 Gaussian들도 탁 트인 공간에서는 D>1일 수 있음",
    "→ 필요한 scene coverage Gaussian까지 투명화",
    "Gaussian 수 급감 → 렌더링 품질 대폭 저하",
    "Lambda=0.10은 photometric을 이길 만큼 강함",
], Inches(6.4), Inches(3.5), Inches(6.5), Inches(2.7),
   title="실패 원인", title_color=C_RED)

# Summary of all MPS experiments
add_text(sl, "MPS 실험 전체 요약",
         Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.3),
         size=13, bold=True, color=C_ACCENT)
add_table(sl,
    ["실험", "핵심 변경", "PSNR @30k", "vs exp08"],
    [
        ["exp08", "baseline", "33.012", "—"],
        ["exp19", "ellipsoidal λ=0.01 fixed", "32.753", "-0.26 dB"],
        ["exp20", "λ schedule (1000→0.1, 7k→0.03, 15k→0)", "31.693", "-1.32 dB"],
        ["exp21", "opacity_weight + λ=0.10 + start=7k", "30.770", "-2.24 dB (최악)"],
    ],
    Inches(0.4), Inches(6.55), Inches(12.5),
    hdr_color=C_BORDER,
)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide: 새 방향 설계 (섹션 브릿지)
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "문제 재정의 — 왜 Loss만으로는 안 되는가?",
                "Photometric loss와의 경쟁에서 plateau loss가 지는 구조적 이유")

# Problem structure
add_bullet_box(sl, [
    "Photometric gradient: 모든 visible Gaussian에 강력하게 작용",
    "Plateau gradient: 8192 샘플 × 37 iter 주기로 약하게 작용",
    "Floater는 photometric으로 이미지를 설명하고 있음 (고opacity)",
    "→ plateau가 당겨도 photometric이 더 강하게 제자리로 유지",
    "opacity weighting으로 강도 올려도 표면 Gaussian도 피해",
], Inches(0.4), Inches(1.4), Inches(6.0), Inches(3.2),
   title="구조적 한계")

# Two strategies
add_bullet_box(sl, [
    "전략 1: Loss 강화 — exp(D-1) kernel",
    "  D=5에서 gradient 6.7× 증가 (quadratic 대비)",
    "  멀수록 지수적으로 커지는 gradient",
    "  단: lambda를 줄여야 표면 Gaussian 보호",
], Inches(6.6), Inches(1.4), Inches(2.9), Inches(3.2),
   title="방향 A: Exp Loss", title_color=C_ACCENT)

add_bullet_box(sl, [
    "전략 2: 직접 Pruning",
    "  확신 floater는 gradient 없이 바로 제거",
    "  d_euc > 1.5m + opacity > 0.5",
    "  + max_scale > 0.15m",
    "  ambiguous는 loss로 유도",
], Inches(9.7), Inches(1.4), Inches(3.5), Inches(3.2),
   title="방향 B: Adaptive Prune", title_color=C_ACCENT2)

# gradient comparison table
add_text(sl, "Loss Kernel 비교 (λ=0.05 기준)",
         Inches(0.4), Inches(4.8), Inches(12.5), Inches(0.32),
         size=14, bold=True, color=C_ACCENT)
add_table(sl,
    ["D (거리)", "Quadratic (D-1)²", "Exp loss exp(D-1)-1", "비율"],
    [
        ["D = 1.5", "0.25", "0.65", "2.6×"],
        ["D = 2.0", "1.00", "1.72", "1.7×"],
        ["D = 5.0", "16.0", "53.6", "3.4×"],
        ["D = 10.0", "81.0", "2981", "36.8×"],
    ],
    Inches(0.4), Inches(5.15), Inches(8.0),
    hdr_color=C_BORDER,
)

add_rect(sl, Inches(8.6), Inches(4.8), Inches(4.7), Inches(2.6), C_PANEL)
add_rect(sl, Inches(8.6), Inches(4.8), Inches(0.04), Inches(2.6), C_GREEN)
add_text(sl, "핵심 통찰",
         Inches(8.75), Inches(4.95), Inches(4.4), Inches(0.32),
         size=14, bold=True, color=C_GREEN)
add_text(sl,
    "두 전략을 결합하면:\n"
    "• 확신 floater → prune (빠르고 직접적)\n"
    "• 불확실 floater → exp loss 유도 (안전)\n"
    "→ Gaussian survival 최대화하면서\n"
    "   floater 제거 효율 극대화",
    Inches(8.75), Inches(5.35), Inches(4.4), Inches(1.9),
    size=13, color=C_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide: exp22~24 실험 계획
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "exp22~24 실험 계획: Exp Loss + Adaptive Pruning",
                "plateau_loss.py에 exp_loss / adaptive_prune 구현 완료 → 순차 실행 중")

# Implementation highlights
add_rect(sl, Inches(0.4), Inches(1.4), Inches(5.8), Inches(2.5), C_PANEL)
add_rect(sl, Inches(0.4), Inches(1.4), Inches(0.04), Inches(2.5), C_ACCENT)
add_text(sl, "구현 변경사항 (plateau_loss.py)",
         Inches(0.55), Inches(1.55), Inches(5.5), Inches(0.32),
         size=14, bold=True, color=C_ACCENT)
add_text(sl,
    "# 새 config 파라미터:\n"
    "exp_loss: bool = False\n"
    "  → exp(clamp(D-1,0,8))-1 kernel\n\n"
    "adaptive_prune: bool = False\n"
    "adaptive_prune_d_euc: float = 1.5   # m\n"
    "adaptive_prune_opacity: float = 0.5\n"
    "adaptive_prune_scale: float = 0.15  # m\n"
    "adaptive_prune_interval: int = 500\n"
    "adaptive_prune_start_iter: int = 7000\n\n"
    "# post_backward에 _adaptive_prune() 통합",
    Inches(0.55), Inches(1.95), Inches(5.5), Inches(1.75),
    size=11, color=C_WHITE)

# Adaptive prune logic
add_rect(sl, Inches(6.4), Inches(1.4), Inches(6.5), Inches(2.5), C_PANEL)
add_rect(sl, Inches(6.4), Inches(1.4), Inches(0.04), Inches(2.5), C_ACCENT2)
add_text(sl, "Adaptive Pruning 조건",
         Inches(6.55), Inches(1.55), Inches(6.2), Inches(0.32),
         size=14, bold=True, color=C_ACCENT2)
add_text(sl,
    "# post_backward에서 매 500 iter:\n"
    "d_min = min_j ||xyz - anchor_j||   # Euclidean\n\n"
    "far    = d_min > 1.5m\n"
    "conf   = (opacity > 0.5)\n"
    "       | (max_scale > 0.15m)\n\n"
    "prune_mask = far & conf\n"
    "gaussians.prune_points(prune_mask)\n\n"
    "# 애매한 Gaussian은 exp loss로 유도",
    Inches(6.55), Inches(1.95), Inches(6.2), Inches(1.75),
    size=11, color=C_WHITE)

# Experiment design table
add_text(sl, "Ablation 설계",
         Inches(0.4), Inches(4.1), Inches(12.5), Inches(0.32),
         size=14, bold=True, color=C_ACCENT)
add_table(sl,
    ["실험", "Exp Loss", "Adaptive Prune", "Opacity Weight", "λ", "Start iter", "목적"],
    [
        ["exp21 (끝남)", "✗", "✗", "✓", "0.10", "7000", "opacity weight 기준선"],
        ["exp22", "✓", "✗", "✓", "0.05", "7000", "exp loss 효과 단독 검증"],
        ["exp23", "✗", "✓", "✓", "0.10", "7000", "adaptive pruning 효과 단독 검증"],
        ["exp24", "✓", "✓", "✓", "0.05", "7000", "두 방법 결합 (full strategy)"],
    ],
    Inches(0.4), Inches(4.45), Inches(12.5),
    hdr_color=C_BORDER,
)

# Timeline / status
add_rect(sl, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.8), C_PANEL)
add_rect(sl, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.04), C_GREEN)
add_text(sl,
    "현재 상태:  exp21 완료(30.770) → exp22 자동 시작 중 → exp23 → exp24 순차 실행 예정  "
    "|  예상 완료: ~3시간 후  |  로그: results/run_seq_22_23_24.log",
    Inches(0.55), Inches(6.65), Inches(12.1), Inches(0.65),
    size=13, color=C_GREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(PPTX_PATH)
print(f"Saved: {PPTX_PATH}")
print(f"Total slides: {len(prs.slides)}")
