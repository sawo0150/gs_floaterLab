#!/usr/bin/env python3
"""
3DGS Floater 제거 실험 보고서 (2026-07-04~05) 슬라이드 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

OUT = "/home/wosas/Desktop/Incremental_mapping_test/gs_floaterLab/context/ppt/report_20260705.pptx"

# ── Colors ────────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
C_ACCENT   = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
C_ACCENT2  = RGBColor(0xFB, 0xBF, 0x24)   # amber
C_GREEN    = RGBColor(0x34, 0xD3, 0x99)   # emerald
C_RED      = RGBColor(0xF8, 0x71, 0x71)   # rose
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
C_PANEL    = RGBColor(0x1E, 0x2D, 0x45)   # slightly lighter navy
C_BORDER   = RGBColor(0x38, 0x4F, 0x6E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def slide_title_bar(slide, title, subtitle=None):
    """Top accent bar + title text."""
    add_rect(slide, 0, 0, W, Inches(0.07), C_ACCENT)
    add_text(slide, title,
             Inches(0.5), Inches(0.18), Inches(12), Inches(0.7),
             size=32, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.88), Inches(12), Inches(0.4),
                 size=16, color=C_ACCENT)

def add_bullet_box(slide, items, l, t, w, h, title=None, title_color=C_ACCENT):
    """Panel box with optional title and bullet list."""
    add_rect(slide, l, t, w, h, C_PANEL)
    add_rect(slide, l, t, Inches(0.04), h, C_ACCENT)
    y = t + Inches(0.15)
    if title:
        add_text(slide, title, l + Inches(0.12), y, w - Inches(0.2), Inches(0.35),
                 size=14, bold=True, color=title_color)
        y += Inches(0.38)
    for item in items:
        add_text(slide, f"• {item}", l + Inches(0.12), y,
                 w - Inches(0.2), Inches(0.35), size=13, color=C_WHITE)
        y += Inches(0.34)

def add_table(slide, headers, rows, l, t, w,
              hdr_color=C_ACCENT, odd_color=C_PANEL, even_color=C_BG):
    """Simple manual table using rectangles + textboxes."""
    n_cols = len(headers)
    col_w  = w / n_cols
    row_h  = Inches(0.40)
    # header
    for ci, hdr in enumerate(headers):
        add_rect(slide, l + col_w * ci, t, col_w, row_h, hdr_color)
        add_text(slide, hdr,
                 l + col_w * ci + Inches(0.05), t + Inches(0.05),
                 col_w - Inches(0.1), row_h - Inches(0.05),
                 size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    # rows
    for ri, row in enumerate(rows):
        bg = odd_color if ri % 2 == 0 else even_color
        for ci, cell in enumerate(row):
            add_rect(slide, l + col_w * ci, t + row_h * (ri + 1), col_w, row_h, bg)
            # color coding for delta column
            txt_color = C_WHITE
            if cell.startswith("-") and "dB" in cell:
                txt_color = C_RED
            elif cell.startswith("+") and "dB" in cell:
                txt_color = C_GREEN
            add_text(slide, cell,
                     l + col_w * ci + Inches(0.05),
                     t + row_h * (ri + 1) + Inches(0.05),
                     col_w - Inches(0.1), row_h - Inches(0.05),
                     size=12, color=txt_color, align=PP_ALIGN.CENTER)
    return t + row_h * (len(rows) + 1)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, W, Inches(0.08), C_ACCENT)
add_rect(sl, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT)

add_text(sl, "3D Gaussian Splatting",
         Inches(1), Inches(1.8), Inches(11), Inches(1.0),
         size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Floater 제거 실험 보고서",
         Inches(1), Inches(2.7), Inches(11), Inches(0.9),
         size=40, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_text(sl, "3D Plateau Loss 설계 · 구현 · 실험 (Round 6)",
         Inches(1), Inches(3.6), Inches(11), Inches(0.6),
         size=22, color=C_GRAY, align=PP_ALIGN.CENTER)
add_text(sl, "2026-07-04 ~ 2026-07-05",
         Inches(1), Inches(4.3), Inches(11), Inches(0.5),
         size=18, color=C_ACCENT2, align=PP_ALIGN.CENTER)

add_text(sl, "Scene: Aria MPS  |  OpenMAVIS 3DGS  |  RTX 5070 Ti",
         Inches(1), Inches(5.5), Inches(11), Inches(0.4),
         size=14, color=C_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — 연구 배경
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "연구 배경: 왜 Floater 제거인가?",
                "3DGS는 고품질 렌더링을 달성하지만 floater artifact가 시각 품질을 저하시킴")

# left panel: what is floater
add_bullet_box(sl, [
    "학습 중 densification으로 생성된 Gaussian이",
    "실제 표면과 관계없는 공중에 위치",
    "Photometric loss가 보이지 않는 곳은 제어 못함",
    "Ray-triggered gradient의 구조적 한계",
], Inches(0.4), Inches(1.5), Inches(4.1), Inches(2.5),
   title="Floater란?")

# right panel: two populations
add_bullet_box(sl, [
    "Pop1: SLAM init outlier (|Z| 수십~수만 m)",
    "  → camera-bound filtering으로 99% 해결",
    "Pop2: densification 중 생성 floater (|Z|=3~6m)",
    "  → 아직 미해결, 이번 실험의 핵심 타겟",
], Inches(4.7), Inches(1.5), Inches(4.1), Inches(2.5),
   title="Floater 두 종류", title_color=C_ACCENT2)

# current best
add_rect(sl, Inches(9.0), Inches(1.5), Inches(3.9), Inches(2.5), C_PANEL)
add_rect(sl, Inches(9.0), Inches(1.5), Inches(0.04), Inches(2.5), C_GREEN)
add_text(sl, "현재 Best", Inches(9.12), Inches(1.65), Inches(3.6), Inches(0.35),
         size=14, bold=True, color=C_GREEN)
add_text(sl, "exp08 (OpenMAVIS MPS)", Inches(9.12), Inches(2.05), Inches(3.6), Inches(0.35),
         size=13, color=C_WHITE, bold=True)
add_text(sl, "PSNR @30k = 33.012 dB", Inches(9.12), Inches(2.45), Inches(3.6), Inches(0.35),
         size=20, bold=True, color=C_ACCENT2)
add_text(sl, "Gaussians: 323,864", Inches(9.12), Inches(2.85), Inches(3.6), Inches(0.35),
         size=13, color=C_GRAY)
add_text(sl, "MPS semidense init (626k pts)", Inches(9.12), Inches(3.2), Inches(3.6), Inches(0.35),
         size=13, color=C_GRAY)

# round 1-5 summary
add_bullet_box(sl, [
    "Round 1-2: Pop1 필터 (Z-clip >42m)  → PSNR 33.012 확보",
    "Round 3-4: Sparse depth prior (lambda=0.002) → -0.43 dB 악화 (미채택)",
    "Round 5: camera-bound init filtering → Pop1 완전 해결, Pop2 잔존 확인",
    "Round 6 (이번): 3D Plateau Loss로 Pop2 제어 시도",
], Inches(0.4), Inches(4.2), Inches(12.5), Inches(2.8),
   title="Round 1~5 요약")


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3 — Plateau Loss 개념
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "3D Plateau Loss 개념",
                "앵커 주변에 plateau 영역을 만들어 Gaussian을 표면으로 끌어당기는 정규화")

# Core idea
add_text(sl, "핵심 아이디어",
         Inches(0.4), Inches(1.4), Inches(6), Inches(0.4),
         size=16, bold=True, color=C_ACCENT)
add_text(sl,
    "Photometric loss는 카메라에 보이는 Gaussian에만 gradient를 줌.\n"
    "→ 보이지 않는 floater는 제어 불가.\n\n"
    "Plateau Loss는 모든 Gaussian에 anchor 기반 regularization을 적용:\n"
    "  • Plateau 내부 (D < 1): gradient = 0  (photometric loss가 자유롭게 작동)\n"
    "  • Plateau 외부 (D ≥ 1): quadratic hinge로 안쪽으로 당김",
    Inches(0.4), Inches(1.85), Inches(6.2), Inches(2.4),
    size=13, color=C_WHITE)

# Formula box
add_rect(sl, Inches(0.4), Inches(4.4), Inches(6.2), Inches(2.7), C_PANEL)
add_rect(sl, Inches(0.4), Inches(4.4), Inches(0.04), Inches(2.7), C_ACCENT2)
add_text(sl, "수식",
         Inches(0.55), Inches(4.55), Inches(6), Inches(0.35),
         size=14, bold=True, color=C_ACCENT2)
add_text(sl,
    "Spherical:  D_min = min_j  ||x - p_j|| / tau_j\n\n"
    "Ellipsoidal:  D_aniso_j = sqrt( (Δ·u_t1/τ_t)² + (Δ·u_t2/τ_t)² + (Δ·u_n/τ_n)² )\n\n"
    "Loss = mean_n  max( min_j D(x_n) - 1, 0 )²   (quadratic hinge)\n\n"
    "τ_n = clip(0.4 × h_j, 0.03, 0.30m)    ← surface normal 방향 (tight)\n"
    "τ_t = clip(0.9 × h_j, 0.03, 0.60m)    ← tangential 방향 (loose)\n"
    "h_j = kNN-5 거리 (anchor spacing 추정)",
    Inches(0.55), Inches(4.95), Inches(5.9), Inches(2.0),
    size=12, color=C_WHITE)

# Right: two types comparison
add_text(sl, "Spherical vs Ellipsoidal",
         Inches(6.9), Inches(1.4), Inches(6), Inches(0.4),
         size=16, bold=True, color=C_ACCENT)

add_bullet_box(sl, [
    "tau = clip(0.6 × h_j, 0.05, 0.60m)",
    "모든 방향 동일 반경 → 구형 plateau",
    "앵커 과밀집 시 Gaussian이 한 점에 몰림",
    "XY coverage 11.9% (layer 4 기준)",
], Inches(6.9), Inches(1.85), Inches(2.9), Inches(2.5),
   title="Spherical (A)", title_color=C_GRAY)

add_bullet_box(sl, [
    "법선 방향(τ_n) tight, 접선(τ_t) loose",
    "표면 방향은 허용, 법선 방향만 제한",
    "Gaussian 분산 → 과밀집 완화",
    "XY coverage 14.0% (layer 4 기준)",
], Inches(10.05), Inches(1.85), Inches(2.9), Inches(2.5),
   title="Ellipsoidal (B)", title_color=C_ACCENT)

# anchor types
add_text(sl, "앵커 종류",
         Inches(6.9), Inches(4.5), Inches(6), Inches(0.35),
         size=14, bold=True, color=C_ACCENT)
add_table(sl,
    ["앵커 소스", "점 수", "특징"],
    [
        ["ORB-SLAM filtered", "6,492 pts", "실제 삼각화, sparse"],
        ["DepthPro virtual", "7,338 pts", "monocular depth 투영 (D_target=0.5m)"],
        ["Metric3D virtual", "9,110 pts", "monocular depth 투영 (D_target=0.5m)"],
    ],
    Inches(6.9), Inches(4.9), Inches(6.0),
)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4 — 구현 및 버그 수정
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "구현 과정에서 발견된 버그 2건",
                "모든 실험이 iter ~5000에서 crash → 디버깅 후 수정")

# Bug 1
add_rect(sl, Inches(0.4), Inches(1.4), Inches(6.0), Inches(2.85), C_PANEL)
add_rect(sl, Inches(0.4), Inches(1.4), Inches(0.04), Inches(2.85), C_RED)
add_text(sl, "Bug #1 — tmp_radii NoneType crash",
         Inches(0.55), Inches(1.55), Inches(5.7), Inches(0.35),
         size=14, bold=True, color=C_RED)
add_text(sl, "gaussian_model.py  line 409",
         Inches(0.55), Inches(1.95), Inches(5.7), Inches(0.3),
         size=12, color=C_GRAY)
add_text(sl,
    "원인:\n"
    "densify_and_prune() 종료 시 self.tmp_radii = None 으로 리셋.\n"
    "이후 Z-clip prune_points() 호출 시 None[mask] → TypeError.\n\n"
    "수정:\n"
    "if self.tmp_radii is not None:\n"
    "    self.tmp_radii = self.tmp_radii[valid_points_mask]",
    Inches(0.55), Inches(2.3), Inches(5.7), Inches(1.65),
    size=12, color=C_WHITE)

# Bug 2
add_rect(sl, Inches(6.9), Inches(1.4), Inches(6.0), Inches(2.85), C_PANEL)
add_rect(sl, Inches(6.9), Inches(1.4), Inches(0.04), Inches(2.85), C_RED)
add_text(sl, "Bug #2 — CUDA device-side assert (크기 불일치)",
         Inches(7.05), Inches(1.55), Inches(5.7), Inches(0.35),
         size=14, bold=True, color=C_RED)
add_text(sl, "train.py  post_backward 실행 순서 오류",
         Inches(7.05), Inches(1.95), Inches(5.7), Inches(0.3),
         size=12, color=C_GRAY)
add_text(sl,
    "원인:\n"
    "Z-clip(post_backward)이 densify_and_prune 이전에 실행.\n"
    "Z-clip: N → N-n_pruned (Gaussian 감소).\n"
    "그 뒤 densify_and_clone이 radii(N) vs xyz(N-n_pruned) 불일치.\n"
    "→ CUDA device-side assert triggered.\n\n"
    "수정:\n"
    "post_backward 호출을 densification 블록 이후로 이동.",
    Inches(7.05), Inches(2.3), Inches(5.7), Inches(1.65),
    size=12, color=C_WHITE)

# Impact
add_rect(sl, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.55), C_PANEL)
add_rect(sl, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.04), C_ACCENT2)
add_text(sl,
    "수정 효과:  두 버그 모두 수정 후 exp15~exp18 전체 30k 완주 성공."
    "  1차 rerun → Bug1 crash, 2차 rerun → Bug2 crash, 3차 rerun → 정상 완료.",
    Inches(0.55), Inches(4.57), Inches(12.1), Inches(0.45),
    size=13, color=C_ACCENT2)

# Implementation details
add_bullet_box(sl, [
    "eval/plateau_loss.py — PlateauLoss, PlateauLossConfig 클래스",
    "Cyclic sampler: sample_size=8192, 모든 Gaussian에 균등하게 gradient 보장",
    "train.py 통합: compute_loss() → backward() 이전, post_backward() → densification 이후",
    "YAML toggle: configs/plateau_loss/*.yaml 로 실험 전환",
    "Pop2 Z-clip: iter 5000~, 1000 iter마다 Z ≥ 2.0m Gaussian 제거",
], Inches(0.4), Inches(5.2), Inches(12.5), Inches(2.1),
   title="구현 세부사항")


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5 — ORB 실험 결과
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "Round 6 실험 결과 — ORB-SLAM Init",
                "Dataset: rgb_3dgs_openmavis_orb_full_301_1253  |  init 7,182pts  |  301 cameras  |  30k iters")

add_table(sl,
    ["실험", "Plateau 타입", "앵커", "PSNR @7k", "PSNR @30k", "vs baseline"],
    [
        ["exp_orb_baseline", "없음", "-", "25.3237", "29.0226", "—"],
        ["exp15_orb_spherical", "Spherical", "ORB 6,492pt", "24.9010", "27.9082", "-1.10 dB"],
        ["exp16_orb_ellipsoidal", "Ellipsoidal", "ORB 6,492pt", "25.0716", "28.9239", "-0.10 dB"],
        ["exp17_orb_metric3d", "Ellipsoidal", "Metric3D 9,110pt", "24.5279", "27.6681", "-1.35 dB"],
        ["exp18_orb_depthpro", "Ellipsoidal", "DepthPro 7,338pt", "25.3857", "28.9343", "-0.09 dB"],
    ],
    Inches(0.4), Inches(1.4), Inches(12.5),
)

# Findings
add_bullet_box(sl, [
    "Ellipsoidal >> Spherical (+1.0 dB, 동일 앵커에서)",
    "DepthPro ≈ ORB(ellipsoidal) >> Metric3D  — 앵커 수가 많다고 좋지 않음",
    "최선(exp18)도 baseline 대비 -0.09 dB  — plateau loss가 개선하지 못함",
    "exp18 PSNR@7k=25.3857 > baseline 25.3237  — 초기는 도움, 후반에 중립화",
], Inches(0.4), Inches(4.0), Inches(6.0), Inches(2.85),
   title="핵심 발견")

# exp15 analysis
add_bullet_box(sl, [
    "low_opacity_ratio: 32% → 62.6% 급증",
    "ORB 6,492 앵커 주변에 Gaussian 과밀집",
    "Photometric loss가 과밀집 Gaussian을 투명화",
    "Z-clip: ORB 좌표계에서 거의 무효 (12개만 제거)",
], Inches(6.6), Inches(4.0), Inches(3.0), Inches(2.85),
   title="exp15 Spherical 악화 원인", title_color=C_RED)

# exp17 analysis
add_bullet_box(sl, [
    "앵커 9,110개 (가장 많음)에도 최악 성능",
    "Metric3D depth 품질 문제 가능성",
    "잘못된 surface normal → plateau 왜곡",
], Inches(9.8), Inches(4.0), Inches(3.1), Inches(2.85),
   title="exp17 Metric3D 최악 원인", title_color=C_RED)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6 — 좌표계 불일치 문제 발견
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "중요 발견: 좌표계 불일치",
                "ORB 실험의 앵커가 MPS world coordinate 기준이었음")

add_rect(sl, Inches(0.4), Inches(1.4), Inches(12.5), Inches(1.8), C_PANEL)
add_rect(sl, Inches(0.4), Inches(1.4), Inches(0.04), Inches(1.8), C_RED)
add_text(sl, "문제 발견",
         Inches(0.55), Inches(1.55), Inches(12), Inches(0.35),
         size=14, bold=True, color=C_RED)
add_text(sl,
    "앵커 생성 스크립트(render_completed_plateau_ellipsoids.py)가 MPS keyframe camera poses로 DepthPro depth를 3D 투영.\n"
    "→ anchors_all_depth_pro.npy의 XYZ 좌표 = MPS world coordinate\n"
    "→ ORB 실험(exp15-18)은 ORB world coordinate에서 학습 → 앵커 위치가 완전히 틀린 공간\n"
    "MPS ↔ ORB: scale / rotation / translation 모두 다름. 앵커가 scene과 무관한 위치에 존재.",
    Inches(0.55), Inches(1.95), Inches(12.1), Inches(1.1),
    size=13, color=C_WHITE)

# Impact analysis
add_bullet_box(sl, [
    "exp16/18이 -0.10 dB에서 큰 피해 없이 끝난 이유:",
    "  → lambda=0.01 약해서 잘못된 앵커로 강하게 못 당김",
    "  → 결과적으로 'plateau loss 사실상 꺼진 상태'로 학습",
    "ORB 실험 결과는 plateau loss의 진짜 효과를 측정하지 못함",
], Inches(0.4), Inches(3.4), Inches(6.0), Inches(2.3),
   title="영향 분석", title_color=C_ACCENT2)

# Solution
add_bullet_box(sl, [
    "앵커가 MPS world space → MPS로 학습하면 좌표 일치",
    "exp19: MPS init + MPS cameras + 동일 DepthPro 앵커",
    "처음으로 올바른 조건에서 plateau loss 효과 측정",
    "MPS semidense init: 626k pts → 323k Gaussians",
], Inches(6.6), Inches(3.4), Inches(6.3), Inches(2.3),
   title="해결책: MPS 실험으로 전환", title_color=C_GREEN)

add_rect(sl, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.6), C_PANEL)
add_rect(sl, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.04), C_ACCENT2)
add_text(sl,
    "교훈: 앵커 기반 regularizer는 앵커와 학습의 world coordinate가 반드시 일치해야 함. "
    "앞으로 anchor 생성 시 학습에 사용할 COLMAP 데이터셋을 명시적으로 지정.",
    Inches(0.55), Inches(5.92), Inches(12.1), Inches(0.5),
    size=13, color=C_ACCENT2)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7 — MPS 실험 (exp19)
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "exp19: MPS Init + DepthPro 앵커 (좌표계 수정)",
                "처음으로 앵커와 학습 좌표계가 일치하는 실험")

# Setup
add_table(sl,
    ["항목", "exp08 (baseline)", "exp19 (MPS+DepthPro)"],
    [
        ["Source", "MPS full (1310 imgs)", "MPS full (1310 imgs)"],
        ["Init pts", "626,811 (semidense)", "626,811 (semidense)"],
        ["Plateau", "없음", "Ellipsoidal, DepthPro 7,338pt"],
        ["Lambda", "-", "0.01 (고정)"],
        ["Start iter", "-", "5000"],
        ["좌표계 일치", "N/A", "✓ (MPS anchor = MPS world)"],
    ],
    Inches(0.4), Inches(1.4), Inches(7.0),
)

# PSNR result
add_rect(sl, Inches(7.6), Inches(1.4), Inches(5.3), Inches(2.65), C_PANEL)
add_rect(sl, Inches(7.6), Inches(1.4), Inches(0.04), Inches(2.65), C_ACCENT)
add_text(sl, "PSNR 결과",
         Inches(7.75), Inches(1.55), Inches(5.0), Inches(0.35),
         size=14, bold=True, color=C_ACCENT)
add_text(sl, "exp08 baseline", Inches(7.75), Inches(2.0), Inches(2.3), Inches(0.35),
         size=13, color=C_GRAY)
add_text(sl, "@7k: 28.3505", Inches(7.75), Inches(2.4), Inches(2.3), Inches(0.35),
         size=14, color=C_WHITE)
add_text(sl, "@30k: 33.0123", Inches(7.75), Inches(2.8), Inches(2.3), Inches(0.35),
         size=20, bold=True, color=C_ACCENT2)
add_text(sl, "exp19 (MPS+DepthPro)", Inches(10.1), Inches(2.0), Inches(2.6), Inches(0.35),
         size=13, color=C_GRAY)
add_text(sl, "@7k: 28.0899", Inches(10.1), Inches(2.4), Inches(2.6), Inches(0.35),
         size=14, color=C_WHITE)
add_text(sl, "@30k: 32.7533", Inches(10.1), Inches(2.8), Inches(2.6), Inches(0.35),
         size=20, bold=True, color=C_RED)
add_text(sl, "Δ = -0.26 dB",
         Inches(7.75), Inches(3.7), Inches(5.0), Inches(0.35),
         size=16, bold=True, color=C_RED, align=PP_ALIGN.CENTER)

# Coverage analysis
add_text(sl, "Plateau Coverage 분석 (30k checkpoint)",
         Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.35),
         size=14, bold=True, color=C_ACCENT)
add_table(sl,
    ["지표", "exp08 (no loss)", "exp19 (λ=0.01)", "변화"],
    [
        ["Gaussians", "323,864", "307,202", "-16,662"],
        ["Inside plateau (D<1)", "9.0%", "10.3%", "+1.3%p"],
        ["D_aniso 중앙값", "3.332", "2.806", "-0.526 ↓"],
        ["Euclidean 중앙값", "1.474 m", "1.273 m", "-0.201m ↓"],
        ["tau_n 중앙값", "0.059 m", "← 동일", ""],
        ["tau_t 중앙값", "0.133 m", "← 동일", ""],
    ],
    Inches(0.4), Inches(4.6), Inches(8.5),
)

add_bullet_box(sl, [
    "Gaussian들이 앵커 쪽으로 당겨짐 (거리 -0.2m)",
    "하지만 plateau 진입 9%→10% 소폭 상승뿐",
    "tau 너무 작음 (6cm / 13cm)",
    "lambda=0.01이 photometric에 밀려 효과 미미",
], Inches(9.1), Inches(4.6), Inches(3.8), Inches(2.6),
   title="분석", title_color=C_ACCENT2)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Coverage 진단 및 Lambda Schedule 설계
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "진단: 왜 plateau가 잘 안 채워지는가?",
                "문제: lambda 너무 약, tau 너무 작, start_iter 너무 늦음")

# Problem analysis
add_bullet_box(sl, [
    "Gaussian → anchor 평균 거리: 1.27 m",
    "Plateau 반경: tau_t=13cm, tau_n=6cm (매우 작음)",
    "30k iter 중 5000~30000 = 25000 iter 동안 loss 적용",
    "lambda=0.01  →  gradient가 photometric loss에 압도됨",
    "Gaussian이 1.27m를 이동해 13cm 안으로 들어가야 함",
    "→ 실질적으로 불가능한 강도",
], Inches(0.4), Inches(1.4), Inches(5.8), Inches(3.3),
   title="현재 문제 구조")

# lambda vs frequency
add_bullet_box(sl, [
    "sample_size=8192 / 307k Gaussians",
    "→ Gaussian 하나당 ~37 iter마다 gradient",
    "→ 25k iter 동안 ~675번 plateau gradient 받음",
    "빈도(sample_size)는 이미 충분",
    "핵심은 gradient 강도(lambda)의 문제",
], Inches(6.4), Inches(1.4), Inches(6.5), Inches(3.3),
   title="Lambda vs 샘플링 빈도")

# Solution: schedule
add_rect(sl, Inches(0.4), Inches(4.85), Inches(12.5), Inches(0.04), C_ACCENT)
add_text(sl, "해결책: 2단계 Lambda Schedule (exp20)",
         Inches(0.4), Inches(4.95), Inches(12.5), Inches(0.4),
         size=16, bold=True, color=C_ACCENT)

add_table(sl,
    ["Iter 범위", "Lambda", "목적"],
    [
        ["iter 1000 ~ 7000", "0.10  (10×↑)", "강한 당김 — densification 중 빠르게 plateau 진입"],
        ["iter 7000 ~ 15000", "0.03  (3×)", "densification 끝, 중간 강도 유지"],
        ["iter 15000 ~ 30000", "0.00  (off)", "photometric loss만 — plateau 내부에서 자유롭게 최적화"],
    ],
    Inches(0.4), Inches(5.4), Inches(12.5),
)

add_text(sl,
    "start_iter=1000  (기존 5000 → 조기 시작, 새로 생긴 Gaussian들도 앵커 근처에서 born되도록)",
    Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.35),
    size=13, color=C_ACCENT2)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9 — exp20 설정 및 구현 변경
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "exp20: Lambda Schedule 구현 및 실험 (진행 중)",
                "plateau_loss.py에 lambda_schedule 파라미터 추가")

# Code change
add_rect(sl, Inches(0.4), Inches(1.4), Inches(6.0), Inches(3.6), C_PANEL)
add_rect(sl, Inches(0.4), Inches(1.4), Inches(0.04), Inches(3.6), C_ACCENT)
add_text(sl, "plateau_loss.py 수정",
         Inches(0.55), Inches(1.55), Inches(5.7), Inches(0.35),
         size=14, bold=True, color=C_ACCENT)
add_text(sl,
    "# PlateauLossConfig 에 추가:\n"
    "lambda_schedule: Optional[list] = None\n"
    "# [[iter, lambda], ...] 오름차순\n\n"
    "# _lambda_at(iteration) 메서드:\n"
    "def _lambda_at(self, iteration):\n"
    "    sched = self.cfg.lambda_schedule\n"
    "    if not sched:\n"
    "        return self.cfg.lambda_plateau\n"
    "    lam = self.cfg.lambda_plateau\n"
    "    for bp_iter, bp_lam in sched:\n"
    "        if iteration >= bp_iter:\n"
    "            lam = bp_lam\n"
    "        else:\n"
    "            break\n"
    "    return lam\n\n"
    "# train.py: lambda_plateau → _lambda_at(iteration)",
    Inches(0.55), Inches(1.95), Inches(5.7), Inches(2.9),
    size=11, color=C_WHITE)

# YAML config
add_rect(sl, Inches(6.6), Inches(1.4), Inches(6.3), Inches(3.6), C_PANEL)
add_rect(sl, Inches(6.6), Inches(1.4), Inches(0.04), Inches(3.6), C_ACCENT2)
add_text(sl, "mps_depthpro_scheduled.yaml",
         Inches(6.75), Inches(1.55), Inches(6.0), Inches(0.35),
         size=14, bold=True, color=C_ACCENT2)
add_text(sl,
    "enabled: true\n"
    "type: ellipsoidal\n"
    "anchor_path: ...anchors_all_depth_pro.npy\n"
    "knn_iso_mult: 0\n"
    "alpha_n: 0.4\n"
    "alpha_t: 0.9\n\n"
    "start_iter: 1000        # 기존 5000 → 조기 시작\n"
    "lambda_plateau: 0.10    # 초기 default\n\n"
    "lambda_schedule:\n"
    "  - [1000,  0.10]       # iter 1000부터 λ=0.10\n"
    "  - [7000,  0.03]       # iter 7000부터 λ=0.03\n"
    "  - [15000, 0.00]       # iter 15000부터 off\n\n"
    "pop2_zclip: true\n"
    "pop2_z_threshold: 2.0",
    Inches(6.75), Inches(1.95), Inches(5.9), Inches(2.9),
    size=11, color=C_WHITE)

# Comparison table
add_text(sl, "exp 비교 요약 (MPS init 기준)",
         Inches(0.4), Inches(5.15), Inches(12.5), Inches(0.35),
         size=14, bold=True, color=C_ACCENT)
add_table(sl,
    ["실험", "Lambda", "Start iter", "PSNR @30k", "상태"],
    [
        ["exp08 (baseline)", "없음", "-", "33.0123", "완료"],
        ["exp19 (fixed λ)", "0.01 고정", "5000", "32.7533 (-0.26dB)", "완료"],
        ["exp20 (scheduled)", "0.10→0.03→0.00", "1000", "TBD", "진행 중"],
    ],
    Inches(0.4), Inches(5.55), Inches(12.5),
)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 10 — 전체 실험 타임라인 & 다음 방향
# ═══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
slide_title_bar(sl, "전체 실험 흐름 & 다음 방향",
                "2026-07-04 ~ 05 작업 요약")

# Timeline
add_text(sl, "실험 타임라인",
         Inches(0.4), Inches(1.4), Inches(8.0), Inches(0.35),
         size=14, bold=True, color=C_ACCENT)

timeline = [
    ("07-04 18:27", "Plateau Loss 설계 완료 (spherical/ellipsoidal 후보 확정)", C_GRAY),
    ("07-04 22:xx", "DepthPro/Metric3D virtual anchor 생성 (v4, D=0.5m)", C_GRAY),
    ("07-05 04:xx", "1차 실험 (exp15-18 ORB) → Bug#1 crash (tmp_radii)", C_RED),
    ("07-05 05:xx", "2차 실험 (exp15-18 ORB) → Bug#2 crash (CUDA assert)", C_RED),
    ("07-05 05:26", "3차 실험 (exp15-18 ORB) → 두 버그 수정 후 성공", C_GREEN),
    ("07-05 07:09", "ORB 실험 완료: exp15=-1.1dB, exp16=-0.10dB, exp17=-1.35dB, exp18=-0.09dB", C_ACCENT2),
    ("07-05 08:xx", "좌표계 불일치 발견 → MPS 실험(exp19) 시작", C_ACCENT2),
    ("07-05 09:xx", "exp19 완료: -0.26dB / Coverage 분석: 9%→10% inside plateau", C_ACCENT2),
    ("07-05 09:41", "Lambda Schedule 설계 → exp20 실행 중 (진행 중)", C_ACCENT),
]

y = Inches(1.85)
for ts, desc, color in timeline:
    add_rect(sl, Inches(0.4), y + Inches(0.09), Inches(0.08), Inches(0.15), color)
    add_text(sl, ts, Inches(0.6), y, Inches(1.6), Inches(0.32), size=11, color=C_GRAY)
    add_text(sl, desc, Inches(2.3), y, Inches(6.0), Inches(0.32), size=11, color=C_WHITE)
    y += Inches(0.34)

# Next steps
add_bullet_box(sl, [
    "exp20 결과 확인: scheduled lambda가 coverage를 올리는지",
    "Tau 확대 실험: alpha_n/alpha_t 키워 plateau 반경 자체를 넓히기",
    "OpenMAVIS + plateau: MPS dense anchor로 더 촘촘한 coverage",
    "Z-clip 재설계: 좌표계 독립적인 카메라 거리 기반 pruning",
    "MPS confidence 활용: dist_std 기반 weighted plateau",
], Inches(8.6), Inches(1.4), Inches(4.7), Inches(4.8),
   title="다음 방향")

# Key takeaways
add_rect(sl, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.04), C_ACCENT)
add_text(sl,
    "핵심 교훈:  앵커-학습 좌표계 일치 필수  |  Ellipsoidal >> Spherical  "
    "|  Lambda schedule로 빠른 수렴 후 photometric 전환이 올바른 방향",
    Inches(0.4), Inches(6.38), Inches(12.5), Inches(0.8),
    size=13, bold=True, color=C_ACCENT2)


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved: {OUT}")
