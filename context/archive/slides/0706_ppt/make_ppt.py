"""
make_ppt.py  —  3DGS Floater 제거 연구 PPT 생성 (42 slides)
"""
from __future__ import annotations
import io, os
from pathlib import Path
from PIL import Image as PILImage
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE   = Path("/home/wosas/Desktop/Incremental_mapping_test/gs_floaterLab")
IMGS   = BASE / "context/ppt/0706_ppt/imgs"
DIAG   = BASE / "results/diagnostic"
OUT    = BASE / "context/ppt/0706_ppt/report_20260706.pptx"

# ── Colours ───────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x0D, 0x1B, 0x2A)
C_BG2     = RGBColor(0x14, 0x26, 0x38)
C_ORANGE  = RGBColor(0xFF, 0x6B, 0x35)
C_BLUE    = RGBColor(0x4E, 0xCD, 0xC4)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_SUB     = RGBColor(0xB0, 0xBE, 0xC5)
C_GREEN   = RGBColor(0x66, 0xBB, 0x6A)
C_YELLOW  = RGBColor(0xFF, 0xA7, 0x26)
C_RED     = RGBColor(0xEF, 0x53, 0x50)
C_DARK    = RGBColor(0x1A, 0x2B, 0x3C)
C_PANEL   = RGBColor(0x0F, 0x22, 0x33)

SW, SH = Inches(13.333), Inches(7.5)   # 16:9

# ── Font helper ───────────────────────────────────────────────────────────────
def _add_run(tf, text, size=18, bold=False, color=None, italic=False):
    run = tf.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_WHITE
    run.font.name = "Noto Sans CJK KR"

def para(tf, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    _add_run(p, text, size=size, bold=bold, color=color)
    return p

# ── Shape helpers ─────────────────────────────────────────────────────────────
def bg(slide, color=None):
    sh = slide.shapes.add_shape(1, 0, 0, SW, SH)
    sh.fill.solid(); sh.fill.fore_color.rgb = color or C_BG
    sh.line.fill.background()
    sh.zorder = 0

def box(slide, x, y, w, h, fill=None, line=None, radius=None):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:  sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:     sh.fill.background()
    if line:  sh.line.color.rgb = line; sh.line.width = Pt(1.5)
    else:     sh.line.fill.background()
    return sh

def textbox(slide, x, y, w, h, text, size=18, bold=False,
            color=None, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.text_frame.word_wrap = wrap
    tf = txb.text_frame
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    _add_run(p, text, size=size, bold=bold, color=color, italic=italic)
    return txb

def title_bar(slide, title, subtitle=None, title_size=30, sub_size=16):
    """Dark bar at top with title."""
    box(slide, 0, 0, 13.333, 1.1, fill=C_DARK)
    textbox(slide, 0.3, 0.08, 12.5, 0.6, title,
            size=title_size, bold=True, color=C_WHITE)
    if subtitle:
        textbox(slide, 0.3, 0.65, 12.5, 0.4, subtitle,
                size=sub_size, color=C_SUB)

def accent_box(slide, x, y, w, h, label, value, label_size=13, val_size=22,
               fill=C_DARK, accent=C_ORANGE):
    box(slide, x, y, w, h, fill=fill)
    # left accent strip
    box(slide, x, y, 0.05, h, fill=accent)
    textbox(slide, x+0.12, y+0.08, w-0.2, 0.3, label, size=label_size, color=C_SUB)
    textbox(slide, x+0.12, y+0.38, w-0.2, h-0.5, value, size=val_size, bold=True, color=C_WHITE)

def img(slide, path, x, y, w, h=None):
    if not Path(path).exists():
        return
    if h:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))

def section_divider(slide, text, color=C_ORANGE):
    """Full-width section label strip."""
    box(slide, 0, 6.9, 13.333, 0.6, fill=color)
    textbox(slide, 0.3, 6.93, 12.5, 0.5, text, size=14, bold=True,
            color=C_BG, align=PP_ALIGN.CENTER)

def insight_box(slide, x, y, w, h, text, size=14, fill=C_DARK, accent=C_BLUE):
    box(slide, x, y, w, h, fill=fill)
    box(slide, x, y, w, 0.04, fill=accent)
    textbox(slide, x+0.15, y+0.12, w-0.3, h-0.2, text, size=size, color=C_WHITE, wrap=True)

# ── PSNR chart ────────────────────────────────────────────────────────────────
def make_psnr_chart():
    out = IMGS / "psnr_chart.png"
    if out.exists(): return out

    labels = ["exp08\nbaseline","exp25\nτ enlarged","exp26\nλ=1→decay",
              "exp19\nλ=0.01","exp13\ncam-filter","exp18\nORB DepthPro",
              "exp20\nscheduled","exp21\nopacity-wt","exp22\nexp-loss",
              "exp23\nadaptive prune"]
    vals   = [33.012, 32.969, 32.674, 32.753, 32.855, 28.934, 31.693, 30.770, 29.917, 26.655]
    base   = 33.012

    # colors
    colors = []
    for v in vals:
        d = v - base
        if d >= -0.05: colors.append("#66BB6A")
        elif d >= -0.35: colors.append("#4ECDC4")
        elif d >= -1.0:  colors.append("#FFA726")
        else:            colors.append("#EF5350")

    BG = "#0D1B2A"
    fig, ax = plt.subplots(figsize=(14, 5.5), facecolor=BG)
    ax.set_facecolor(BG)
    bars = ax.barh(labels[::-1], vals[::-1], color=colors[::-1],
                   height=0.65, edgecolor='none')
    ax.axvline(base, color="#FFD700", linestyle="--", linewidth=1.5, alpha=0.8, label="baseline 33.012")
    for bar, v in zip(bars, vals[::-1]):
        ax.text(v + 0.05, bar.get_y() + bar.get_height()/2,
                f"{v:.3f}", va='center', color='white', fontsize=10, fontweight='bold')
    ax.set_xlim(25, 33.8)
    ax.set_xlabel("PSNR @ 30k (dB)", color='white', fontsize=12)
    ax.tick_params(colors='white', labelsize=10)
    for spine in ax.spines.values(): spine.set_edgecolor('#334455')
    ax.legend(fontsize=10, facecolor='#1A2B3C', labelcolor='white', framealpha=0.8)
    ax.set_title("MPS Plateau Loss 실험 결과 비교", color='white', fontsize=14, fontweight='bold', pad=10)
    plt.tight_layout()
    fig.savefig(str(out), dpi=150, bbox_inches='tight', facecolor=BG)
    plt.close()
    return out

# ── Plateau Loss diagram ──────────────────────────────────────────────────────
def make_plateau_diagram():
    out = IMGS / "plateau_diagram.png"
    if out.exists(): return out
    BG = "#0D1B2A"
    fig, axes = plt.subplots(1, 2, figsize=(13, 5), facecolor=BG)
    fig.patch.set_facecolor(BG)

    # LEFT: Phi(D) plot
    ax = axes[0]; ax.set_facecolor(BG)
    D = np.linspace(0, 4, 300)
    phi_quad = np.where(D > 1, (D-1)**2, 0)
    phi_exp  = np.where(D > 1, np.exp(np.clip(D-1, 0, 8)) - 1, 0)
    ax.fill_between(D, 0, 1.2, where=(D <= 1), alpha=0.15, color='#4ECDC4', label='Plateau (gradient=0)')
    ax.plot(D, phi_quad, color='#4ECDC4', lw=2.5, label='Quadratic (D-1)²')
    ax.plot(D, phi_exp,  color='#FF6B35', lw=2.5, label='Exp (used in exp22)')
    ax.axvline(1.0, color='#FFD700', ls='--', lw=1.5, alpha=0.8)
    ax.text(0.5, 0.55, 'Plateau\n(inside)', ha='center', color='#4ECDC4', fontsize=11, transform=ax.transAxes)
    ax.text(0.88, 0.7, 'Floater\nzone', ha='center', color='#FF6B35', fontsize=11, transform=ax.transAxes)
    ax.set_xlabel("D (normalized distance)", color='white', fontsize=11)
    ax.set_ylabel("Loss Φ(D)", color='white', fontsize=11)
    ax.set_title("Plateau Loss 형태", color='white', fontsize=12, fontweight='bold')
    ax.set_ylim(0, 6)
    ax.tick_params(colors='white')
    for s in ax.spines.values(): s.set_edgecolor('#334455')
    ax.legend(fontsize=9, facecolor='#1A2B3C', labelcolor='white', framealpha=0.8)

    # RIGHT: Gradient comparison at each D
    ax2 = axes[1]; ax2.set_facecolor(BG)
    D2 = np.linspace(1.01, 4, 200)
    g_quad = 2*(D2-1)
    g_exp  = np.exp(np.clip(D2-1, 0, 8))
    ax2.plot(D2, g_quad, color='#4ECDC4', lw=2.5, label='dΦ/dD quadratic')
    ax2.plot(D2, g_exp,  color='#FF6B35', lw=2.5, label='dΦ/dD exp')
    ax2.fill_between(D2, g_quad, g_exp, alpha=0.15, color='#FF6B35', label='Gradient gain')
    for d_mark, label in [(2,'D=2\n1.7×'), (3,'D=3\n3.7×'), (4,'D=4\n9.7×')]:
        gq = 2*(d_mark-1); ge = np.exp(d_mark-1)
        ax2.annotate('', xy=(d_mark, ge), xytext=(d_mark, gq),
                     arrowprops=dict(arrowstyle='<->', color='white', lw=1.2))
        ax2.text(d_mark+0.06, (gq+ge)/2, label, color='white', fontsize=9, va='center')
    ax2.set_xlabel("D (normalized distance from plateau)", color='white', fontsize=11)
    ax2.set_ylabel("Gradient magnitude", color='white', fontsize=11)
    ax2.set_title("Gradient 크기 비교 (exp loss vs quadratic)", color='white', fontsize=12, fontweight='bold')
    ax2.tick_params(colors='white')
    for s in ax2.spines.values(): s.set_edgecolor('#334455')
    ax2.legend(fontsize=9, facecolor='#1A2B3C', labelcolor='white', framealpha=0.8)
    ax2.set_ylim(0, 25)

    plt.tight_layout()
    fig.savefig(str(out), dpi=150, bbox_inches='tight', facecolor=BG)
    plt.close()
    return out

# ── Ray gradient diagram ──────────────────────────────────────────────────────
def make_ray_diagram():
    out = IMGS / "ray_gradient_diagram.png"
    if out.exists(): return out
    BG = "#0D1B2A"
    fig, ax = plt.subplots(figsize=(10, 5.5), facecolor=BG)
    ax.set_facecolor(BG)
    ax.set_xlim(0, 10); ax.set_ylim(0, 7); ax.axis('off')

    # Surface Gaussian
    ax.add_patch(plt.Circle((5, 2), 0.4, color='#4ECDC4', alpha=0.8, zorder=3))
    ax.text(5, 2, 'Surface\nGaussian', ha='center', va='center', fontsize=8, color='white', zorder=4)
    # Camera rays to surface
    for cx, cy in [(1,6),(3,6.5),(7,6.5),(9,6)]:
        ax.annotate('', xy=(5,2.4), xytext=(cx,cy),
                    arrowprops=dict(arrowstyle='->', color='#4ECDC4', lw=1.5, alpha=0.7))
    ax.text(5, 4.8, 'Many rays → Strong photometric gradient ✓',
            ha='center', color='#4ECDC4', fontsize=11, fontweight='bold')

    # Floater
    ax.add_patch(plt.Circle((2, 5), 0.45, color='#FF6B35', alpha=0.8, zorder=3))
    ax.text(2, 5, 'Floater\n(Pop2)', ha='center', va='center', fontsize=8, color='white', zorder=4)
    ax.text(0.2, 3.2, '≈ 0 rays\n→ No photometric\n   gradient ✗',
            color='#FF6B35', fontsize=10, fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#1A0A05', alpha=0.8))
    ax.annotate('', xy=(1.6, 5), xytext=(0.8, 3.7),
                arrowprops=dict(arrowstyle='->', color='#FF6B35', lw=1.5))

    # Plateau force arrow
    ax.annotate('', xy=(4.2, 2.5), xytext=(2.4, 4.6),
                arrowprops=dict(arrowstyle='->', color='#FFD700', lw=2.5))
    ax.text(2.8, 3.9, 'Plateau\nLoss', ha='center', color='#FFD700', fontsize=11, fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#1A1A00', alpha=0.8))

    ax.text(5, 0.5, 'Plateau Loss: 3D regularizer that acts on ALL Gaussians — no ray needed',
            ha='center', color='#B0BEC5', fontsize=11, style='italic')

    fig.savefig(str(out), dpi=150, bbox_inches='tight', facecolor=BG)
    plt.close()
    return out


# ── PDF page helper ───────────────────────────────────────────────────────────
def get_pdf_img(folder_name, page_num):
    """Get PNG path for PDF page (1-indexed)."""
    p = IMGS / f"pdf_{folder_name}" / f"page-{page_num}.png"
    return p if p.exists() else None


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    # Large center gradient strip
    box(sl, 0, 2.8, 13.333, 2.0, fill=RGBColor(0x10, 0x22, 0x35))
    box(sl, 0, 2.78, 13.333, 0.05, fill=C_ORANGE)
    box(sl, 0, 4.78, 13.333, 0.05, fill=C_ORANGE)

    textbox(sl, 1, 3.0, 11.3, 1.2,
            "3DGS Floater 제거 연구",
            size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, 1, 4.05, 11.3, 0.7,
            "Photometric Loss의 한계와 Sparse Support Plateau Loss 접근",
            size=22, color=C_BLUE, align=PP_ALIGN.CENTER)
    textbox(sl, 1, 5.2, 11.3, 0.5,
            "2026-07-04 ~ 07-06    |    OpenMAVIS + 3DGS Custom",
            size=16, color=C_SUB, align=PP_ALIGN.CENTER)

    # Bottom stat row
    for i, (lbl, val) in enumerate([
        ("기준 PSNR", "33.012 dB"),
        ("총 실험", "26 runs"),
        ("Best plateau", "exp25  -0.04 dB"),
    ]):
        accent_box(sl, 1.5 + i*3.5, 6.3, 3.0, 0.95, lbl, val, val_size=20)


def slide_02_background(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "연구 배경 & 목표", "OpenMAVIS 기반 3DGS에서 Floater 발생 원인 규명 및 제거")

    # Left panel
    box(sl, 0.3, 1.3, 6.0, 5.5, fill=C_DARK)
    box(sl, 0.3, 1.3, 6.0, 0.04, fill=C_RED)
    textbox(sl, 0.5, 1.4, 5.7, 0.4, "문제", size=16, bold=True, color=C_RED)
    for t in [
        "• OpenMAVIS 기반 3DGS 학습 시 floater 다수 발생",
        "• Photometric loss는 rendered pixel에만 gradient",
        "  → floater는 gradient를 거의 받지 못함",
        "• 단순 파라미터 조정으로는 해결 한계",
    ]:
        p = sl.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(5.7), Inches(0.45))
        p.text_frame.word_wrap = True
        _add_run(p.text_frame.paragraphs[0], t, size=15, color=C_WHITE)

    # Right panel: baseline stats
    box(sl, 6.7, 1.3, 6.3, 5.5, fill=C_DARK)
    box(sl, 6.7, 1.3, 6.3, 0.04, fill=C_GREEN)
    textbox(sl, 6.9, 1.4, 6.0, 0.4, "현재 Best Baseline: exp08", size=16, bold=True, color=C_GREEN)

    stats = [
        ("PSNR @ 30k", "33.012 dB", C_ORANGE),
        ("Gaussian count", "323,864", C_BLUE),
        ("densify_until_iter", "7,000", C_WHITE),
        ("densify_interval", "200", C_WHITE),
        ("min_opacity_prune", "0.01", C_WHITE),
        ("optimizer_beta1", "0.85", C_WHITE),
    ]
    for i, (k, v, vc) in enumerate(stats):
        y = 1.9 + i * 0.75
        box(sl, 6.9, y, 5.9, 0.65, fill=C_PANEL)
        textbox(sl, 7.05, y+0.08, 3.5, 0.3, k, size=13, color=C_SUB)
        textbox(sl, 7.05, y+0.35, 5.6, 0.25, v, size=15, bold=True, color=vc)

    # Research question
    box(sl, 0.3, 6.55, 12.7, 0.75, fill=RGBColor(0x1A, 0x10, 0x00))
    box(sl, 0.3, 6.55, 0.06, 0.75, fill=C_ORANGE)
    textbox(sl, 0.5, 6.62, 12.4, 0.55,
            "연구 질문: Photometric loss가 닿지 않는 floater를 어떻게 제거할 수 있는가?",
            size=15, bold=True, color=C_ORANGE)


def slide_03_ply(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "exp08 Gaussian 분포 — Side & Top View", "Z-colored | 오렌지 = Pop2 floater (Z > +2m)")

    ply_img = IMGS / "slide03_ply_render.png"
    if ply_img.exists():
        sl.shapes.add_picture(str(ply_img), Inches(0.2), Inches(1.2), Inches(12.9))

    insight_box(sl, 0.2, 6.35, 12.9, 0.9,
                "Pop2 floater: 카메라 위쪽(Z>+2m)에 1,763개 (0.5%) 존재 — 천장 방향, photometric ray 거의 없음",
                size=14, accent=C_ORANGE)


def slide_04_vggt(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "VGGT vs OpenMAVIS 비교 + EVO 궤적 정확도",
              "64-frame subset 기준 | MPS 좌표계 기준 EVO")

    # Left table: 3DGS
    box(sl, 0.3, 1.25, 6.1, 4.8, fill=C_DARK)
    box(sl, 0.3, 1.25, 6.1, 0.04, fill=C_BLUE)
    textbox(sl, 0.45, 1.32, 5.8, 0.4, "64-frame 3DGS 비교 (@7k iter)", size=15, bold=True, color=C_BLUE)
    rows = [("항목","VGGT64","OpenMAVIS64"),
            ("Test PSNR @7k","17.04","18.65 ✓"),
            ("Train PSNR @7k","30.26","34.25 ✓"),
            ("Gaussian count","557k","792k"),
            ("Large scale ratio","0.00018","0.02987"),
            ("Low opacity ratio","0.5575","0.7309")]
    for ri, row in enumerate(rows):
        yy = 1.75 + ri * 0.6
        fill = C_PANEL if ri > 0 else RGBColor(0x1E, 0x35, 0x50)
        box(sl, 0.35, yy, 6.0, 0.55, fill=fill)
        for ci, cell in enumerate(row):
            cc = C_BLUE if ri==0 else (C_GREEN if (ci==2 and ri>0 and '✓' in cell) else C_WHITE)
            textbox(sl, 0.4+ci*1.95, yy+0.1, 1.85, 0.35, cell, size=13, bold=(ri==0), color=cc)

    # Right table: EVO
    box(sl, 6.9, 1.25, 6.1, 4.8, fill=C_DARK)
    box(sl, 6.9, 1.25, 6.1, 0.04, fill=C_ORANGE)
    textbox(sl, 7.05, 1.32, 5.8, 0.4, "EVO 궤적 정확도 (MPS 기준)", size=15, bold=True, color=C_ORANGE)
    evo_rows = [("항목","OpenMAVIS ORB","VGGT64"),
                ("APE RMSE","0.56 m ✓","2.68 m"),
                ("APE median","0.40 m ✓","1.23 m"),
                ("RPE RMSE","1.01 m ✓","1.11 m"),
                ("RPE median","1.03 m","0.17 m ✓")]
    for ri, row in enumerate(evo_rows):
        yy = 1.75 + ri * 0.75
        fill = C_PANEL if ri > 0 else RGBColor(0x1E, 0x35, 0x50)
        box(sl, 6.95, yy, 6.0, 0.68, fill=fill)
        for ci, cell in enumerate(row):
            cc = C_ORANGE if ri==0 else (C_GREEN if '✓' in cell else C_WHITE)
            textbox(sl, 7.0+ci*1.95, yy+0.12, 1.85, 0.42, cell, size=13, bold=(ri==0), color=cc)

    # Bottom conclusion
    box(sl, 0.3, 6.25, 12.7, 0.7, fill=C_PANEL)
    box(sl, 0.3, 6.25, 0.06, 0.7, fill=C_YELLOW)
    textbox(sl, 0.5, 6.32, 12.4, 0.55,
            "결론: VGGT는 compact하지만 render 품질과 궤적 정확도 모두 OpenMAVIS 열세. 현재 대체 불가.",
            size=14, bold=True, color=C_YELLOW)


def slide_05_photo_limit1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "Photometric Loss의 한계 (1): Ray-Triggered Gradient",
              "어떤 Gaussian이 gradient를 받는가?")

    ray_img = make_ray_diagram()
    sl.shapes.add_picture(str(ray_img), Inches(0.3), Inches(1.2), Inches(8.0))

    # Right text panel
    box(sl, 8.6, 1.2, 4.5, 5.5, fill=C_DARK)
    box(sl, 8.6, 1.2, 0.05, 5.5, fill=C_ORANGE)
    textbox(sl, 8.75, 1.3, 4.2, 0.5, "핵심 문제", size=16, bold=True, color=C_ORANGE)
    for t in [
        "Photometric loss는\nrendered pixel에 기여하는\nGaussian에만 gradient 전달.",
        "→ 카메라에서 보이지 않는\n   floater는 이동 신호 없음",
        "→ Ray 분포가 Z 높이별로\n   극단적으로 다름",
        "→ 천장 방향 floater는\n   구조적으로 gradient 없음",
    ]:
        p = sl.shapes.add_textbox(Inches(8.75), Inches(1.85+[0,1.0,1.9,2.8][
            ["Photometric loss는\nrendered pixel에 기여하는\nGaussian에만 gradient 전달.",
             "→ 카메라에서 보이지 않는\n   floater는 이동 신호 없음",
             "→ Ray 분포가 Z 높이별로\n   극단적으로 다름",
             "→ 천장 방향 floater는\n   구조적으로 gradient 없음"].index(t)]), Inches(4.2), Inches(0.85))
        p.text_frame.word_wrap = True
        _add_run(p.text_frame.paragraphs[0], t, size=14,
                 color=C_WHITE if '→' not in t else C_BLUE)

    insight_box(sl, 0.3, 6.35, 12.7, 0.9,
                "다음 슬라이드: Z layer별 ray density를 정량화하여 천장 방향의 gradient 공백을 실제 데이터로 증명",
                size=13, accent=C_BLUE)


def slide_06_photo_limit2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "Photometric Loss의 한계 (2): 공간적 제약 없음",
              "Loss가 표면 근처를 보장하지 않는 이유")

    cards = [
        ("Photometric Loss\n(2D)", "#4ECDC4",
         "이미지 픽셀 색상만 맞추면 됨\n\n공간적 이동 최소화\n\n→ 어디 있어도 색이 맞으면 OK\n\n→ 먼 거리 이동 없음"),
        ("Depth Regularization\n(2.5D)", "#FFA726",
         "2D 이미지 depth map 기반\n\n역시 ray-triggered\n\n표면 구조 일부 반영\n\n→ 근본 한계 동일"),
        ("Plateau Loss\n(3D)  ← 목표", "#66BB6A",
         "Gaussian 3D 위치에 직접 gradient\n\nRay와 무관하게 모든 Gaussian 작용\n\n표면 앵커 주변: gradient = 0\n\n표면 밖 floater: 강한 gradient →"),
    ]
    for i, (title, col, body) in enumerate(cards):
        x = 0.3 + i * 4.35
        box(sl, x, 1.2, 4.1, 5.5, fill=C_DARK)
        box(sl, x, 1.2, 4.1, 0.06, fill=RGBColor(int(col[1:3],16), int(col[3:5],16), int(col[5:],16)))
        textbox(sl, x+0.15, 1.32, 3.8, 0.55, title, size=16, bold=True,
                color=RGBColor(int(col[1:3],16), int(col[3:5],16), int(col[5:],16)))
        textbox(sl, x+0.15, 2.0, 3.8, 4.5, body, size=14, color=C_WHITE, wrap=True)

    insight_box(sl, 0.3, 6.35, 12.7, 0.9,
                "핵심: 표면 근처(plateau 내부)는 3D gradient = 0 → photometric 자유. 밖은 강한 gradient → floater 유도.",
                size=14, accent=C_GREEN)


def slide_07_photo_limit3(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "Photometric Loss의 한계 (3): Gradient 이동 범위 제약",
              "Floater가 local minimum에 갇히는 메커니즘")

    box(sl, 0.3, 1.2, 12.7, 4.9, fill=C_DARK)

    items = [
        ("1", "Photometric gradient 크기", "#4ECDC4",
         "Rendered image와의 pixel 오차에 비례\n→ floater가 색이 잘 맞으면 gradient ≈ 0\n→ 멀리 이동할 이유 없음"),
        ("2", "이동 범위 한계", "#FFA726",
         "3DGS position LR ≈ 1.6e-4 per step\n× 30,000 steps = 4.8 units 최대\n→ 실제로는 loss gradient 균형으로 훨씬 작음"),
        ("3", "Local minimum", "#EF5350",
         "Floater가 특정 위치에서 photometric을\n설명하면 이미 local minimum\n→ 어떤 gradient도 이를 깨지 못함"),
        ("4", "Plateau Loss의 역할", "#66BB6A",
         "앵커(surface)와의 거리가 클수록\n강한 geometric gradient 생성\n→ photometric local minimum 탈출 가능"),
    ]
    for i, (num, title, col, body) in enumerate(items):
        x = 0.5 + i * 3.18
        col_rgb = RGBColor(int(col[1:3],16), int(col[3:5],16), int(col[5:],16))
        box(sl, x, 1.35, 2.95, 4.5, fill=C_PANEL)
        box(sl, x, 1.35, 0.05, 4.5, fill=col_rgb)
        textbox(sl, x+0.12, 1.42, 2.75, 0.4, f"  {num}.", size=22, bold=True, color=col_rgb)
        textbox(sl, x+0.12, 1.85, 2.75, 0.45, title, size=14, bold=True, color=C_WHITE)
        textbox(sl, x+0.12, 2.38, 2.75, 3.3, body, size=13, color=C_SUB, wrap=True)

    insight_box(sl, 0.3, 6.3, 12.7, 0.95,
                "다음 3개 섹션: 실제 데이터로 증명 — (1) Ray density Z-layer, (2) ORB-SLAM points Z-layer, (3) MPS points Z-layer",
                size=13, accent=C_BLUE)


def slide_08_evidence_intro(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "증거 데이터 개요: 3개 분석",
              "Ray density + ORB-SLAM sparse + MPS semi-dense")

    for i, (icon, title, col, desc, note) in enumerate([
        ("📊", "Ray Density Z-layer", "#4ECDC4",
         "카메라 ray가 각 Z 슬랩을 몇 번 통과하는가\n→ Photometric gradient 세기 proxy",
         "Layer 4(카메라 높이) 최강, Layer 8(천장위) 급감"),
        ("🔵", "ORB-SLAM Sparse Points", "#FF6B35",
         "ORB-SLAM map points Z-layer별 분포\n→ Plateau 앵커 후보의 공간 커버리지",
         "Layer 8 (Pop2 구간): 고confidence 점 7개 (사실상 0)"),
        ("🟦", "MPS Semi-Dense Points", "#66BB6A",
         "Aria MPS semi-dense 626k pts Z-layer별 분포\n→ ORB 대비 커버리지 및 좌표계 정합 근거",
         "Layer 8: 14,738 pts → DepthPro 앵커 사용 가능"),
    ]):
        x = 0.5 + i * 4.28
        box(sl, x, 1.2, 4.0, 5.5, fill=C_DARK)
        col_rgb = RGBColor(int(col[1:3],16), int(col[3:5],16), int(col[5:],16))
        box(sl, x, 1.2, 4.0, 0.06, fill=col_rgb)
        textbox(sl, x+0.15, 1.32, 3.7, 0.5, title, size=16, bold=True, color=col_rgb)
        textbox(sl, x+0.15, 1.95, 3.7, 1.8, desc, size=14, color=C_WHITE, wrap=True)
        box(sl, x+0.15, 3.9, 3.7, 1.5, fill=C_PANEL)
        box(sl, x+0.15, 3.9, 0.05, 1.5, fill=col_rgb)
        textbox(sl, x+0.25, 4.0, 3.5, 1.35, f"핵심 발견:\n{note}", size=13, color=C_SUB, wrap=True)

    insight_box(sl, 0.3, 6.35, 12.7, 0.9,
                "이 세 데이터를 조합하면: '천장 방향에는 ray도 없고 앵커도 없다 → ORB에서 MPS로 전환이 필수'",
                size=14, accent=C_ORANGE)


def slide_pdf(prs, pdf_folder, page, section_label, insight_text, accent=C_BLUE):
    """Embed a PDF page PNG as fullscreen with overlays."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, color=RGBColor(0x08, 0x14, 0x20))

    p = get_pdf_img(pdf_folder, page)
    if p:
        sl.shapes.add_picture(str(p), Inches(0), Inches(0.5), width=SW)

    # Top banner
    box(sl, 0, 0, 13.333, 0.52, fill=RGBColor(0x0A, 0x18, 0x28))
    box(sl, 0, 0, 0.06, 0.52, fill=accent)
    textbox(sl, 0.15, 0.05, 13.0, 0.42, section_label, size=14, bold=True, color=C_SUB)

    # Bottom insight
    box(sl, 0, 6.6, 13.333, 0.9, fill=RGBColor(0x08, 0x16, 0x28))
    box(sl, 0, 6.6, 0.06, 0.9, fill=accent)
    textbox(sl, 0.15, 6.65, 13.0, 0.8, insight_text, size=13, color=C_WHITE, wrap=True)


def slide_plateau_concept(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "Plateau Loss: 핵심 아이디어",
              "표면 근처는 자유, 표면 밖은 geometric gradient")

    dia = make_plateau_diagram()
    sl.shapes.add_picture(str(dia), Inches(0.2), Inches(1.2), Inches(8.5))

    # Right: formula
    box(sl, 9.0, 1.2, 4.1, 5.8, fill=C_DARK)
    box(sl, 9.0, 1.2, 4.1, 0.05, fill=C_ORANGE)
    textbox(sl, 9.15, 1.3, 3.8, 0.4, "Loss 수식", size=16, bold=True, color=C_ORANGE)

    formulas = [
        ("Plateau 내부 (D ≤ 1)", "gradient = 0\nphotometric 자유", C_GREEN),
        ("Plateau 외부 (D > 1)", "L = (D - 1)²\n→ floater를 surface로", C_ORANGE),
        ("학습 통합", "loss = L_photo\n     + λ × L_plateau", C_BLUE),
        ("λ 역할", "λ = 0: 기존 3DGS 동일\nλ > 0: geometric force 추가", C_WHITE),
    ]
    for i, (title, body, col) in enumerate(formulas):
        y = 1.85 + i * 1.3
        box(sl, 9.1, y, 3.9, 1.15, fill=C_PANEL)
        box(sl, 9.1, y, 0.05, 1.15, fill=col)
        textbox(sl, 9.2, y+0.05, 3.7, 0.35, title, size=12, bold=True, color=col)
        textbox(sl, 9.2, y+0.42, 3.7, 0.65, body, size=13, color=C_WHITE, wrap=True, italic=True)


def slide_anchor_pipeline(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "앵커 소스: DepthPro v4 생성 파이프라인",
              "MPS 좌표계 정합 — 7,338 virtual anchors")

    # Pipeline boxes
    steps = [
        ("57 MPS\nKeyframes", C_BLUE, "MPS 학습에 쓰이는\n57개 키프레임 선택"),
        ("DepthPro\n깊이 추정", C_ORANGE, "Metric depth estimation\n(절대 스케일)"),
        ("Back-projection\n3D 복원", C_BLUE, "MPS intrinsics +\nextrinsics 사용"),
        ("Poisson-disk\nsampling", C_ORANGE, "간격 0.5m\n중복 제거"),
        ("7,338\nVirtual Anchors", C_GREEN, "MPS world 좌표계\n→ 3DGS init과 일치 ✓"),
    ]
    for i, (label, col, sub) in enumerate(steps):
        x = 0.4 + i * 2.56
        box(sl, x, 1.3, 2.3, 1.4, fill=C_DARK)
        box(sl, x, 1.3, 2.3, 0.05, fill=col)
        textbox(sl, x+0.12, 1.42, 2.08, 0.65, label, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        textbox(sl, x+0.1, 2.05, 2.1, 0.58, sub, size=11, color=C_SUB, align=PP_ALIGN.CENTER, wrap=True)
        if i < 4:
            textbox(sl, x+2.35, 1.7, 0.18, 0.5, "→", size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Coordinate system comparison
    box(sl, 0.3, 3.0, 5.8, 3.7, fill=C_DARK)
    box(sl, 0.3, 3.0, 5.8, 0.05, fill=C_GREEN)
    textbox(sl, 0.45, 3.1, 5.5, 0.4, "MPS 학습: 좌표계 정합 ✓", size=15, bold=True, color=C_GREEN)
    for t in ["MPS semi-dense init   →  MPS world 좌표계",
              "DepthPro anchors      →  MPS world 좌표계",
              "                           ↕ 동일 좌표계 → plateau 작동 ✓"]:
        p = sl.shapes.add_textbox(Inches(0.5), Inches(3.6+[0,0.55,1.1][
            ["MPS semi-dense init   →  MPS world 좌표계",
             "DepthPro anchors      →  MPS world 좌표계",
             "                           ↕ 동일 좌표계 → plateau 작동 ✓"].index(t)]),
            Inches(5.5), Inches(0.5))
        col = C_GREEN if '✓' in t else C_WHITE
        _add_run(p.text_frame.paragraphs[0], t, size=14, color=col, italic=('✓' in t))

    box(sl, 6.6, 3.0, 6.4, 3.7, fill=C_DARK)
    box(sl, 6.6, 3.0, 6.4, 0.05, fill=C_RED)
    textbox(sl, 6.75, 3.1, 6.1, 0.4, "ORB 학습: 좌표계 불일치 ✗  (exp15-18 교훈)", size=14, bold=True, color=C_RED)
    for t in ["ORB 학습             →  ORB world 좌표계",
              "DepthPro anchors     →  MPS world 좌표계",
              "                          ↕ 완전히 다른 공간 → exp15-18 무효"]:
        p = sl.shapes.add_textbox(Inches(6.8), Inches(3.6+[0,0.55,1.1][
            ["ORB 학습             →  ORB world 좌표계",
             "DepthPro anchors     →  MPS world 좌표계",
             "                          ↕ 완전히 다른 공간 → exp15-18 무효"].index(t)]),
            Inches(6.0), Inches(0.5))
        col = C_RED if '무효' in t else C_WHITE
        _add_run(p.text_frame.paragraphs[0], t, size=14, color=col)


def slide_spherical_vs_ellipsoidal(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "Plateau 형태: 구형 vs 타원체",
              "kNN PCA로 표면 법선 추정 → 법선 방향 tight, 접선 방향 loose")

    for i, (t, col, tau, coverage, result, body) in enumerate([
        ("구형 (Spherical)", "#4ECDC4",
         "τ = clip(0.6·h_j, 0.05m, 0.60m)\nh_j = kNN-5 3D 거리",
         "Layer4 XY: 11.9%",
         "exp15: -1.11 dB",
         "D = ||x-p|| / τ\n\n방향 무관, 구현 단순\n\n과밀집 문제 발생 가능"),
        ("타원체 (Ellipsoidal)", "#FF6B35",
         "τ_n = clip(0.4·h_j, 0.03m, 0.30m)  [법선, tight]\nτ_t = clip(0.9·h_j, 0.03m, 0.60m)  [접선, loose]",
         "Layer4 XY: 14.0% (+2.1%)",
         "exp18: -0.09 dB  ✓",
         "D = sqrt((Δ·u_t/τ_t)² + (Δ·u_n/τ_n)²)\n\n표면 법선 방향 tight\n접선 방향 loose → 자연스러운 분포"),
    ]):
        x = 0.4 + i * 6.5
        box(sl, x, 1.2, 6.2, 5.5, fill=C_DARK)
        col_rgb = RGBColor(int(col[1:3],16), int(col[3:5],16), int(col[5:],16))
        box(sl, x, 1.2, 6.2, 0.05, fill=col_rgb)
        textbox(sl, x+0.15, 1.3, 5.9, 0.45, t, size=17, bold=True, color=col_rgb)
        textbox(sl, x+0.15, 1.85, 5.9, 0.7, tau, size=12, color=C_SUB, wrap=True, italic=True)
        box(sl, x+0.15, 2.65, 5.9, 0.6, fill=C_PANEL)
        textbox(sl, x+0.3, 2.72, 5.6, 0.45, f"Coverage: {coverage}", size=14, bold=True, color=col_rgb)
        box(sl, x+0.15, 3.35, 5.9, 0.6, fill=C_PANEL)
        rc = C_RED if '-1' in result else C_GREEN
        textbox(sl, x+0.3, 3.42, 5.6, 0.45, f"실험 결과: {result}", size=14, bold=True, color=rc)
        textbox(sl, x+0.15, 4.1, 5.9, 1.5, body, size=13, color=C_WHITE, wrap=True)

    insight_box(sl, 0.3, 6.35, 12.7, 0.9,
                "타원체 > 구형: 같은 앵커에서 1 dB 이상 차이. 표면 법선 추정이 핵심.",
                size=14, accent=C_GREEN)


def slide_coverage(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "Plateau Coverage 분석 & Pop2 구간 문제",
              "Z-layer별 plateau 내부 Gaussian 비율")

    # Coverage table
    box(sl, 0.3, 1.2, 7.5, 5.5, fill=C_DARK)
    textbox(sl, 0.45, 1.25, 7.2, 0.5, "Z-layer별 Coverage (ORB 6,492 앵커 기준)", size=15, bold=True, color=C_BLUE)
    hdrs = ["Layer", "Z 범위", "앵커 수", "Sphere XY%", "Ellipsoid XY%"]
    rows_data = [
        ("2", "[-2.30,-1.54)", "134", "2.1%", "2.5%"),
        ("3", "[-1.54,-0.78)", "1,540", "9.7%", "11.9%"),
        ("4", "[-0.78,-0.02)", "1,896", "11.9%", "14.0%"),
        ("5", "[-0.02,+0.74)", "1,334", "7.8%", "8.9%"),
        ("6", "[+0.74,+1.50)", "1,316", "10.6%", "11.9%"),
        ("7", "[+1.50,+2.26)", "263", "5.4%", "6.5%"),
        ("8 ★", "[+2.26,+3.02)", "9", "0%", "0%"),
    ]
    col_w = [0.6, 1.6, 1.0, 1.35, 1.55]
    col_x = [0.35, 0.95, 2.55, 3.55, 4.9]
    for ci, h in enumerate(hdrs):
        box(sl, col_x[ci], 1.8, col_w[ci], 0.45, fill=RGBColor(0x1E,0x35,0x50))
        textbox(sl, col_x[ci]+0.05, 1.86, col_w[ci]-0.1, 0.33, h, size=12, bold=True, color=C_BLUE)
    for ri, row in enumerate(rows_data):
        y = 2.3 + ri * 0.6
        is_pop2 = ri == 6
        fill = RGBColor(0x2A,0x08,0x08) if is_pop2 else C_PANEL
        box(sl, 0.35, y, 7.1, 0.55, fill=fill)
        for ci, cell in enumerate(row):
            cc = C_RED if is_pop2 else (C_GREEN if cell.endswith('%') and float(cell.replace('%','') or 0) > 10 else C_WHITE)
            textbox(sl, col_x[ci]+0.05, y+0.1, col_w[ci]-0.1, 0.35, cell, size=12, color=cc, bold=is_pop2)

    # Right: solutions
    box(sl, 8.1, 1.2, 5.0, 5.5, fill=C_DARK)
    box(sl, 8.1, 1.2, 5.0, 0.05, fill=C_RED)
    textbox(sl, 8.25, 1.3, 4.7, 0.4, "Pop2 구간 대응 전략", size=15, bold=True, color=C_RED)
    solutions = [
        ("pop2_zclip", "Z ≥ 2.0m Gaussian 주기 hard pruning\n매 1000 iter, start=7000", C_YELLOW),
        ("DepthPro anchors\n(MPS 좌표계)", "MPS 57 keyframes → Pop2 구간도\n14,738 pts 커버 가능", C_GREEN),
        ("Enlarged tau\n(exp25)", "tau_t_max: 0.60 → 2.00m\n인접 앵커 영향범위 확대", C_BLUE),
    ]
    for i, (s_title, s_body, sc) in enumerate(solutions):
        y = 1.9 + i * 1.6
        box(sl, 8.2, y, 4.8, 1.4, fill=C_PANEL)
        box(sl, 8.2, y, 0.05, 1.4, fill=sc)
        textbox(sl, 8.3, y+0.08, 4.5, 0.45, s_title, size=13, bold=True, color=sc)
        textbox(sl, 8.3, y+0.55, 4.5, 0.75, s_body, size=12, color=C_WHITE, wrap=True)

    insight_box(sl, 0.3, 6.35, 12.7, 0.9,
                "Layer 8 (Pop2 구간): ORB 앵커로는 coverage = 0%. → DepthPro(MPS) + enlarged tau + pop2_zclip 병행 필수.",
                size=14, accent=C_RED)


def slide_exp_results(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "MPS 실험 결과 전체 비교",
              "기준: exp08 = 33.012 dB  |  모든 실험 MPS init + DepthPro v4 anchors")

    chart = make_psnr_chart()
    sl.shapes.add_picture(str(chart), Inches(0.2), Inches(1.15), Inches(8.5))

    # Right: key table
    box(sl, 9.0, 1.2, 4.1, 5.8, fill=C_DARK)
    textbox(sl, 9.1, 1.25, 3.9, 0.4, "핵심 설정 요약", size=14, bold=True, color=C_WHITE)
    tbl = [
        ("exp25", "τ 확대+schedule", "32.969", C_GREEN, "-0.04"),
        ("exp26", "τ 확대+λ=1decay", "32.674", C_BLUE, "-0.34"),
        ("exp19", "DepthPro λ=0.01", "32.753", C_BLUE, "-0.26"),
        ("exp20", "scheduled(early)", "31.693", C_YELLOW, "-1.32"),
        ("exp21", "opacity-weight", "30.770", C_YELLOW, "-2.24"),
        ("exp22", "exp-loss", "29.917", C_RED, "-3.10"),
        ("exp23", "adaptive prune", "26.655", C_RED, "-6.36"),
    ]
    for i, (exp, setting, psnr, col, delta) in enumerate(tbl):
        y = 1.7 + i * 0.73
        box(sl, 9.05, y, 3.95, 0.65, fill=C_PANEL)
        box(sl, 9.05, y, 0.04, 0.65, fill=col)
        textbox(sl, 9.12, y+0.05, 0.75, 0.28, exp, size=11, bold=True, color=col)
        textbox(sl, 9.12, y+0.35, 2.5, 0.25, setting, size=10, color=C_SUB)
        textbox(sl, 11.6, y+0.05, 0.9, 0.28, psnr, size=12, bold=True, color=C_WHITE)
        textbox(sl, 11.6, y+0.35, 0.9, 0.25, delta, size=11, color=col)


def slide_analysis_tau(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "핵심 발견: Enlarged Tau (exp25)의 효과",
              "Loss 강화보다 plateau 범위 설정이 더 중요했다")

    # Left: config comparison
    box(sl, 0.3, 1.2, 5.8, 5.5, fill=C_DARK)
    box(sl, 0.3, 1.2, 5.8, 0.05, fill=C_GREEN)
    textbox(sl, 0.45, 1.3, 5.5, 0.45, "exp25 설정 (PSNR: 32.969, -0.04 dB)", size=15, bold=True, color=C_GREEN)
    settings = [
        ("alpha_n", "0.4 → 0.8", "+100%"),
        ("alpha_t", "0.9 → 1.8", "+100%"),
        ("tau_n_max", "0.30 → 0.80 m", "+167%"),
        ("tau_t_max", "0.60 → 2.00 m", "+233%"),
        ("lambda_schedule", "[7k: 0.10] → [15k: 0.03]", "2-phase"),
        ("start_iter", "7000", "densif. 이후"),
    ]
    for i, (k, v, note) in enumerate(settings):
        y = 1.85 + i * 0.62
        box(sl, 0.4, y, 5.6, 0.55, fill=C_PANEL)
        textbox(sl, 0.55, y+0.08, 1.6, 0.35, k, size=12, bold=True, color=C_BLUE)
        textbox(sl, 2.2, y+0.08, 2.4, 0.35, v, size=13, color=C_ORANGE)
        textbox(sl, 4.7, y+0.08, 1.2, 0.35, note, size=11, color=C_SUB)

    # Right: interpretation
    box(sl, 6.5, 1.2, 6.5, 5.5, fill=C_DARK)
    box(sl, 6.5, 1.2, 6.5, 0.05, fill=C_ORANGE)
    textbox(sl, 6.65, 1.3, 6.2, 0.45, "왜 효과가 있었는가?", size=15, bold=True, color=C_ORANGE)
    interpretations = [
        ("Tau 확대 효과",
         "더 많은 surface Gaussian이 plateau 내부에\n→ 이 Gaussian들에 gradient 0\n→ photometric loss 방해 최소화", C_GREEN),
        ("Lambda schedule",
         "초반(iter 7k-15k): λ=0.10 강하게 당기기\n후반(iter 15k-30k): λ=0.03 유지\n→ densification 간섭 없이 수렴", C_BLUE),
        ("exp20 실패와 비교",
         "exp20: start_iter=1000 (densif. 중)\n→ λ=0.10이 densification과 충돌\n→ Gaussian 수 감소 → PSNR -1.32dB", C_YELLOW),
    ]
    for i, (title, body, col) in enumerate(interpretations):
        y = 1.85 + i * 1.62
        box(sl, 6.6, y, 6.3, 1.45, fill=C_PANEL)
        box(sl, 6.6, y, 0.05, 1.45, fill=col)
        textbox(sl, 6.7, y+0.08, 6.1, 0.38, title, size=13, bold=True, color=col)
        textbox(sl, 6.7, y+0.5, 6.1, 0.88, body, size=12, color=C_WHITE, wrap=True)

    insight_box(sl, 0.3, 6.35, 12.7, 0.9,
                "핵심 교훈: Loss를 강화하거나 직접 pruning하는 것보다 plateau의 범위(tau)를 올바르게 설정하는 것이 더 중요하다.",
                size=14, accent=C_GREEN)


def slide_analysis_failure(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "실패 분석: Loss 강화 / Pruning 전략",
              "exp21-23 — 왜 강한 loss와 pruning이 역효과를 냈는가")

    cases = [
        ("exp21\nOpacity-weighted\nλ=0.10", "-2.24 dB", C_YELLOW,
         "High-opacity surface Gaussian까지\nopacity gradient 받음\n→ 필요한 Gaussian들이 투명해짐\n→ PSNR 손실"),
        ("exp22\nExp-loss\nλ=0.05", "-3.10 dB", C_RED,
         "exp(D-1) 커널이 D=1 근방에서도\n강한 gradient → 수렴 불안정\n→ Surface Gaussian도 이동\n→ PSNR 급락"),
        ("exp23\nAdaptive pruning\nd_euc>1.5m", "-6.36 dB", C_RED,
         "d_euc > 1.5m는 DepthPro anchor와\n멀지만 정상인 Gaussian도 포함\n→ 대량 pruning (정상 구조물 제거)\n→ 재앙적 PSNR 손실"),
    ]
    for i, (title, delta, col, analysis) in enumerate(cases):
        x = 0.4 + i * 4.28
        box(sl, x, 1.2, 4.0, 5.5, fill=C_DARK)
        box(sl, x, 1.2, 4.0, 0.05, fill=col)
        textbox(sl, x+0.15, 1.3, 3.7, 0.7, title, size=15, bold=True, color=col)
        box(sl, x+0.15, 2.1, 3.7, 0.65, fill=C_PANEL)
        textbox(sl, x+0.3, 2.2, 3.4, 0.45, f"PSNR: {delta}", size=18, bold=True, color=col)
        textbox(sl, x+0.15, 2.9, 3.7, 2.7, analysis, size=13, color=C_WHITE, wrap=True)

    insight_box(sl, 0.3, 6.35, 12.7, 0.9,
                "공통 원인: 앵커(DepthPro v4)의 밀도/위치가 scene geometry를 완벽히 커버하지 못하는 상황에서 aggressive한 전략은 역효과.",
                size=14, accent=C_ORANGE)


def slide_future_related(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "향후 읽어볼 관련 연구",
              "SplatFace (CVPR 2025)  &  CoMapGS (CVPR 2025)")

    for i, (title, venue, similarity, difference, takeaway, col_rgb) in enumerate([
        ("SplatFace", "CVPR 2025",
         "Splat-to-surface distance loss\n→ 표면 기준 Gaussian 거리 loss\n→ '표면 밖 force' 아이디어와 유사",
         "얼굴 mesh(explicit surface) 기준\ndead-zone(gradient=0 구간) 없음\n→ 항상 당기는 방향 (plateau 없음)",
         "Explicit distance loss의 실제\n구현/효과 학습 목적", C_BLUE),
        ("CoMapGS", "CVPR 2025",
         "proximity MLP classifier\nGaussian↔initial geometry 근접도 예측\ncovisibility 기반 loss weight 차등",
         "학습된 classifier 기반\nplateau (gradient=0 구간) 없음\n명시적 scalar potential field 아님",
         "Sparse point confidence를\nloss weight에 반영하는 방법 참고", C_ORANGE),
    ]):
        x = 0.4 + i * 6.5
        box(sl, x, 1.2, 6.2, 5.7, fill=C_DARK)
        box(sl, x, 1.2, 6.2, 0.06, fill=col_rgb)
        textbox(sl, x+0.15, 1.3, 5.9, 0.5, f"{title}   [{venue}]", size=18, bold=True, color=col_rgb)

        box(sl, x+0.15, 1.95, 5.9, 1.5, fill=C_PANEL)
        box(sl, x+0.15, 1.95, 0.05, 1.5, fill=C_GREEN)
        textbox(sl, x+0.25, 2.0, 5.7, 0.3, "유사점 (우리 아이디어와)", size=11, bold=True, color=C_GREEN)
        textbox(sl, x+0.25, 2.35, 5.7, 1.05, similarity, size=13, color=C_WHITE, wrap=True)

        box(sl, x+0.15, 3.55, 5.9, 1.5, fill=C_PANEL)
        box(sl, x+0.15, 3.55, 0.05, 1.5, fill=C_YELLOW)
        textbox(sl, x+0.25, 3.6, 5.7, 0.3, "차이점", size=11, bold=True, color=C_YELLOW)
        textbox(sl, x+0.25, 3.95, 5.7, 1.05, difference, size=13, color=C_WHITE, wrap=True)

        box(sl, x+0.15, 5.15, 5.9, 1.5, fill=RGBColor(0x10, 0x20, 0x10))
        box(sl, x+0.15, 5.15, 0.05, 1.5, fill=col_rgb)
        textbox(sl, x+0.25, 5.2, 5.7, 0.3, "읽을 이유", size=11, bold=True, color=col_rgb)
        textbox(sl, x+0.25, 5.55, 5.7, 1.0, takeaway, size=13, color=C_WHITE, wrap=True)


def slide_future_gradient(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "향후 계획: Floater Gradient Tracking 분석",
              "Floater가 photometric gradient를 얼마나 받는지 직접 측정")

    box(sl, 0.3, 1.2, 12.7, 1.1, fill=C_DARK)
    box(sl, 0.3, 1.2, 0.06, 1.1, fill=C_ORANGE)
    textbox(sl, 0.5, 1.25, 12.3, 0.9,
            "동기: Floater가 photometric loss로 이동하지 않는다는 것은 알지만,\n"
            "      실제로 gradient 크기가 surface Gaussian 대비 얼마나 작은지는 아직 측정하지 않았다.",
            size=15, color=C_WHITE, wrap=True)

    # Implementation plan
    box(sl, 0.3, 2.45, 7.5, 4.5, fill=C_DARK)
    box(sl, 0.3, 2.45, 0.06, 4.5, fill=C_BLUE)
    textbox(sl, 0.5, 2.52, 7.2, 0.45, "구현 계획", size=15, bold=True, color=C_BLUE)
    code = (
        "# 각 iteration에서:\n"
        "1. floater 후보 식별\n"
        "   (d_euc > threshold  OR  Z > 2m)\n\n"
        "2. loss.backward() 후\n"
        "   grad = gaussians.get_xyz.grad\n\n"
        "3. floater vs surface grad 비교\n"
        "   ratio = |grad_floater| / |grad_surface|\n\n"
        "4. W&B 로깅\n"
        "   'floater_grad_ratio' per iteration"
    )
    box(sl, 0.5, 3.05, 7.1, 3.8, fill=RGBColor(0x0A, 0x12, 0x1C))
    textbox(sl, 0.65, 3.1, 6.9, 3.7, code, size=13, color=C_BLUE, wrap=True, italic=True)

    # Analysis questions
    box(sl, 8.1, 2.45, 5.0, 4.5, fill=C_DARK)
    box(sl, 8.1, 2.45, 0.06, 4.5, fill=C_GREEN)
    textbox(sl, 8.3, 2.52, 4.7, 0.45, "기대 분석 질문", size=15, bold=True, color=C_GREEN)
    questions = [
        "Floater의 photometric gradient가\nsurface Gaussian 대비 얼마나 약한가?",
        "λ를 어느 수준으로 올려야\nphotometric gradient를 이길 수 있는가?",
        "Plateau Loss gradient가 실제로\nfloater에 dominant하게 작용하는가?",
        "Tau 확대가 surface Gaussian의\ngradient를 실제로 줄이는가?",
    ]
    for i, q in enumerate(questions):
        box(sl, 8.2, 3.05+i*0.85, 4.8, 0.75, fill=C_PANEL)
        textbox(sl, 8.35, 3.1+i*0.85, 4.6, 0.65, f"Q{i+1}. {q}", size=12, color=C_WHITE, wrap=True)

    insight_box(sl, 0.3, 6.35, 12.7, 0.9,
                "이 분석을 통해 adaptive λ 설계, 더 나은 tau 기준, 그리고 plateau loss의 실제 작동 메커니즘을 이해할 수 있다.",
                size=13, accent=C_GREEN)


def slide_roadmap(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    title_bar(sl, "향후 연구 로드맵",
              "단기 → 중기 → 장기")

    for i, (phase, col_rgb, items) in enumerate([
        ("단기 (진행 중)", C_ORANGE, [
            "exp24 (exp-loss+pruning) 결과 분석",
            "exp25 PLY 시각화로 floater 실제 감소 확인",
            "Gradient tracking 도구 개발 (W&B 통합)",
            "tau 확대 + opacity-weight 조합 실험",
        ]),
        ("중기", C_BLUE, [
            "SplatFace / CoMapGS 논문 review → 아이디어 통합",
            "Pop2 구간 전용 앵커 보강\n(DepthPro 추가 키프레임, Z=2-4m 타겟)",
            "MPS confidence(dist_std)를\nplateau λ 가중치로 활용",
            "Gradient tracking 결과 기반 adaptive λ 설계",
        ]),
        ("장기", C_GREEN, [
            "Plateau Loss + depth prior 동시 적용\n(exp12 실패 원인 분석 후 재설계)",
            "Floater 제거율 직접 측정하는 새 metric\n(Z>+2m Gaussian 수 @ 30k)",
            "Mesh-based loss (SplatFace 스타일) 검토\n→ plateau dead-zone과 결합",
            "실시간 floater pruning 지표 자동 리포팅",
        ]),
    ]):
        x = 0.3 + i * 4.35
        box(sl, x, 1.2, 4.1, 5.7, fill=C_DARK)
        box(sl, x, 1.2, 4.1, 0.06, fill=col_rgb)
        textbox(sl, x+0.15, 1.3, 3.9, 0.5, phase, size=17, bold=True, color=col_rgb)
        for j, item in enumerate(items):
            box(sl, x+0.15, 1.95+j*1.18, 3.9, 1.05, fill=C_PANEL)
            box(sl, x+0.15, 1.95+j*1.18, 0.05, 1.05, fill=col_rgb)
            textbox(sl, x+0.25, 2.02+j*1.18, 3.7, 0.9, item, size=12, color=C_WHITE, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    IMGS.mkdir(exist_ok=True)
    make_plateau_diagram()
    make_ray_diagram()
    make_psnr_chart()

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Building slides...")

    # ── Section 1: Title & Background ─────────────────────────────────────────
    slide_01_title(prs);         print("  01 title")
    slide_02_background(prs);   print("  02 background")

    # ── Section 2: Baseline & VGGT ───────────────────────────────────────────
    slide_03_ply(prs);           print("  03 PLY render")
    slide_04_vggt(prs);          print("  04 VGGT/EVO")

    # ── Section 3: Problem motivation ─────────────────────────────────────────
    slide_05_photo_limit1(prs);  print("  05 photometric limit 1")
    slide_06_photo_limit2(prs);  print("  06 photometric limit 2")
    slide_07_photo_limit3(prs);  print("  07 photometric limit 3")
    slide_08_evidence_intro(prs); print("  08 evidence intro")

    # ── Section 4: Ray Density PDF (9 pages) ──────────────────────────────────
    ray_insights = [
        "커버 슬라이드: 총 1,300,500 rays, 8개 Z 슬랩, log scale 공유 → 층간 비교 가능",
        "Layer 1/8  Z∈[-3.02,-2.27)m  바닥 아래 — 총 372,687 rays. 씬 하한. 카메라 ray 거의 없음.",
        "Layer 2/8  Z∈[-2.27,-1.51)m  바닥면 부근 — 465,245 rays. 바닥 슬라브, 아래 방향 ray만 통과.",
        "Layer 3/8  Z∈[-1.51,-0.75)m  바닥/하부 벽 — 560,711 rays. 씬 핵심 구조, photometric gradient 집중.",
        "Layer 4/8  Z∈[-0.75, 0.00)m  카메라 눈높이 — 813,594 rays. 가장 많은 ray. Gradient 최강.",
        "Layer 5/8  Z∈[ 0.00,+0.76)m  카메라 위 — ray 감소 시작. 상부 벽면.",
        "Layer 6/8  Z∈[+0.76,+1.52)m  상부 벽/천장 — ray 급감. Photometric gradient 약화.",
        "Layer 7/8  Z∈[+1.52,+2.28)m  천장 위 — 소수 ray만 통과. Pop2 floater 시작 구간.",
        "Layer 8/8  Z∈[+2.28,+3.03)m  Pop2 구간 — ray 거의 없음. Photometric gradient 구조적으로 불가.",
    ]
    for pg in range(1, 10):
        slide_pdf(prs, "ray_density", pg,
                  f"[증거 1: Ray Density]  Page {pg}/9",
                  ray_insights[pg-1], accent=C_BLUE)
        print(f"  {8+pg:02d} ray_density p{pg}")

    # ── Section 5: ORB-SLAM PDF (9 pages) ─────────────────────────────────────
    orb_insights = [
        "커버: ORB-SLAM map points 7,182개. 고confidence(obs≥10): 835개(11.6%). MPS semi-dense와 완전히 다른 출처.",
        "Layer 1/8  Z∈[-3.06,-2.30)m  3pts (obs≥10: 0). 사실상 비어있음. ORB가 이 구간 삼각화 못함.",
        "Layer 2/8  Z∈[-2.30,-1.54)m  185pts (obs≥10: 19). 바닥 슬라브. 고confidence 점 생기기 시작.",
        "Layer 3/8  Z∈[-1.54,-0.78)m  1,645pts (obs≥10: 157). 바닥+하부 벽. 고confidence 집중.",
        "Layer 4/8  Z∈[-0.78,-0.02)m  2,007pts (obs≥10: 232). 가장 잘 관측되는 표면. 고obs 점 최다.",
        "Layer 5/8  Z∈[-0.02,+0.74)m  1,440pts (obs≥10: 194). 카메라 통과 높이 직상.",
        "Layer 6/8  Z∈[+0.74,+1.50)m  1,481pts (obs≥10: 204). 천장+상부 벽. 관측 수 급감 시작.",
        "Layer 7/8  Z∈[+1.50,+2.26)m  380pts (obs≥10: 22). Pop2 floater 시작. ORB 삼각화 각도 부족.",
        "★ Layer 8/8  Z∈[+2.26,+3.02)m  30pts (obs≥10: 7). ORB 고confidence 점 사실상 0개. Plateau 앵커 없음 → ORB 기반으로 Pop2 floater 제거 불가.",
    ]
    for pg in range(1, 10):
        slide_pdf(prs, "orb_zlayers", pg,
                  f"[증거 2: ORB-SLAM Sparse Points]  Page {pg}/9",
                  orb_insights[pg-1], accent=C_ORANGE)
        print(f"  {17+pg:02d} orb_zlayers p{pg}")

    # ── Section 6: MPS PDF (9 pages) ──────────────────────────────────────────
    mps_insights = [
        "커버: MPS semi-dense 626,811pts (ORB의 87배). SLAM outlier 7.4% → camera-bound filter로 제거.",
        "Layer 1/8  Z∈[-3.05,-2.29)m  1,193pts(0.2%). ORB(3pts)와 달리 MPS는 이 구간도 점 존재. 단 outlier 혼재.",
        "Layer 2/8  Z∈[-2.29,-1.53)m  11,562pts(2.0%). 바닥 슬라브 구조 보임. 밀도 증가.",
        "Layer 3/8  Z∈[-1.53,-0.77)m  118,243pts(20.4%). 핵심 구조면. MPS가 dense하게 커버.",
        "Layer 4/8  Z∈[-0.77,-0.01)m  133,581pts(23.0%). 가장 잘 보이는 표면. 최고 밀도.",
        "Layer 5/8  Z∈[-0.01,+0.75)m  114,665pts(19.8%). 카메라 통과 높이 직상. 동일하게 충분.",
        "Layer 6/8  Z∈[+0.75,+1.51)m  127,379pts(21.9%). 천장+상부 벽. ORB(1,481pts) 대비 압도적.",
        "Layer 7/8  Z∈[+1.51,+2.27)m  59,174pts(10.2%). ORB(380pts)와 달리 MPS는 충분. Pop2 시작 구간.",
        "★ Layer 8/8  Z∈[+2.27,+3.03)m  14,738pts(2.5%). ORB(30pts)의 491배. MPS 좌표계 + DepthPro anchors → Pop2 구간 coverage 가능. → MPS 전환의 핵심 근거.",
    ]
    for pg in range(1, 10):
        slide_pdf(prs, "mps_zlayers", pg,
                  f"[증거 3: MPS Semi-Dense Points]  Page {pg}/9",
                  mps_insights[pg-1], accent=C_GREEN)
        print(f"  {26+pg:02d} mps_zlayers p{pg}")

    # ── Section 7: Plateau Loss Design ────────────────────────────────────────
    slide_plateau_concept(prs);         print("  36 plateau concept")
    slide_anchor_pipeline(prs);         print("  37 anchor pipeline")
    slide_spherical_vs_ellipsoidal(prs); print("  38 sphere vs ellipsoid")
    slide_coverage(prs);                print("  39 coverage")

    # ── Section 8: Experiment Results ─────────────────────────────────────────
    slide_exp_results(prs);             print("  40 experiment results")
    slide_analysis_tau(prs);            print("  41 analysis: tau")
    slide_analysis_failure(prs);        print("  42 analysis: failure")

    # ── Section 9: Future Work ─────────────────────────────────────────────────
    slide_future_related(prs);          print("  43 future: related work")
    slide_future_gradient(prs);         print("  44 future: gradient tracking")
    slide_roadmap(prs);                 print("  45 roadmap")

    prs.save(str(OUT))
    print(f"\nSaved: {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
