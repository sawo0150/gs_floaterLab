#!/usr/bin/env python3
"""build_ppt.py (ppt0720) — VIGS-SLAM 실시간성 타임라인 분석, 단독 덱.

기존 ppt0714/0717 덱과 달리 이 덱은 GPU 프로파일러/타임라인이 주제라 다크 테마로
통일(의도적 이탈, make_figures.py의 다크 배경 그림들과 시각적으로 맞춤).
md 파서를 재사용하지 않고 이 주제 전용으로 직접 슬라이드를 구성.

실행: python build_ppt.py
출력: vigs_realtime_timeline_0720.pptx  (PDF 변환은 셸에서 soffice로 별도)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = Path(__file__).parent
IMG = HERE / "img"
OUT = HERE / "vigs_realtime_timeline_0720.pptx"

FONT = "Noto Sans CJK KR"
MONO = "Noto Sans Mono CJK KR"

# ── 팔레트 (make_figures.py와 동일) ──────────────────────────────────
BG = RGBColor(0x12, 0x15, 0x1c)
PANEL = RGBColor(0x1a, 0x1f, 0x2b)
GRID = RGBColor(0x2a, 0x32, 0x42)
TEXT = RGBColor(0xe8, 0xec, 0xf4)
MUTED = RGBColor(0x8b, 0x95, 0xab)
BUDGET = RGBColor(0xff, 0xb4, 0x54)
BLUE = RGBColor(0x4c, 0x86, 0xea)
CORAL = RGBColor(0xef, 0x6a, 0x5c)

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)
BODY_W = SW - 2 * MARGIN


def new_prs():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return bg


def add_text(slide, x, y, w, h, text, size, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return box


def eyebrow_title(slide, eyebrow, title):
    add_text(slide, MARGIN, Inches(0.35), BODY_W, Inches(0.32), eyebrow, 12, BUDGET,
              bold=True)
    add_text(slide, MARGIN, Inches(0.62), BODY_W, Inches(0.6), title, 22, TEXT, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.18), BODY_W, Pt(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GRID
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_image_fit(slide, path, x, y, w, h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    ww, hh = w, w / ar
    if hh > h:
        hh = h
        ww = h * ar
    left = x + (w - ww) / 2
    top = y + (h - hh) / 2
    slide.shapes.add_picture(str(path), left, top, width=ww)


def qa_card(slide, x, y, w, h, q, a):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = GRID
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.16)
    tf.margin_bottom = Inches(0.16)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = q
    r.font.size = Pt(13.5); r.font.bold = True; r.font.color.rgb = TEXT; r.font.name = FONT
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    p2.line_spacing = 1.18
    r2 = p2.add_run(); r2.text = a
    r2.font.size = Pt(11.5); r2.font.color.rgb = MUTED; r2.font.name = FONT


def stat_pill(slide, x, y, w, h, label, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.08
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = color
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.16)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = value
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = color; r.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(10.5); r2.font.color.rgb = MUTED; r2.font.name = FONT


def action_row(slide, x, y, w, tag, tag_color, text):
    tagbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.55), Inches(0.32))
    tagbox.adjustments[0] = 0.5
    tagbox.fill.solid(); tagbox.fill.fore_color.rgb = PANEL
    tagbox.line.color.rgb = tag_color; tagbox.line.width = Pt(1)
    tagbox.shadow.inherit = False
    tf = tagbox.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag
    r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = tag_color; r.font.name = FONT
    add_text(slide, x + Inches(1.75), y - Inches(0.03), w - Inches(1.75), Inches(0.5), text,
              12.5, TEXT, anchor=MSO_ANCHOR.MIDDLE)


def build():
    prs = new_prs()

    # ── Slide 1: Title ──────────────────────────────────────────
    s = blank(prs)
    set_bg(s)
    add_text(s, Inches(1.0), Inches(2.55), Inches(11.3), Inches(0.4),
              "EXP52 / EXP53 · GS_FLOATERLAB", 13, BUDGET, bold=True)
    add_text(s, Inches(1.0), Inches(2.95), Inches(11.3), Inches(1.6),
              "VIGS-SLAM 실시간성 타임라인", 38, TEXT, bold=True)
    add_text(s, Inches(1.0), Inches(3.85), Inches(11.3), Inches(0.6),
              "직렬 체인과 GS Mapping, 얼마나 겹치고 어디가 남는가", 18, MUTED)
    add_text(s, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.5),
              "1253 시퀀스 · 65.1초 실녹화 · imu_cpp+TensorRT 가속 적용 · 2026-07-20", 11.5,
              MUTED)

    # ── Slide 2: 세 가지 질문 ────────────────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "SUMMARY", "세 가지 질문에 대한 답")
    cy = Inches(1.5)
    ch = Inches(1.65)
    gap = Inches(0.18)
    qa_card(s, MARGIN, cy, BODY_W, ch,
            "RTX 5090이면 실시간이 될까?",
            "부분적으로만. VIGS 저자 공식 벤치마크(RTX 5090)에서도 tracking만은 39.83fps(여유)지만 "
            "tracking+mapping은 12.02fps로 여전히 목표(30fps) 미달. 27.0초 오버헤드 중 21.1초(순수 "
            "Python 루프·IPC)는 GPU 세대와 거의 무관하게 남을 가능성이 큼 — GPU 업그레이드는 "
            "필요조건이지 충분조건이 아니다.")
    cy2 = cy + ch + gap
    qa_card(s, MARGIN, cy2, BODY_W, ch,
            "직렬 체인이 gs_mapping보다 작은가?",
            "거의 같다 — 89.6초 vs 90.5초. 어느 한쪽만 0으로 줄여도 나머지 하나가 이미 예산(65.1초)을 "
            "넘는다. 두 축을 동시에 줄여야만 실시간에 도달한다.")
    cy3 = cy2 + ch + gap
    qa_card(s, MARGIN, cy3, BODY_W, ch,
            "27.0초 오버헤드의 정체는?",
            "model_loading 1.5초(1회성) + queue_get_wait 4.4초(리더 프로세스 IPC 대기) + 미계측 잔여 "
            "21.1초(pbar 문자열 포맷·이미지 텐서 복사·save_trajectory의 GPU→CPU 전송+파일 I/O로 추정 "
            "— 아직 세분 계측 안 함, exp53 신규 축 후보).")

    # ── Slide 3: 시나리오 A ─────────────────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "TIMELINE · SCENARIO A", "순차 실행 — gs_parallel 없이 돌리면")
    add_image_fit(s, IMG / "fig_timeline_serial.png", MARGIN, Inches(1.5), BODY_W, Inches(3.0))
    stats_y = Inches(4.75)
    stat_pill(s, MARGIN, stats_y, Inches(3.9), Inches(1.15), "총합", "180.10s", TEXT)
    stat_pill(s, MARGIN + Inches(4.1), stats_y, Inches(3.9), Inches(1.15), "실시간 배수", "2.77×", BUDGET)
    stat_pill(s, MARGIN + Inches(8.2), stats_y, Inches(4.0), Inches(1.15), "gs_mapping 비중", "50.2%", CORAL)
    add_text(s, MARGIN, Inches(6.15), BODY_W, Inches(0.9),
              "한 프레임당 motion_filter→frontend→PGBA→gs_mapping이 전부 같은 스레드에서 순서대로 "
              "실행된다고 가정했을 때, 전체 시퀀스에 걸쳐 각 단계가 누적한 시간의 구성.", 12, MUTED)

    # ── Slide 4: 시나리오 B ─────────────────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "TIMELINE · SCENARIO B", "gs_parallel — 두 스레드가 GPU를 나눠 쓰면")
    add_image_fit(s, IMG / "fig_timeline_parallel.png", MARGIN, Inches(1.5), BODY_W, Inches(3.3))
    stats_y = Inches(5.05)
    stat_pill(s, MARGIN, stats_y, Inches(3.9), Inches(1.15), "총합(실측)", "133.04s", TEXT)
    stat_pill(s, MARGIN + Inches(4.1), stats_y, Inches(3.9), Inches(1.15), "실시간 배수", "2.04×", BUDGET)
    stat_pill(s, MARGIN + Inches(8.2), stats_y, Inches(4.0), Inches(1.15), "순차 대비", "−26.1%", BLUE)
    add_text(s, MARGIN, Inches(6.4), BODY_W, Inches(0.7),
              "두 스레드 다 133.04초에 함께 끝난다 — 한쪽이 먼저 끝나고 기다리는 구조가 아니라, "
              "서로가 서로를 GPU 경합으로 늦추며 같은 시점에 수렴한다.", 12, MUTED)

    # ── Slide 4b: frontend가 병렬일 때 왜 느려지나 (신규) ────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "CONTENTION IS BIDIRECTIONAL", "왜 frontend가 병렬일 때 오히려 느려지나")
    add_image_fit(s, IMG / "fig_frontend_contention.png", MARGIN, Inches(1.5), BODY_W, Inches(3.35))
    add_text(s, MARGIN, Inches(5.05), BODY_W, Inches(0.85),
              "GPU 경합은 매핑 쪽에만 일어나는 게 아니라 양방향이다 — bundle_adjust(BA solve)와 "
              "update_op_forward(GRU)는 frontend에서 가장 무거운 GPU 커널인데, rasterize/backward처럼 "
              "gs_mapping이 던지는 무거운 커널과 같은 SM·메모리 대역폭을 놓고 경쟁한다.", 12.5, TEXT)
    add_text(s, MARGIN, Inches(5.85), BODY_W, Inches(0.85),
              "반면 PGBA(pgba_run)는 sparse pose graph 최적화라 상대적으로 가벼워 거의 안 느려짐(+12%) "
              "— mapping 자체도 거의 그대로(90.5→86.24s). 가설(미검증, nsys 트레이스 필요): mapping이 "
              "던지는 크고 지속시간 긴 커널 뒤에서 frontend의 작고 빈번한 커널들이 줄서서 기다린다.",
              12.5, MUTED)

    # ── Slide 5: 오버랩 효율 ────────────────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "HOW PARALLEL, EXACTLY", "오버랩 효율 — 실제로 얼마나 겹치는가")
    add_image_fit(s, IMG / "fig_overlap_efficiency.png", MARGIN, Inches(1.4), BODY_W, Inches(4.15))
    stats_y = Inches(5.7)
    stat_pill(s, MARGIN, stats_y, Inches(3.9), Inches(1.05), "오버랩 효율", "66.0%", BLUE)
    stat_pill(s, MARGIN + Inches(4.1), stats_y, Inches(3.9), Inches(1.05), "GPU 경합으로 손실", "34.0%", CORAL)
    stat_pill(s, MARGIN + Inches(8.2), stats_y, Inches(4.0), Inches(1.05), "Tracking 실측(같은 런)", "103.72s", MUTED)
    add_text(s, MARGIN, Inches(6.9), BODY_W, Inches(0.5),
              "※ 103.72s는 \"경합이 없었을 때의 tracking 비용\"이 아니라 이 gs_parallel 런에서 "
              "실측된 값 — 그 자체로 이미 위 슬라이드의 경합을 포함하고 있다.", 10.5, MUTED)

    # ── Slide 6: 이론적 최선 ────────────────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "THE FLOOR", "완벽한 병렬화를 가정해도 안 되는 이유")
    add_image_fit(s, IMG / "fig_theoretical_floor.png", MARGIN, Inches(1.45), BODY_W, Inches(2.15))
    add_text(s, MARGIN, Inches(3.75), BODY_W, Inches(0.55),
              "병렬화(exp52)만으로는 구조적으로 도달 불가 — 직렬 체인과 gs_mapping을 각각 "
              "65.1초 아래로 줄여야 한다.", 13.5, TEXT, bold=True)
    ay = Inches(4.15)
    rows = [
        ("EXP53 · 최우선", BUDGET, "Frontend 반복 축소(iters1/iters2) — bundle_adjust+update_op_forward가 frontend 43.2s의 80%. 직렬 체인을 줄이는 가장 직접적인 레버."),
        ("EXP53", BLUE, "keyframe 밀도 억제(motion_filter.thresh) — keyframe 발생이 fps와 무관하다는 게 확인됐으니, 밀도 자체를 낮추는 게 fps를 낮추는 것보다 유효."),
        ("EXP53 · 신규", CORAL, "27.0초 오버헤드 세분 계측 — 21.1초 미계측 잔여의 정체 파악. GPU와 무관하게 남을 항목이라 별도 대응 필요."),
        ("EXP53 · 신규", TEXT, "축A~D를 gs_parallel 켠 상태/끈 상태 둘 다에서 측정 — frontend가 경합 시 +72% 느려짐이 확인됐으니, 순차 측정만으론 병렬 실전 효과를 과대추정할 수 있음."),
        ("EXP52 · 계속", MUTED, "gs_mapping 자체 연산량 감소 — iteration 수·해상도·densify 빈도. RTX 5090에서도 저자 자신이 12.02fps(미달)였던 부분이라 GPU 업그레이드만으론 부족."),
    ]
    for i, (tag, color, text) in enumerate(rows):
        action_row(s, MARGIN, ay + Inches(0.58) * i, BODY_W, tag, color, text)

    # ── Slide 7: 구성요소 세부 분해 ─────────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "COMPONENT BREAKDOWN", "구성요소별 세부 분해")
    add_image_fit(s, IMG / "fig_components.png", MARGIN, Inches(1.45), BODY_W, Inches(5.6))

    prs.save(OUT)
    print(f"[saved] {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
