#!/usr/bin/env python3
"""build_ppt.py — ppt_outline_*.md 를 그대로 파싱해 PPTX 생성.

- md의 `### 슬라이드 N — 제목`을 슬라이드 경계로, `## §..`/`## 백업..`/`## 데이터..`를 섹션 구분 슬라이드로.
- 불릿(들여쓰기=레벨), 마크다운 표, 코드펜스(수식/다이어그램), [시각자료], [발표 노트] 지원.
- [발표 노트] → 슬라이드 노트. [시각자료] → 본문에 색 구분 박스. 내용이 넘치면 "(계속)" 슬라이드로 자동 분할.
- 이미지: img/<파일>.png 가 있으면 해당 슬라이드에 삽입(파일명 규약은 IMG_MAP 참조). 없으면 무시.

사용: python build_ppt.py  (같은 폴더의 ppt_outline_20260711_13.md 사용)
출력: gs_floater_deck.pptx  → soffice로 PDF 변환은 셸에서 별도.
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
MD = HERE / "ppt_outline_20260711_13.md"
OUT = HERE / "gs_floater_deck.pptx"
IMG = HERE / "img"

FONT = "Noto Sans CJK KR"
MONO = "Noto Sans Mono CJK KR"

# 색
C_TITLE = RGBColor(0x1A, 0x1A, 0x2E)
C_ACCENT = RGBColor(0x0F, 0x4C, 0x81)
C_TEXT = RGBColor(0x22, 0x22, 0x22)
C_SUB = RGBColor(0x55, 0x55, 0x55)
C_VIS = RGBColor(0x8A, 0x5A, 0x00)      # 시각자료
C_VISBG = RGBColor(0xFF, 0xF6, 0xE0)
C_CODEBG = RGBColor(0xF2, 0xF3, 0xF7)
C_HDRBG = RGBColor(0x0F, 0x4C, 0x81)
C_HDRTX = RGBColor(0xFF, 0xFF, 0xFF)
C_ROW0 = RGBColor(0xFF, 0xFF, 0xFF)
C_ROW1 = RGBColor(0xEE, 0xF2, 0xF7)
C_SECTIONBG = RGBColor(0x0F, 0x4C, 0x81)

# 슬라이드 제목 앞부분(정규화) → 이미지 파일명(들). img/<name>.png 존재 시 우측에 삽입.
IMG_MAP = {
    "슬라이드 4": ["fig_landscape", "fig_raycoverage"],
    "슬라이드 5": ["fig_label_highlight", "fig_region"],
    "슬라이드 6": ["fig_auc_plateau", "fig_grad_asym"],
    "슬라이드 7": ["fig_rho_section"],
    "슬라이드 10": ["fig_waterfall", "fig_ab_render"],
    "슬라이드 12": ["fig_pareto"],
    "슬라이드 14": ["fig_sobel_ppm"],
    "슬라이드 19": ["fig_crossscene_auc"],
    "슬라이드 21": ["fig_huber"],
    "슬라이드 22": ["fig_anchors"],
    "슬라이드 25": ["fig_slamfree_ladder"],
    "슬라이드 27": ["fig_scenes"],
    "슬라이드 28": ["fig_landscape"],
}

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
BODY_TOP = Inches(1.15)
BODY_BOTTOM = Inches(7.15)
BODY_W = SW - 2 * MARGIN


# ── md 파싱 ────────────────────────────────────────────────
def parse(md_text):
    """섹션/슬라이드 리스트 반환. 각 원소: ('section', title) 또는 ('slide', title, blocks)."""
    lines = md_text.splitlines()
    items = []
    i = 0
    cur = None          # 현재 슬라이드 blocks
    in_code = False
    code_buf = []

    def flush_slide():
        nonlocal cur
        if cur is not None:
            items.append(cur)
            cur = None

    while i < len(lines):
        ln = lines[i]
        s = ln.strip()

        # 코드펜스
        if s.startswith("```"):
            if not in_code:
                in_code = True; code_buf = []
            else:
                in_code = False
                if cur is not None:
                    cur[2].append(("code", code_buf[:]))
            i += 1; continue
        if in_code:
            code_buf.append(ln)
            i += 1; continue

        # 슬라이드 제목
        m = re.match(r"^###\s+슬라이드\s+([\w\-–]+)\s*[—-]\s*(.*)$", s)
        if m:
            flush_slide()
            cur = ("slide", f"슬라이드 {m.group(1)} — {m.group(2)}", [])
            i += 1; continue

        # 섹션 구분
        if s.startswith("## "):
            flush_slide()
            title = s[3:].strip()
            if title.startswith("§") or title.startswith("백업") or title.startswith("데이터 소스"):
                items.append(("section", title))
                # 백업/데이터 소스는 뒤따르는 불릿을 담을 슬라이드도 생성
                if title.startswith("백업") or title.startswith("데이터 소스"):
                    cur = ("slide", title, [])
            i += 1; continue

        if cur is None:
            i += 1; continue

        # 표
        if s.startswith("|"):
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                rows.append(lines[i].strip())
                i += 1
            # 구분행 제거
            parsed = []
            for r in rows:
                cells = [c.strip() for c in r.strip("|").split("|")]
                if all(re.fullmatch(r":?-{2,}:?", c or "-") for c in cells):
                    continue
                parsed.append(cells)
            if parsed:
                cur[2].append(("table", parsed))
            continue

        # 불릿
        m = re.match(r"^(\s*)-\s+(.*)$", ln)
        if m:
            indent = len(m.group(1))
            level = min(indent // 2, 3)
            text = m.group(2).rstrip()
            if text.startswith("[시각자료]"):
                cur[2].append(("visual", text[len("[시각자료]"):].strip(" :")))
            elif text.startswith("[발표 노트]"):
                cur[2].append(("note", text[len("[발표 노트]"):].strip(" :")))
            else:
                cur[2].append(("bullet", level, text))
            i += 1; continue

        i += 1

    flush_slide()
    return items


# ── 텍스트 렌더 헬퍼 ────────────────────────────────────────
def add_runs(para, text, size, color, base_bold=False):
    """**bold**·`code` 인라인 처리. 백틱은 mono 폰트."""
    text = text.replace("→", "→")
    # 토큰: **...** 또는 `...` 또는 일반
    tokens = re.split(r"(\*\*.*?\*\*|`[^`]*`)", text)
    for tok in tokens:
        if not tok:
            continue
        bold = base_bold
        mono = False
        t = tok
        if tok.startswith("**") and tok.endswith("**"):
            bold = True; t = tok[2:-2]
        elif tok.startswith("`") and tok.endswith("`"):
            mono = True; t = tok[1:-1]
        r = para.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = MONO if mono else FONT
    if not tokens:
        r = para.add_run(); r.text = text
        r.font.size = Pt(size); r.font.name = FONT; r.font.color.rgb = color


def est_lines(text, size, width_in):
    chars = max(10, int(width_in * 72 / (size * 0.58)))
    plain = re.sub(r"[*`]", "", text)
    return max(1, -(-len(plain) // chars))


# ── 높이 추정 ───────────────────────────────────────────────
BODY_W_IN = BODY_W / 914400.0

def h_bullet(level, text):
    size = 15 if level == 0 else (13 if level == 1 else 12)
    w = BODY_W_IN - 0.3 * level
    return Inches(0.30 * est_lines(text, size, w) + 0.05)

def h_visual(text):
    return Inches(0.30 * est_lines(text, 12, BODY_W_IN - 0.4) + 0.20)

def h_table(rows):
    return Inches(0.34 * len(rows) + 0.12)

def h_code(codelines):
    return Inches(0.24 * len(codelines) + 0.20)


# ── 슬라이드 빌드 ───────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_title(slide, title):
    box = slide.shapes.add_textbox(MARGIN, Inches(0.28), BODY_W, Inches(0.8))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    add_runs(p, title, 24, C_TITLE, base_bold=True)
    # 밑줄 바
    bar = slide.shapes.add_shape(1, MARGIN, Inches(1.02), BODY_W, Pt(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

def section_slide(prs, title):
    s = blank(prs)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_SECTIONBG; bg.line.fill.background()
    box = s.shapes.add_textbox(Inches(1), Inches(3), SW - Inches(2), Inches(1.5))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_runs(p, title, 34, RGBColor(0xFF, 0xFF, 0xFF), base_bold=True)

def title_slide(prs, title, sub):
    s = blank(prs)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_TITLE; bg.line.fill.background()
    box = s.shapes.add_textbox(Inches(1), Inches(2.6), SW - Inches(2), Inches(2))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    add_runs(p, title, 40, RGBColor(0xFF, 0xFF, 0xFF), base_bold=True)
    p2 = tf.add_paragraph()
    add_runs(p2, sub, 18, RGBColor(0xBB, 0xCC, 0xEE))

def render_table(slide, rows, y):
    nr, nc = len(rows), max(len(r) for r in rows)
    h = h_table(rows)
    gt = slide.shapes.add_table(nr, nc, MARGIN, y, BODY_W, h).table
    for ci in range(nc):
        for ri, row in enumerate(rows):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            txt = row[ci] if ci < len(row) else ""
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = C_HDRBG
                col = C_HDRTX; bold = True
            else:
                cell.fill.fore_color.rgb = C_ROW1 if ri % 2 else C_ROW0
                col = C_TEXT; bold = txt.startswith("**")
            para = cell.text_frame.paragraphs[0]
            add_runs(para, txt, 9, col if ri == 0 else C_TEXT, base_bold=(ri == 0))
    return h

def render_code(slide, codelines, y):
    h = h_code(codelines)
    box = slide.shapes.add_shape(1, MARGIN, y, BODY_W, h)
    box.fill.solid(); box.fill.fore_color.rgb = C_CODEBG
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); box.line.width = Pt(0.5)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.06)
    for j, cl in enumerate(codelines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = cl
        r.font.name = MONO; r.font.size = Pt(9); r.font.color.rgb = C_TEXT
    return h

def render_visual(slide, text, y):
    h = h_visual(text)
    box = slide.shapes.add_shape(1, MARGIN, y, BODY_W, h)
    box.fill.solid(); box.fill.fore_color.rgb = C_VISBG
    box.line.color.rgb = RGBColor(0xE0, 0xC0, 0x70); box.line.width = Pt(0.75)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "🎨 시각자료  "; r.font.bold = True
    r.font.size = Pt(11); r.font.color.rgb = C_VIS; r.font.name = FONT
    add_runs(p, text, 11, C_VIS)
    return h

def render_bullet(slide, level, text, y):
    h = h_bullet(level, text)
    x = MARGIN + Inches(0.3 * level)
    w = BODY_W - Inches(0.3 * level)
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    bullet = ["● ", "– ", "· ", "· "][level]
    size = 15 if level == 0 else (13 if level == 1 else 12)
    col = C_TEXT if level == 0 else C_SUB
    rb = p.add_run(); rb.text = bullet
    rb.font.size = Pt(size); rb.font.color.rgb = C_ACCENT if level == 0 else C_SUB
    rb.font.name = FONT
    add_runs(p, text, size, col)
    return h


def figure_slide(prs, title, names):
    from PIL import Image as PILImage
    paths = [IMG / f"{n}.png" for n in names]
    paths = [p for p in paths if p.exists()]
    if not paths:
        return
    slide = blank(prs)
    add_title(slide, title + " — 시각자료")
    area_top = 1.35
    area_h = 5.6
    if len(paths) == 1:
        boxes = [(MARGIN / 914400.0, area_top, BODY_W_IN, area_h)]
    else:
        half = (BODY_W_IN - 0.3) / 2
        x0 = MARGIN / 914400.0
        boxes = [(x0, area_top, half, area_h), (x0 + half + 0.3, area_top, half, area_h)]
    for (bx, by, bw, bh), p in zip(boxes, paths):
        iw, ih = PILImage.open(p).size
        ar = iw / ih
        w = bw; h = w / ar
        if h > bh:
            h = bh; w = h * ar
        left = Inches(bx + (bw - w) / 2)
        top = Inches(by + (bh - h) / 2)
        slide.shapes.add_picture(str(p), left, top, width=Inches(w))


def build():
    md = MD.read_text()
    items = parse(md)
    prs = new_prs()

    # 표지: 첫 '슬라이드 1'의 내용에서 제목/부제 추출
    title_txt = "Aria Glass 실시간 3DGS에서 Floater 제거하기"
    sub_txt = "3D Carve Loss · 평가 지표 재정비 · Init 엔지니어링 · 교차 장면 일반화 · 전자동 파이프라인"
    title_slide(prs, title_txt, sub_txt)

    notes_buf = []
    for it in items:
        if it[0] == "section":
            section_slide(prs, it[1])
            continue
        _, title, blocks = it
        # 표지 슬라이드(슬라이드 1)는 이미 표지로 대체 → 건너뜀
        if title.startswith("슬라이드 1 —") and "표지" in title:
            continue

        # 플로우 배치 (넘치면 계속 슬라이드)
        idx = 0
        cont = False
        while idx < len(blocks) or idx == 0:
            slide = blank(prs)
            add_title(slide, title + ("  (계속)" if cont else ""))
            y = BODY_TOP
            notes = []
            any_block = False
            while idx < len(blocks):
                b = blocks[idx]
                if b[0] == "note":
                    notes.append(b[1]); idx += 1; continue
                if b[0] == "bullet":
                    bh = h_bullet(b[1], b[2]); render = lambda s, yy, b=b: render_bullet(s, b[1], b[2], yy)
                elif b[0] == "visual":
                    bh = h_visual(b[1]); render = lambda s, yy, b=b: render_visual(s, b[1], yy)
                elif b[0] == "table":
                    bh = h_table(b[1]); render = lambda s, yy, b=b: render_table(s, b[1], yy)
                elif b[0] == "code":
                    bh = h_code(b[1]); render = lambda s, yy, b=b: render_code(s, b[1], yy)
                else:
                    idx += 1; continue
                if y + bh > BODY_BOTTOM and any_block:
                    break  # 다음 슬라이드로
                render(slide, y)
                y = y + bh + Inches(0.06)
                any_block = True
                idx += 1
            # 노트 기록
            if notes:
                ns = slide.notes_slide
                ns.notes_text_frame.text = "\n".join(notes)
            cont = True
            if idx >= len(blocks):
                break

        # 이 슬라이드 제목에 매핑된 시각자료가 있으면 그림 슬라이드 추가
        m = re.match(r"^(슬라이드\s+[\w\-–]+)", title)
        if m and m.group(1) in IMG_MAP:
            figure_slide(prs, title, IMG_MAP[m.group(1)])

    prs.save(OUT)
    print(f"[saved] {OUT}  (슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장)")


if __name__ == "__main__":
    build()
