#!/usr/bin/env python3
"""build_ppt.py (ppt0812) — exp63 예산 경쟁 구조 분석 + 다음 4방향 사전 브리핑, 단독 덱.

ppt0720(다크 테마, GPU 타임라인)과 동일 스타일 재사용. 이 덱은 완료 보고가 아니라
"다음에 뭘 할지" 방향 설정용 — 지금까지 실측된 증거 + 계측 인프라 현황 + 4방향 계획.

실행: python build_ppt.py
출력: vigs_budget_briefing_0812.pptx (PDF 변환은 셸에서 soffice로 별도)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = Path(__file__).parent
IMG = HERE / "img"
OUT = HERE / "vigs_budget_briefing_0812.pptx"

FONT = "Noto Sans CJK KR"
MONO = "Noto Sans Mono CJK KR"

# ── 팔레트 (make_figures.py와 동일) ──────────────────────────────────
BG = RGBColor(0x12, 0x15, 0x1c)
PANEL = RGBColor(0x1a, 0x1f, 0x2b)
GRID = RGBColor(0x2a, 0x32, 0x42)
TEXT = RGBColor(0xe8, 0xec, 0xf4)
MUTED = RGBColor(0x8b, 0x95, 0xab)
BUDGET = RGBColor(0xff, 0xb4, 0x54)
BLUE = RGBColor(0x4c, 0x86, 0xea)
CORAL = RGBColor(0xef, 0x6a, 0x5c)
GREEN = RGBColor(0x5f, 0xd6, 0x8a)
PURPLE = RGBColor(0xb4, 0x8c, 0xf0)

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)
BODY_W = SW - 2 * MARGIN


def new_prs():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return bg


def add_text(slide, x, y, w, h, text, size, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return box


def eyebrow_title(slide, eyebrow, title):
    add_text(slide, MARGIN, Inches(0.35), BODY_W, Inches(0.32), eyebrow, 12, BUDGET, bold=True)
    add_text(slide, MARGIN, Inches(0.62), BODY_W, Inches(0.6), title, 22, TEXT, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.18), BODY_W, Pt(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GRID
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_image_fit(slide, path, x, y, w, h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    ww, hh = w, w / ar
    if hh > h:
        hh = h
        ww = h * ar
    left = x + (w - ww) / 2
    top = y + (h - hh) / 2
    slide.shapes.add_picture(str(path), left, top, width=ww)


def qa_card(slide, x, y, w, h, q, a):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = GRID
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.16)
    tf.margin_bottom = Inches(0.16)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = q
    r.font.size = Pt(13.5); r.font.bold = True; r.font.color.rgb = TEXT; r.font.name = FONT
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    p2.line_spacing = 1.18
    r2 = p2.add_run(); r2.text = a
    r2.font.size = Pt(11.5); r2.font.color.rgb = MUTED; r2.font.name = FONT


def stat_pill(slide, x, y, w, h, label, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.08
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = color
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.16)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = value
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = color; r.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(10.5); r2.font.color.rgb = MUTED; r2.font.name = FONT


def action_row(slide, x, y, w, tag, tag_color, text):
    tagbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.75), Inches(0.32))
    tagbox.adjustments[0] = 0.5
    tagbox.fill.solid(); tagbox.fill.fore_color.rgb = PANEL
    tagbox.line.color.rgb = tag_color; tagbox.line.width = Pt(1)
    tagbox.shadow.inherit = False
    tf = tagbox.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag
    r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = tag_color; r.font.name = FONT
    add_text(slide, x + Inches(1.95), y - Inches(0.03), w - Inches(1.95), Inches(0.5), text,
              12.5, TEXT, anchor=MSO_ANCHOR.MIDDLE)


def build():
    prs = new_prs()

    # ── Slide 1: Title ──────────────────────────────────────────
    s = blank(prs); set_bg(s)
    add_text(s, Inches(1.0), Inches(2.35), Inches(11.3), Inches(0.4),
              "EXP63 · GS_FLOATERLAB · 방향 설정 브리핑", 13, BUDGET, bold=True)
    add_text(s, Inches(1.0), Inches(2.75), Inches(11.3), Inches(1.6),
              "예산 경쟁 구조와 다음 4방향", 36, TEXT, bold=True)
    add_text(s, Inches(1.0), Inches(3.65), Inches(11.3), Inches(1.0),
              "tracking·PGBA·map()·background_polish가 하나의 GPU를 어떻게 나눠 쓰는가\n"
              "— 왜 같은 파라미터가 어떤 장면은 살리고 어떤 장면은 죽이는가",
              18, MUTED, line_spacing=1.3)
    add_text(s, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.5),
              "exp63 축A2 12F 회귀 발견 · 계측 인프라 현황 · 다음 실험 로드맵 · 2026-08-12",
              11.5, MUTED)

    # ── Slide 2: 핵심 질문 ──────────────────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "WHY THIS BRIEFING", "지금까지 봐온 패턴 — 결과만 알고 원인은 몰랐다")
    cy = Inches(1.5)
    ch = Inches(1.6)
    gap = Inches(0.2)
    qa_card(s, MARGIN, cy, BODY_W, ch,
            "지금까지 방식",
            "파라미터를 바꾸고 → dB가 오르내리는 걸 보고 → \"이게 좋은가 보다\"로 채택/기각. "
            "축A2(thresh 2.6·iters1=2)가 정확히 이렇게 채택됐다 — 305에서 +7.0dB 보고 좋다고 "
            "판단, 근데 12F는 한 번도 검증 안 하고 방치했다가 나중에 -4.8dB 회귀를 발견했다.")
    cy2 = cy + ch + gap
    qa_card(s, MARGIN, cy2, BODY_W, ch,
            "이번에 요구되는 방식",
            "각 함수(motion_filter/frontend/PGBA/map()/background_polish)가 실제로 몇 번 도는지, "
            "1회당 몇 ms인지, 그 시간이 어떤 파라미터에 어떻게 스케일되는지를 먼저 재고, "
            "그 지식으로 \"왜 이 장면은 오르고 저 장면은 내렸는지\"를 설명할 수 있어야 다음 "
            "파라미터를 결정한다.")

    # ── Slide 3: 아키텍처 — 스레드/자원 경쟁 구조 ────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "ARCHITECTURE", "두 스레드가 하나의 GPU를 나눠 쓰는 구조")
    add_image_fit(s, IMG / "fig_thread_arch.png", MARGIN, Inches(1.4), BODY_W, Inches(4.55))
    add_text(s, MARGIN, Inches(6.15), BODY_W, Inches(1.0),
              "tracking 스레드(motion_filter→frontend→PGBA)는 동기·고정 반복(iters1)으로 돈다 — "
              "시간 예산 캡이 코드 어디에도 없다. gs_worker 스레드(map()/background_polish)는 "
              "큐가 빌 때만 polish를 시도하는 idle-gated 구조라, tracking이 바빠질수록 polish 몫이 "
              "일방적으로 줄어든다. 이 비대칭이 오늘 발견한 회귀의 근본 구조다.",
              12.5, MUTED, line_spacing=1.25)

    # ── Slide 4: 실측 증거 — A2가 12F를 몰래 죽였다 ──────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "EVIDENCE #1", "같은 파라미터, 정반대 결과 — 실측")
    add_image_fit(s, IMG / "fig_a2_scene_flip.png", MARGIN, Inches(1.4), BODY_W, Inches(4.15))
    stats_y = Inches(5.75)
    stat_pill(s, MARGIN, stats_y, Inches(3.9), Inches(1.1), "305 (원래 부족)", "+7.00 dB", GREEN)
    stat_pill(s, MARGIN + Inches(4.1), stats_y, Inches(3.9), Inches(1.1), "12F (원래 충분)", "-4.84 dB", CORAL)
    stat_pill(s, MARGIN + Inches(8.2), stats_y, Inches(4.0), Inches(1.1), "검증 상태", "12F 미검증 방치", BUDGET)

    # ── Slide 5: 실측 증거 — polish 붕괴 ──────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "EVIDENCE #2", "12F: 왜 떨어졌나 — background_polish가 굶었다")
    add_image_fit(s, IMG / "fig_polish_collapse.png", MARGIN, Inches(1.4), BODY_W, Inches(4.3))
    add_text(s, MARGIN, Inches(5.9), BODY_W, Inches(1.1),
              "A2가 keyframe을 298→352개로 촘촘히 뽑으면서 tracking/mapping이 더 오래 바빠졌고, "
              "그만큼 background_polish가 돌 idle 시간이 사라졌다(6,586→772 step, -88%). "
              "12F는 원래(sparse 설정) 이미 keyframe 밀도가 충분한 장면(58.7개/프레임, "
              "305는 29.7개/프레임)이라 — 촘촘하게 뽑는 \"이득\"은 거의 없이 \"polish 손해\"만 봤다.",
              12.5, MUTED, line_spacing=1.25)

    # ── Slide 6: 개념 정리 — 이득 vs 손해 ─────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "THE MENTAL MODEL", "같은 다이얼, 장면마다 반대로 작용하는 이유")
    add_image_fit(s, IMG / "fig_tradeoff_concept.png", MARGIN, Inches(1.4), BODY_W, Inches(4.05))
    add_text(s, MARGIN, Inches(5.75), BODY_W, Inches(1.15),
              "※ 개념도 — 12F의 \"이득 0\"은 정밀 분해가 아니라 관측된 순효과(-4.84dB)를 "
              "설명하는 단순화된 예시다. 실제 이득/손해를 분리해서 재려면 아래 계측 정비가 "
              "선행돼야 한다. 핵심은 \"keyframe 밀도\"라는 하나의 다이얼이 두 개의 서로 다른 "
              "효과(커버리지 개선 vs polish 굶주림)를 동시에 만든다는 것.",
              12, MUTED, line_spacing=1.25)

    # ── Slide 7: 계측 인프라 현황 (+ 완료 표시) ──────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "INSTRUMENTATION · DONE", "있는 계측 vs 없는 계측 — 빈 구멍은 메웠다")
    add_image_fit(s, IMG / "fig_instrumentation.png", MARGIN, Inches(1.4), BODY_W, Inches(4.3))
    add_text(s, MARGIN, Inches(5.9), BODY_W, Inches(1.1),
              "`VIGS_TIMING_LOG`는 이미 있었지만 이번 세션 내내 한 번도 안 켠 상태였다 — 지금 켜서 "
              "4개 런(12F pre-A2/A2, 305 A2, 12F A2 재검증)을 재분석했다. 유일한 빈 구멍이었던 "
              "`background_polish_step` 내부에도 `_Sect` 패턴으로 계측을 추가(`gs_backend.py:1607,2099`"
              " 근처, n_gauss/scope/batch_size 컨텍스트 포함)해서 이제 완전히 커버된다.",
              12.5, GREEN, line_spacing=1.25)

    # ── Slide 7b: 실측 결과 — frontend 비용 스케일링 ──────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "MEASURED · STEP 1", "실측 결과 ① — frontend가 진짜 지배적이다")
    add_image_fit(s, IMG / "fig_frontend_scaling.png", MARGIN, Inches(1.4), BODY_W, Inches(4.15))
    stats_y = Inches(5.75)
    stat_pill(s, MARGIN, stats_y, Inches(3.9), Inches(1.1), "12F A2 frontend 비중", "73.4%", CORAL)
    stat_pill(s, MARGIN + Inches(4.1), stats_y, Inches(3.9), Inches(1.1), "iters1 1→2 배율", "×2.17", BUDGET)
    stat_pill(s, MARGIN + Inches(8.2), stats_y, Inches(4.0), Inches(1.1), "gs_mapping(디스패치) 비중", "1.2~1.7%", MUTED)

    # ── Slide 7c: 실측 결과 — polish 기회 vs 비용 분리 ────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "MEASURED · STEP 2", "실측 결과 ② — polish는 \"느려진 게\" 아니라 \"기회가 없었다\"")
    add_image_fit(s, IMG / "fig_polish_opportunity.png", MARGIN, Inches(1.4), BODY_W, Inches(4.15))
    add_text(s, MARGIN, Inches(5.75), BODY_W, Inches(1.1),
              "핵심 확인: call당 비용(3.9~6.4ms)은 gaussian 수(40k~150k)에 걸쳐 거의 고정 — 미세한 "
              "n_gauss 의존성은 있지만(약 1.5배 범위) 기회 횟수 차이(13배)에 비하면 훨씬 작다. "
              "즉 exp카드에 남겼던 \"772 step이 idle 부족 때문인지 step 자체가 느려서인지\" 질문의 "
              "답은 확정적으로 idle 부족 쪽이다.",
              12.5, TEXT, line_spacing=1.25)

    # ── Slide 7d: 실측 결과 — map() 내부 구성 ─────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "MEASURED · BONUS", "실측 결과 ③ — map() 자체는 두 씬 모두 같은 구성")
    add_image_fit(s, IMG / "fig_map_breakdown.png", MARGIN, Inches(1.4), BODY_W, Inches(4.05))
    add_text(s, MARGIN, Inches(5.75), BODY_W, Inches(1.1),
              "loss_compute+backward가 두 씬 모두 ~90%로 지배적, rasterize는 ~7~8%뿐 — map() 자체의 "
              "내부 구성은 씬에 따라 거의 안 변한다. 즉 12F가 나빠진 원인은 map() 내부가 아니라 "
              "map()이 dispatch되기 전 단계(frontend가 얼마나 자주·오래 도는가)에 있다는 게 다시 확인된다.",
              12.5, MUTED, line_spacing=1.25)

    # ── Slide 7e: 실측 결과 — replay_time_scale 스윕 (③ 직접 증거) ──
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "MEASURED · STEP 3", "실측 결과 ④ — 예산을 33% 늘리면 회귀가 거의 회복된다 (역-U자형)")
    add_image_fit(s, IMG / "fig_scale_sweep.png", MARGIN, Inches(1.4), BODY_W, Inches(4.15))
    stats_y = Inches(5.75)
    stat_pill(s, MARGIN, stats_y, Inches(3.9), Inches(1.1), "scale 1.5→2.0", "+4.26 dB", GREEN)
    stat_pill(s, MARGIN + Inches(4.1), stats_y, Inches(3.9), Inches(1.1), "scale 2.0→3.0", "-2.13 dB", CORAL)
    stat_pill(s, MARGIN + Inches(8.2), stats_y, Inches(4.0), Inches(1.1), "polish 3.0 상태", "10,000 캡 도달", BUDGET)

    # ── Slide 7f: 실측 결과 — polish "양자화" 메커니즘 규명 ─────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "MEASURED · STEP 4", "실측 결과 ⑤ — polish 횟수가 \"계단식\"으로 보였던 진짜 이유")
    add_image_fit(s, IMG / "fig_gap_quantization.png", MARGIN, Inches(1.4), BODY_W, Inches(4.15))
    add_text(s, MARGIN, Inches(5.85), BODY_W, Inches(1.1),
              "gs_worker_dispatch/background_polish_call에 epoch 타임스탬프를 추가해 디스패치 사이 "
              "gap을 직접 재구성했다. 총 polish 횟수는 각 gap 길이를 step당 비용(4~6ms)으로 나눈 정수 "
              "몫의 합 — gap의 43~57%는 아예 1개도 못 낄 만큼 짧아 \"부드러운 비례\"가 아니라 본질적으로 "
              "계단식이다. scale이 커질수록 긴 gap(최대 526→1582→2478ms)이 늘어나며 총량이 늘지만, "
              "3.0의 역행은 이 메커니즘만으론 설명 안 됨 — 좁은 후보 풀(stride-5 rgb_dense) 과적합 "
              "가설 필요(미검증).",
              12, MUTED, line_spacing=1.25)

    # ── Slide 8: 4개 방향 ─────────────────────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "FOUR DIRECTIONS", "다음에 정해야 할 4가지 — 전부 같은 문제의 다른 얼굴")
    ay = Inches(1.55)
    rows = [
        ("① FREEZE 기준", BUDGET,
         "`mapping_freeze_after_frac` 고정 비율(0.6140)은 1253/305 튜닝값일 뿐 원리적 근거 없음 "
         "— 12F 회귀가 직접 반례. coverage/densify-rate/loss-plateau 중 robust한 기준 필요."),
        ("② CARVE LOSS 이식", BLUE,
         "배치에서 검증된 방법론(exp44d2, 33.8dB)을 density/pruning 결정에 이식. 단, 12F 재튜닝 "
         "전이라 3장면 다 27dB 안정화 전 — CLAUDE.md 순서 제약과 맞물림."),
        ("③ 예산 robust화", CORAL,
         "tracking에 시간 캡을 걸 것인가, polish에 최소 하한을 보장할 것인가 — 1253+305+12F "
         "세 장면 교차검증 필수(A2가 두 장면만 보고 실패한 패턴 반복 금지)."),
        ("④ FREEZE 스케줄", PURPLE,
         "①과 본질적으로 같은 문제. iteration/frame count가 아닌 기준(coverage 신호 등)으로 "
         "대체 — 지금 죽어있는 `transmittance` 계산이 후보 재료."),
    ]
    for i, (tag, color, text) in enumerate(rows):
        action_row(s, MARGIN, ay + Inches(1.05) * i, BODY_W, tag, color, text)
    add_text(s, MARGIN, ay + Inches(1.05) * 4 + Inches(0.15), BODY_W, Inches(0.6),
              "①④(freeze 기준)와 ③(예산 배분)은 사실 하나의 문제 — \"제한된 실시간 예산을 구조 "
              "생성과 정제에 어떻게 나누는가\". ②(carve)는 그 정제 단계의 품질을 결정한다.",
              12, TEXT, bold=True)

    # ── Slide 9: 로드맵 진행 상황 ───────────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "ROADMAP · PROGRESS", "계획했던 4단계 중 어디까지 왔나")
    add_image_fit(s, IMG / "fig_next_steps.png", MARGIN, Inches(1.5), BODY_W, Inches(2.9))
    ay = Inches(4.7)
    rows2 = [
        ("1 완료", GREEN, "A2 전/후 × 12F/305 4개 조합을 `VIGS_TIMING_LOG` 켜서 재실행·재분석 완료 — "
                        "frontend 73.4%/64.0%, polish 5,683/822/9,866회 등 실측 확보."),
        ("2 완료", GREEN, "`background_polish_step`에 `_Sect` 계측 추가 완료 — call당 3.9~6.4ms, "
                        "n_gauss 의존성 약함(주 원인 아님) 확정."),
        ("3 보류", MUTED, "시간축 병합 타임라인 재구성 — 이번엔 phase별 총합/스케일링으로 핵심 질문에 "
                        "이미 답이 나와서 우선순위 낮춤. `get_lock()` 경합의 직접 증거는 아직 없음."),
        ("4 완료", GREEN, "파라미터 스윕 — iters1 1↔2 한 쌍(×2.17) + `replay_time_scale` 1.5/2.0/3.0 "
                        "3점 스윕(23.84/28.10/25.97dB, 역-U자형) 확보. gap 기반 양자화 메커니즘도 "
                        "규명 완료. 3.0 역행의 근본원인(과적합 가설)만 미검증으로 남음."),
    ]
    for i, (tag, color, text) in enumerate(rows2):
        action_row(s, MARGIN, ay + Inches(0.58) * i, BODY_W, f"STEP {tag}", color, text)

    # ── Slide 10: 마무리 ────────────────────────────────────────
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "CONCLUSIONS", "이번 조사로 확정된 것 · 다음 단계")
    cy = Inches(1.6)
    ch = Inches(1.5)
    gap = Inches(0.22)
    qa_card(s, MARGIN, cy, BODY_W, ch,
            "확정된 것 (실측 완료)",
            "① frontend가 프레임 단계 시간의 64~73%로 압도적 지배 — 12F는 305보다 절대시간 2배. "
            "② background_polish는 \"느려진 게\" 아니라 \"기회가 없었다\"(call당 비용 거의 고정, "
            "기회 횟수만 13배 차이) — exp카드 미해결 질문 해소. ③ map() 내부 구성(loss/backward "
            "~90%)은 두 씬 모두 동일 — 문제는 map() 안이 아니라 그 앞 단계에 있다. ④ 예산을 33% "
            "늘리면(scale 1.5→2.0) 12F 회귀가 거의 회복된다(+4.26dB) — 예산 부족이 인과관계임을 "
            "직접 실험으로 확인. ⑤ polish 횟수의 \"계단식\" 양상은 gap÷step비용의 정수 몫 합이라는 "
            "구조적 필연 — 신비가 아니라 산수.")
    cy2 = cy + ch + gap
    qa_card(s, MARGIN, cy2, BODY_W, ch,
            "아직 안 밝혀진 것",
            "iters1 1→2가 왜 정확히 2.17배인지(keyframe 수 증가와 순수 iters 효과가 섞여있어 분리 "
            "안 됨) — 3장면 동시 통제 스윕이 필요. `get_lock()` GPU 경합이 실측 병목인지도 직접 "
            "증거는 아직 없음(현재는 순차 phase 합산 기준). scale 3.0의 역행(28.10→25.97dB, "
            "polish는 오히려 최대인데도 하락)은 좁은 후보 풀 과적합 가설뿐 — held-out 분리 재검증 "
            "필요.")
    cy3 = cy2 + ch + gap
    qa_card(s, MARGIN, cy3, BODY_W, ch,
            "4방향에 대한 함의",
            "①④(freeze 기준)와 ③(예산 robust화) 설계는 이제 \"frontend를 줄이거나 polish에 최소 "
            "몫을 직접 보장\" 둘 중 하나로 좁혀졌다 — 감이 아니라 실측 근거로. scale 스윕은 ③이 "
            "실제로 유효한 레버임을 증명했지만, 무작정 늘리는 게 답이 아니라(3.0 역행) \"적정 하한 "
            "보장\" 설계가 필요함도 같이 보여줬다. ②(carve)는 여전히 12F 재튜닝 이후로 순서 유지.")

    prs.save(OUT)
    print(f"[saved] {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
