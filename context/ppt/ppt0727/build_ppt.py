#!/usr/bin/env python3
"""build_ppt.py (ppt0727) — VIGS-SLAM 실시간화 여정 exp52->exp56, 단독 덱.

ppt0720/build_ppt.py와 동일한 다크 테마·헬퍼 함수 패턴을 재사용(코드 보존 지침).
md 파서 재사용 안 함 — ppt_outline_20260727.md 19장 구성을 이 파일에서 직접 조립.

실행: python build_ppt.py
출력: vigs_realtime_journey_0727.pptx  (PDF 변환은 셸에서 soffice로 별도)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = Path(__file__).parent
IMG = HERE / "img"
OUT = HERE / "vigs_realtime_journey_0727.pptx"

FONT = "Noto Sans CJK KR"
MONO = "Noto Sans Mono CJK KR"

# ── 팔레트 (make_figures.py와 동일) ──────────────────────────────────
BG = RGBColor(0x12, 0x15, 0x1c)
PANEL = RGBColor(0x1a, 0x1f, 0x2b)
GRID = RGBColor(0x2a, 0x32, 0x42)
TEXT = RGBColor(0xe8, 0xec, 0xf4)
MUTED = RGBColor(0x8b, 0x95, 0xab)
BUDGET = RGBColor(0xff, 0xb4, 0x54)
BLUE = RGBColor(0x4c, 0x86, 0xea)
CORAL = RGBColor(0xef, 0x6a, 0x5c)
GREEN = RGBColor(0x4c, 0xc9, 0x8a)
PURPLE = RGBColor(0xa4, 0x79, 0xe0)

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)
BODY_W = SW - 2 * MARGIN


def new_prs():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return bg


def add_text(slide, x, y, w, h, text, size, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return box


def eyebrow_title(slide, eyebrow, title):
    add_text(slide, MARGIN, Inches(0.35), BODY_W, Inches(0.32), eyebrow, 12, BUDGET, bold=True)
    add_text(slide, MARGIN, Inches(0.62), BODY_W, Inches(0.6), title, 22, TEXT, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.18), BODY_W, Pt(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GRID
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_image_fit(slide, path, x, y, w, h):
    iw, ih = Image.open(path).size
    ar = iw / ih
    ww, hh = w, w / ar
    if hh > h:
        hh = h
        ww = h * ar
    left = x + (w - ww) / 2
    top = y + (h - hh) / 2
    slide.shapes.add_picture(str(path), left, top, width=ww)


def bullet_block(slide, x, y, w, h, items, size=13.5, gap=Pt(8)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        p.line_spacing = 1.22
        r = p.add_run()
        r.text = f"•  {it}"
        r.font.size = Pt(size)
        r.font.color.rgb = TEXT
        r.font.name = FONT
    return box


def note(slide, x, y, w, h, text):
    add_text(slide, x, y, w, h, text, 11, MUTED, line_spacing=1.2)


def qa_card(slide, x, y, w, h, q, a):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = GRID
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.16)
    tf.margin_bottom = Inches(0.16)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = q
    r.font.size = Pt(13.5); r.font.bold = True; r.font.color.rgb = TEXT; r.font.name = FONT
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    p2.line_spacing = 1.18
    r2 = p2.add_run(); r2.text = a
    r2.font.size = Pt(11.5); r2.font.color.rgb = MUTED; r2.font.name = FONT


def stat_pill(slide, x, y, w, h, label, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.08
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = color
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.16)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = value
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = color; r.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(10.5); r2.font.color.rgb = MUTED; r2.font.name = FONT


def action_row(slide, x, y, w, tag, tag_color, text):
    tagbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.7), Inches(0.32))
    tagbox.adjustments[0] = 0.5
    tagbox.fill.solid(); tagbox.fill.fore_color.rgb = PANEL
    tagbox.line.color.rgb = tag_color; tagbox.line.width = Pt(1)
    tagbox.shadow.inherit = False
    tf = tagbox.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag
    r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = tag_color; r.font.name = FONT
    add_text(slide, x + Inches(1.85), y - Inches(0.06), w - Inches(1.85), Inches(0.55), text,
              12, TEXT, anchor=MSO_ANCHOR.MIDDLE)


def formula_card(slide, x, y, w, h, metric, formula_lines, r2, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.05
    card.fill.solid(); card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = color; card.line.width = Pt(1.25)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.14)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"{metric}   (R² = {r2:.3f})"
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = color; r.font.name = FONT
    for line in formula_lines:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(7)
        p2.line_spacing = 1.25
        r2run = p2.add_run(); r2run.text = line
        r2run.font.size = Pt(13); r2run.font.color.rgb = TEXT; r2run.font.name = MONO


def simple_table(slide, x, y, w, rows, col_widths, header=True, row_h=Inches(0.5)):
    n_rows = len(rows)
    n_cols = len(rows[0])
    h = row_h * n_rows
    gshape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = gshape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if (ri == 0 and header) else BG
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(12.5 if (ri == 0 and header) else 12)
                    r.font.bold = (ri == 0 and header) or ci == 0
                    r.font.color.rgb = BUDGET if (ri == 0 and header) else TEXT
                    r.font.name = FONT
    return table


def build():
    prs = new_prs()

    # ── §0 도입 ──────────────────────────────────────────────
    # 슬라이드 1 — 표지
    s = blank(prs); set_bg(s)
    add_text(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(0.4),
              "GS_FLOATERLAB · EXP52 → EXP56", 13, BUDGET, bold=True)
    add_text(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.7),
              "Incremental 3DGS Mapping 실시간화", 34, TEXT, bold=True)
    add_text(s, Inches(1.0), Inches(3.55), Inches(11.3), Inches(0.9),
              "VIGS-SLAM 채택 리캡 · 실시간 최초 돌파 · 품질까지 함께 개선\n\"왜 안 빨라지나\" 심층 규명",
              17, MUTED, line_spacing=1.3)
    add_text(s, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.5),
              "1253 시퀀스 · exp52~56 · 07/20 → 07/27", 11.5, MUTED)

    # 슬라이드 2 — 07/20 리캡
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "RECAP · 07/20", "지난 발표 리캡: 어디까지 끝냈나")
    bullet_block(s, MARGIN, Inches(1.45), BODY_W, Inches(3.4), [
        "VIGS-SLAM(ECCV2026) 채택 — 단안 RGB+IMU, DROID-SLAM식 dense correlation 트래킹. 우리 목표(흑백 SLAM 트래킹 + RGB incremental 매핑)에 가장 가까운 참조 아키텍처",
        "20초 타이밍 버그 발견·수정 — reader 프로세스 sleep(20) 순서 문제로 모든 \"온라인 루프 총합\"에 인위적 20초가 섞여있었음. 수정 후 실시간 배수 2.77배 → 1.52배로 정정",
        "구조적 개선 — IMU 프리적분 C++ 빌드 + TensorRT 3종(−14%), tracking/mapping 비동기 오버랩 gs_parallel 도입(업스트림 레이스 컨디션 버그 발견·수정, −26.1% 추가)",
        "핵심 진단 — mapping(rasterize+backward+loss_compute)이 온라인 루프의 81.4%로 압도적 병목. 이번 사이클(exp53~56) 전체가 이 진단 위에서 진행됨",
    ], size=13.5)
    note(s, MARGIN, Inches(6.7), BODY_W, Inches(0.5),
         "오늘은 이 지점(1.52배, 아직 실시간 아님) 이후 일주일간 일어난 일만 다룹니다.")

    # 슬라이드 3 — TL;DR
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "SUMMARY", "이번 사이클 요약 (TL;DR)")
    add_image_fit(s, IMG / "fig_tldr_ratio.png", MARGIN, Inches(1.4), Inches(6.4), Inches(4.6))
    bx = MARGIN + Inches(6.7)
    bw = BODY_W - Inches(6.7)
    bullet_block(s, bx, Inches(1.5), bw, Inches(4.6), [
        "exp53+54: 트래킹·매핑 양쪽 경량화 → 실시간 최초 돌파(0.94배)",
        "exp55: 내용-적응 gaussian 예산 + carve loss 온라인 이식 → 시간 유지하며 평균 gaussian −35.9%, 가시 floater −7.5%",
        "exp56: \"왜 안 빨라지나\"를 회귀분석으로 규명 → 숨어있던 진짜 병목 2건 발견·수정 → 0.70배, PSNR까지 +0.88/+0.93dB 개선",
    ], size=13)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(5.9), bw, Inches(1.15))
    box.adjustments[0] = 0.08
    box.fill.solid(); box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = BUDGET; box.line.width = Pt(1.25)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "1.52배(비실시간) → 0.70배(여유 19초), 그 사이 품질은 계속 좋아짐"
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = TEXT; r.font.name = FONT

    # ── §1 실시간 최초 돌파 ──────────────────────────────────
    # 슬라이드 4 — exp53
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§1 · EXP53", "프론트엔드(트래킹) 경량화")
    add_image_fit(s, IMG / "fig_exp53_axes.png", MARGIN, Inches(1.4), BODY_W, Inches(3.6))
    bullet_block(s, MARGIN, Inches(5.15), BODY_W, Inches(1.6), [
        "07/20 진단(mapping 81.4%)과 별개로 트래킹도 손댈 여지가 컸음 — 나중 exp56 부록에서 \"순수 트래킹 절감폭(−47.5%)이 매핑 절감폭(−12~15%)보다 훨씬 크다\"는 게 재확인됨",
        "축D(correlation 해상도)는 사전학습 가중치에 shape가 고정 결합돼 조사 후 구현 불가 판정 — 재학습 없인 못 건드림, 정직하게 기각 기록",
    ], size=12)
    note(s, MARGIN, Inches(6.75), BODY_W, Inches(0.5),
         "재학습 없이 안 되는 것도 시도해보고 명확히 접었다는 게 중요 — 막연히 미룬 게 아니라 판정을 끝냄.")

    # 슬라이드 4b (신규) — 트래킹 축이 실제로 뭘 바꾸는 건지 설명
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§1 · EXP53 · 용어", "트래킹 축A~D, 실제로 뭘 바꾸는 건가")
    add_text(s, MARGIN, Inches(1.35), BODY_W, Inches(0.5),
              "파이프라인: 카메라 프레임 → motion_filter(이게 keyframe인가?) → frontend(correlation → GRU 반복정제 → local BA)",
              12.5, MUTED)
    tk_rows = [
        ("축A", BLUE, "correlation iters — 두 프레임의 optical flow(correlation)를 GRU로 여러 번 반복해서 다듬는 횟수(iters1/iters2). 시간 −20.7%, ΔPSNR +0.21/+0.21dB(오히려 개선) — 4/2회를 1/0회로 줄여도 정확도 손실 없었음"),
        ("축B", BUDGET, "motion_filter.thresh — 새 프레임이 이전 keyframe 대비 \"이만큼은 움직여야 keyframe으로 인정\"하는 임계값. 시간 −15.4%, ΔPSNR +0.41/+0.55dB — keyframe 자체가 덜 생겨 트래킹·매핑 작업량이 동시에 줄어듦, 이 세션 최대 레버"),
        ("축C", GREEN, "frontend_window/radius — local BA가 한 번에 붙잡고 최적화하는 \"최근 keyframe 그래프\"의 크기(window)와 이웃 범위(radius). 시간 −1.7%, ΔPSNR +0.02/+0.12dB — 줄이면 BA solve 행렬이 작아져 계산이 가벼워짐"),
        ("축D", MUTED, "correlation 해상도 — correlation volume(두 프레임 유사도 맵)을 뽑는 특징맵의 해상도. 사전학습된 신경망 가중치의 입력 shape에 못박혀 있어 재학습 없인 못 건드림 — 미실행"),
    ]
    ay = Inches(2.1)
    row_h = Inches(1.1)
    for i, (tag, color, text) in enumerate(tk_rows):
        action_row(s, MARGIN, ay + row_h * i, BODY_W, tag, color, text)

    # 슬라이드 5 — exp54 (원래 막대그래프 유지, 겹침 버그만 수정)
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§1 · EXP54", "매핑 경량화 — 7축 스캔")
    add_image_fit(s, IMG / "fig_exp54_axes.png", MARGIN, Inches(1.4), BODY_W, Inches(4.6))
    note(s, MARGIN, Inches(6.35), BODY_W, Inches(0.9),
         "속도 축만 스캔한 게 아니라 축3에서 우연히 \"이 시점엔 트래킹이 91%\"라는 걸 발견 — "
         "이게 exp53 우선순위의 직접 근거가 됨(서로 다른 실험이 서로의 근거를 만든 사례). "
         "각 축이 정확히 뭘 바꾸는 실험인지는 다음 슬라이드에서 설명.")

    # 슬라이드 5b (신규) — 매핑 축이 실제로 뭘 바꾸는 건지 설명
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§1 · EXP54 · 용어", "매핑 축1~7, 실제로 뭘 바꾸는 건가")
    add_text(s, MARGIN, Inches(1.32), BODY_W, Inches(0.5),
              "파이프라인: keyframe 도착 → 초기 gaussian 샘플링(축1/2/7) → map() 반복 최적화(축3) → rasterize·densify(축4/5/6)",
              12, MUTED)
    map_rows = [
        ("축1", GREEN, "pcd_downsample — 매 keyframe마다 RGB-D에서 새 gaussian을 뽑을 때의 샘플링 간격. 64→128은 점 개수 절반. 시간 −3.3%, ΔPSNR −0.05/−0.14dB(거의 동일) — 채택"),
        ("축2", CORAL, "pcd_downsample_init — 축1과 같은 개념인데 \"시퀀스 맨 첫 keyframe\"에만 적용. 시간 +1.2%(역효과), ΔPSNR +0.24/+0.19dB — densify가 나중에 보충해버려 시간 이득이 상쇄됨, 기각"),
        ("축3", CORAL, "map() iters — keyframe 하나를 처리할 때 gaussian을 SGD로 최적화하는 반복 횟수. 시간 −1.5%, ΔPSNR −0.46/−0.48dB — 이 시점엔 트래킹이 병목(91%)이라 ROI 나쁨, 기각"),
        ("축4", PURPLE, "render_downsample — 매핑용 렌더링 자체를 낮은 해상도로 그림(원래는 트래킹 해상도의 8배까지 업샘플). 시간 −4.2%, ΔPSNR −0.83/−0.68dB — 유효하나 이미 실시간이라 손실 감수 안 함, 보류"),
        ("축5", CORAL, "max_viewpoints — map() 한 번의 호출에서 동시에 보는 카메라 뷰 개수의 상한. 시간 −0.7%(거의 무변화), ΔPSNR −1.73/−1.98dB — 최악의 ROI, 기각"),
        ("축6", CORAL, "densify_grad_threshold — gaussian이 \"부족하다\"고 판단해 새로 쪼개거나 복제하는(densify) 공격성을 정하는 임계값. 시간 +0.2%, ΔPSNR −0.07/−0.29dB — 개수는 줄여도 시간은 그대로, 기각"),
        ("축7", GREEN, "PPM 적응 샘플링 — 균일 샘플링 대신 이미지 디테일(Sobel gradient)에 따라 점을 더/덜 배분. 시간 +0.9%(오차범위), ΔPSNR +0.16/+0.11dB — 같은 예산에서 품질만 공짜로 오름, 채택"),
    ]
    ay = Inches(1.95)
    row_h = Inches(0.66)
    for i, (tag, color, text) in enumerate(map_rows):
        action_row(s, MARGIN, ay + row_h * i, BODY_W, tag, color, text)

    # 슬라이드 6 — 실시간 최초 돌파
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§1 · RESULT", "결과: 실시간 최초 돌파")
    add_image_fit(s, IMG / "fig_breakthrough.png", MARGIN, Inches(1.5), BODY_W, Inches(4.3))
    note(s, MARGIN, Inches(6.1), BODY_W, Inches(0.9),
         "exp53+exp54 통합 레시피: 61.34초, 실시간 배수 0.94배 — 이 프로젝트 최초로 1.0배 미만 달성. "
         "여기서 멈춰도 목표(실시간)는 달성인데, 다음 두 사이클(exp55/56)에서 여유를 더 벌고 품질까지 올림.")

    # ── §2 품질까지 함께 ─────────────────────────────────────
    # 슬라이드 7
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§2 · EXP55", "내용-적응 gaussian 예산")
    add_image_fit(s, IMG / "fig_exp55_adaptive.png", MARGIN, Inches(1.4), BODY_W, Inches(4.1))
    note(s, MARGIN, Inches(5.7), BODY_W, Inches(1.4),
         "가설: 디테일 많은(Sobel gradient 큰) keyframe엔 gaussian이 더 필요하고, 단조로운 keyframe엔 덜 필요할 것. "
         "Phase1(캘리브레이션)에서 상관관계로 먼저 검증한 뒤 Phase2에서 배율곡선+keyframe별 cap을 구현 → "
         "평균 gaussian 수 −35.9%, 최종 −35.3%, PSNR·궤적 손실 없음(오히려 소폭 개선). "
         "사용자 아이디어를 감으로 가지 않고 상관관계로 먼저 검증한 뒤 구현했다는 게 핵심.")

    # 슬라이드 8
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§2 · EXP55", "carve loss 온라인 이식")
    add_image_fit(s, IMG / "fig_exp55_carve.png", MARGIN, Inches(1.4), BODY_W, Inches(3.2))
    bullet_block(s, MARGIN, Inches(4.85), BODY_W, Inches(1.9), [
        "배치 트랙에서 이미 검증된(AUC 0.98) carve loss를, VIGS의 BA-정제 추적 depth(disps_up)를 신뢰 표면 삼아 depth-violation 전용 온라인 근사로 재설계",
        "region GT가 VIGS 좌표계에 안 맞아 — carve_loss.py 자신의 신호 설계를 그대로 재구현한 오프라인 진단 지표를 새로 제작해 검증",
    ], size=12)
    note(s, MARGIN, Inches(6.75), BODY_W, Inches(0.5),
         "표준 지표가 안 맞으면 새로 만들어서라도 채택/기각을 숫자로 판단한다는 원칙.")

    # 슬라이드 9 — 부록
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§2 · APPENDIX", "부록: 직렬 실행으로 밝혀진 진짜 구조")
    add_image_fit(s, IMG / "fig_exp55_appendix.png", MARGIN, Inches(1.4), BODY_W, Inches(3.7))
    bullet_block(s, MARGIN, Inches(5.3), BODY_W, Inches(1.4), [
        "\"지금 상태로 직렬 돌리면 순수 시간이 어떻게 되나\" 실측 → 순수 tracking 27.9초 vs 순수 mapping 80.1초. "
        "이게 exp54의 \"tracking-bound\" 결론과 정면 모순",
        "파고든 결과 두 가지 발견: ① GPU 경합이 병렬 tracking을 거의 2배로 부풀림 ② _gs_queue가 가득 차면 "
        "오래된 keyframe을 버리는 정책 때문에 병렬 모드가 mapping 업데이트의 약 80%를 그냥 스킵(직렬 110회 vs 병렬 22회)",
    ], size=12)
    note(s, MARGIN, Inches(6.9), BODY_W, Inches(0.5),
         "이 발견이 exp56 전체의 출발점 — \"병렬 배수만 보면 안 되고 안에서 뭐가 스킵되는지 봐야 한다\".")

    # ── §3 exp56 ─────────────────────────────────────────────
    # 슬라이드 10
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§3 · EXP56", "문제 제기: gaussian↓인데 시간은 그대로")
    add_image_fit(s, IMG / "fig_exp56_problem.png", MARGIN, Inches(1.4), Inches(6.6), Inches(4.6))
    bx = MARGIN + Inches(6.9); bw = BODY_W - Inches(6.9)
    bullet_block(s, bx, Inches(1.6), bw, Inches(4.2), [
        "exp55에서 평균 gaussian 수를 35.9% 줄였는데 순수 mapping 시간은 12.2%밖에 안 줄었음 — \"gaussian 수를 줄이는\" 접근이 한계에 도달한 신호",
        "exp54 축6+2에서도 이미 같은 현상 관측(gaussian 수를 더 눌러도 시간 그대로)",
    ], size=12.5)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(5.6), bw, Inches(1.3))
    box.adjustments[0] = 0.08
    box.fill.solid(); box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = BUDGET; box.line.width = Pt(1.25)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "질문: gaussian 수가 시간을 안 줄이면, 무엇이 시간을 지배하는가?"
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = TEXT; r.font.name = FONT

    # 슬라이드 11 — Phase 0
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§3 · PHASE 0", "기존 계측 재분석으로 원인 진단")
    add_image_fit(s, IMG / "fig_exp56_phase0.png", MARGIN, Inches(1.35), Inches(5.6), Inches(5.3))
    bx = MARGIN + Inches(5.9); bw = BODY_W - Inches(5.9)
    bullet_block(s, bx, Inches(1.8), bw, Inches(3.5), [
        "새 실험 없이 기존 타이밍 로그를 처음으로 세부 태그별 집계 — 직렬 순수 map() 68.16초 중 rasterize 40% + backward 34% + loss_compute 24%",
        "loss_compute는 순수 픽셀 연산(N-무관), rasterize/backward도 이 gaussian 규모(85k~130k)에선 \"고정비\"가 데이터量-비례 항을 압도함을 확인",
    ], size=12.5)
    note(s, bx, Inches(5.7), bw, Inches(1.0), "새 실험을 돌리기 전에 이미 있는 데이터부터 다시 봤다 — 이게 다음 발견들의 방법론.")

    # 슬라이드 12 — Phase 1
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§3 · PHASE 1", "iters 하향 — 첫 큰 개선")
    add_image_fit(s, IMG / "fig_exp56_phase1.png", MARGIN, Inches(1.4), BODY_W, Inches(4.5))
    note(s, MARGIN, Inches(6.05), BODY_W, Inches(1.2),
         "\"고정비가 지배적\"이라면 진짜 변수는 gaussian 수가 아니라 \"반복 횟수(iters)\" — keyframe당 SGD 반복 10→7→5 스캔. "
         "iters=7: 시간 −16.1%, PSNR mean/kf 둘 다 +0.21dB, map() 처리 keyframe 수도 22→26회 증가 — 전 지표 동시 개선. "
         "(참고) 반대 방향(iters를 올려보는 실험)은 오히려 coverage가 줄어 역효과.")

    # 슬라이드 13 — Phase 4
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§3 · PHASE 4", "숨어있던 진짜 병목 발견")
    add_image_fit(s, IMG / "fig_exp56_phase4.png", MARGIN, Inches(1.35), Inches(5.4), Inches(5.3))
    bx = MARGIN + Inches(5.7); bw = BODY_W - Inches(5.7)
    bullet_block(s, bx, Inches(1.8), bw, Inches(3.6), [
        "map() 호출을 세부 로그로 종류별 분해해보니 — 26회 중 단 2~3회(맵 최초 초기화 + IMU 재초기화)가 전체 mapping 시간의 49%를 차지",
        "원인: IMU 재초기화 시 맵 전체를 삭제하고 처음부터 다시 채우는데, 이때 반복 횟수가 90~131회(정규 호출의 10배 이상)",
        "init_itr_num 1050→600으로 낮춰 추가 시간 −6.2%, PSNR 사실상 무손실(300까지 내리면 실손실 확인 후 기각)",
    ], size=12)
    note(s, bx, Inches(5.8), bw, Inches(0.9), "호출 26개 중 딱 2~3개에 절반이 몰려있었다는 게 이번 사이클 최대 발견.")

    # 슬라이드 13b (신규) — Phase 5: 파라미터 vs 시간 회귀분석 피팅 실물
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§3 · PHASE 5", "회귀분석 피팅 — 무엇이 시간을 결정하는가")
    add_image_fit(s, IMG / "fig_phase5_fit.png", MARGIN, Inches(1.35), BODY_W, Inches(5.1))
    note(s, MARGIN, Inches(6.6), BODY_W, Inches(0.8),
         "왼쪽: 548개 실제 map() 호출을 iters×n_view(한 호출이 처리하는 \"view-op\" 수) 하나로 설명하면 "
         "R²=0.992 — gaussian 수(점 색깔)는 같은 x값 안에서도 시간을 거의 안 바꿈. "
         "오른쪽: 그 회귀식을 항목별로 분해하면 \"view-op당 고정비\"가 gaussian-비례항의 3배, 픽셀-비례항은 사실상 0에 가까움.")

    # 슬라이드 13c (신규) — Phase 5: 회귀식 자체 (수식)
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§3 · PHASE 5 · 수식", "실제로 피팅된 4개 관계식")
    add_text(s, MARGIN, Inches(1.3), BODY_W, Inches(0.35),
              "직렬(GPU 경합 없음) 330개 실제 map_call 로그, 최소제곱(least-squares) 피팅",
              12, MUTED)
    # 파라미터 glossary — 수식에 나오는 기호가 각각 뭔지
    gloss = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(1.68), BODY_W, Inches(1.0))
    gloss.adjustments[0] = 0.06
    gloss.fill.solid(); gloss.fill.fore_color.rgb = PANEL
    gloss.line.color.rgb = GRID; gloss.line.width = Pt(1)
    gloss.shadow.inherit = False
    tf = gloss.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
    glossary_items = [
        ("iters", "그 map() 호출에서 gaussian을 SGD로 최적화하는 반복 횟수 (정규 keyframe=7, 초기화=90~131 등)"),
        ("n_view", "그 반복 1회마다 렌더링해서 loss를 재는 카메라 뷰 개수 (window+global 곁눈질 합)"),
        ("n_gauss", "그 시점의 전체 gaussian 개수"),
        ("pixratio", "렌더 해상도 배율 = 1/render_downsample² (원해상도면 1)"),
        ("ncalls", "= iters × n_view — 한 번의 map() 호출이 처리하는 총 \"view-op\" 수 (rasterize/loss_compute는 view마다 도니까 이 단위가 자연스러움)"),
    ]
    for i, (sym, desc) in enumerate(glossary_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        r1 = p.add_run(); r1.text = f"{sym}  "
        r1.font.size = Pt(11.5); r1.font.bold = True; r1.font.color.rgb = BUDGET; r1.font.name = MONO
        r2 = p.add_run(); r2.text = desc
        r2.font.size = Pt(11); r2.font.color.rgb = TEXT; r2.font.name = FONT
    cw = (BODY_W - Inches(0.3)) / 2
    ch = Inches(1.55)
    gx = Inches(0.3)
    gy = Inches(0.22)
    fy0 = Inches(2.85)
    formula_card(s, MARGIN, fy0, cw, ch, "rasterize",
                 ["= 1.473·ncalls", "+ 0.00545·ncalls·(n_gauss/1000)", "− 0.0346·ncalls·pixratio   [ms]"],
                 0.9921, BLUE)
    formula_card(s, MARGIN + cw + gx, fy0, cw, ch, "loss_compute",
                 ["= 0.956·ncalls", "+ 0.00207·ncalls·(n_gauss/1000)", "− 0.0241·ncalls·pixratio   [ms]"],
                 0.9935, BUDGET)
    formula_card(s, MARGIN, fy0 + ch + gy, cw, ch, "backward",
                 ["= −1.623·iters", "+ 1.126·iters·n_view", "+ 0.126·iters·(n_gauss/1000)   [ms]"],
                 0.9294, PURPLE)
    formula_card(s, MARGIN + cw + gx, fy0 + ch + gy, cw, ch, "optimizer_step",
                 ["= 1.739·iters", "− 0.0946·iters·n_view", "+ 0.00093·iters·(n_gauss/1000)   [ms]"],
                 0.9975, GREEN)
    note(s, MARGIN, fy0 + 2 * ch + gy + Inches(0.15), BODY_W, Inches(0.6),
         "네 식 전부 \"iters·n_view\" 항의 계수가 gaussian/픽셀 항보다 압도적으로 큼(또는 그 항 자체가 없음) — "
         "그래서 \"n_view(카메라 수)가 시간을 지배한다\"는 결론이 감이 아니라 계수로 확정됨.")

    # 슬라이드 14 — Phase 5~7
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§3 · PHASE 5~7", "회귀분석으로 규명한 관계식, 그리고 품질까지")
    add_image_fit(s, IMG / "fig_exp56_phase67.png", MARGIN, Inches(1.35), BODY_W, Inches(3.4))
    bullet_block(s, MARGIN, Inches(4.95), BODY_W, Inches(1.9), [
        "548개 실제 map() 호출 데이터로 회귀분석 — iters×n_view(반복 횟수×카메라 수)가 시간을 지배, gaussian 수·해상도는 부차적(R²=0.93~0.998)",
        "\"카메라 뷰를 늘리면 품질이 좋아지지 않을까?\" → window(프론티어) 자체를 키우면 오히려 PSNR 악화(−1.1~3.5dB) — 기각. "
        "대신 프론티어는 그대로 두고 과거-keyframe 곁눈질 비중만 늘리기 → PSNR 개선, 시간 비용 무시할 수준 — 채택",
    ], size=12)
    note(s, MARGIN, Inches(6.85), BODY_W, Inches(0.5),
         "같은 \"뷰를 늘린다\"는 아이디어인데 방법에 따라 정반대 결과 — 다음 로드맵(dense-frame supervision)에 대한 경고이기도 함.")

    # 슬라이드 15 — Phase 8
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§3 · PHASE 8", "\"공짜\" 최적화 발견")
    add_image_fit(s, IMG / "fig_exp56_phase8.png", MARGIN, Inches(1.6), BODY_W, Inches(3.2))
    bullet_block(s, MARGIN, Inches(5.1), BODY_W, Inches(1.6), [
        "프로파일링 중 우연히 발견: 카메라 pose가 안 바뀌는데도 매 렌더링 호출마다 행렬 역산을 처음부터 다시 계산하고 있었음",
        "pose가 바뀌는 지점이 코드 전체에 딱 한 곳뿐임을 확인한 뒤 캐싱 추가 — 그래디언트 수학은 전혀 안 건드리는 무위험 변경",
    ], size=12.5)
    note(s, MARGIN, Inches(6.85), BODY_W, Inches(0.5),
         "제일 큰 인프라 작업(다음 슬라이드)을 준비하다가 발견한 제일 작은데 제일 확실한 개선.")

    # 슬라이드 16 — Phase 8b
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§3 · PHASE 8b", "CUDA 커널 batch화 시도 — 정직한 실패 기록")
    add_image_fit(s, IMG / "fig_exp56_phase8b.png", MARGIN, Inches(1.5), BODY_W, Inches(2.4))
    bullet_block(s, MARGIN, Inches(4.15), BODY_W, Inches(2.2), [
        "\"카메라 여러 대를 한 번의 CUDA 커널 호출로 묶으면 더 빠르지 않을까\" — 실제로 구현·검증까지 진행",
        "기존 커널(그래디언트 수학 포함)은 안 건드리고 C++에서 안전하게 감싸는 방식으로 구현, forward는 완전 일치, backward도 수치적으로 검증 통과",
        "실전 적용 1차 시도에서 PSNR 붕괴 발견 → 원인(텐서 shape 불일치) 진단·수정 → 재검증 통과",
        "하지만 최종 시간 측정 결과 개선 없음 — 진짜 병목이 \"CUDA 커널 실행 자체\"라는 사전 진단이 재확인됨, 채택하지 않고 안전하게 원복",
    ], size=11.8)
    note(s, MARGIN, Inches(6.85), BODY_W, Inches(0.55),
         "성공한 것만 보여주면 신뢰가 안 감 — 위험한 시도를 안전하게 검증하고, 안 되면 되돌리는 과정 자체가 방법론.")

    # ── §4 종합 ─────────────────────────────────────────────
    # 슬라이드 17 — 전체 타임라인
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§4 · TIMELINE", "전체 타임라인")
    add_image_fit(s, IMG / "fig_overall_timeline.png", MARGIN, Inches(1.4), Inches(7.6), Inches(4.6))
    tx = MARGIN + Inches(7.9)
    tw = BODY_W - Inches(7.9)
    rows = [
        ["시점", "실시간 배수", "PSNR(mean/kf)"],
        ["07/20 (정정 후)", "1.52배", "—"],
        ["exp53+54", "0.94배", "—"],
        ["exp55", "0.92배", "22.61/22.95"],
        ["exp56 최종", "0.70배", "23.49/23.88"],
    ]
    simple_table(s, tx, Inches(1.6), tw, rows,
                 col_widths=[Inches(1.9), Inches(1.0), Inches(1.3)], row_h=Inches(0.62))
    note(s, MARGIN, Inches(6.35), BODY_W, Inches(0.9),
         "실시간을 달성한 뒤에도 멈추지 않고 추가로 시간을 1/4 줄이면서 품질까지 올렸다.")

    # 슬라이드 18 — 스코어카드
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§4 · SCORECARD", "현재 상태 스코어카드 + 다음 후보")
    stats_y = Inches(1.55)
    stat_pill(s, MARGIN, stats_y, Inches(3.9), Inches(1.15), "온라인 루프 총합\n(예산 65.1s 대비 여유 19s)", "45.79s", TEXT)
    stat_pill(s, MARGIN + Inches(4.1), stats_y, Inches(3.9), Inches(1.15), "PSNR (mean/kf)", "23.49/23.88", GREEN)
    stat_pill(s, MARGIN + Inches(8.2), stats_y, Inches(4.0), Inches(1.15), "실시간 배수", "0.70×", BUDGET)
    add_text(s, MARGIN, Inches(3.0), BODY_W, Inches(0.4), "다음 후보 (우선순위 논의 필요)", 14, BUDGET, bold=True)
    rows = [
        ("① 고위험", CORAL, "CUDA 커널 레벨 batch화(Phase 8b가 실패한 진짜 원인 해결) — 이득 불확실"),
        ("② 저위험", BLUE, "init_itr_num/iters 더 세밀한 스캔 — 이득 작음"),
        ("③ 신규 후보", PURPLE, "오프라인 색정제를 실시간 호환으로 재설계 — 다음 슬라이드, +3~6dB 확인됨"),
        ("④ 다음 단계", GREEN, "Localization(흑백 SLAM 트래킹) 실제 통합 — exp50과 합류"),
    ]
    ay = Inches(3.5)
    for i, (tag, color, text) in enumerate(rows):
        action_row(s, MARGIN, ay + Inches(0.75) * i, BODY_W, tag, color, text)
    note(s, MARGIN, Inches(6.5), BODY_W, Inches(0.9),
         "carve loss로 floater 억제는 유지 중. ①번은 리스크 대비 이득이 불확실해서 팀 논의가 필요한 지점.")

    # 슬라이드 18b — 후처리(색정제)가 사주는 PSNR
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§4 · NEW FINDING", "부록: 오프라인 후처리가 사주는 PSNR")
    add_image_fit(s, IMG / "fig_ply_compare.png", MARGIN, Inches(1.35), BODY_W, Inches(3.5))
    stats_y = Inches(5.0)
    stat_pill(s, MARGIN, stats_y, Inches(3.9), Inches(1.0), "후처리 고정비용\n(26,000 iter, 프레임수 무관)", "~196s", CORAL)
    stat_pill(s, MARGIN + Inches(4.1), stats_y, Inches(3.9), Inches(1.0), "PSNR 이득(mean/kf)", "+3.04 / +6.45dB", GREEN)
    stat_pill(s, MARGIN + Inches(8.2), stats_y, Inches(4.0), Inches(1.0), "exp56 최종 레시피 위 실측", "26.53 / 30.33dB", TEXT)
    note(s, MARGIN, Inches(6.15), BODY_W, Inches(1.1),
         "07/18 구버전 설정에서 처음 발견된 패턴(+4.12/+7.95dB)이 exp56 최종 레시피에서도 재현됨(2026-07-27 재검증). "
         "26k iteration은 시퀀스 길이와 무관한 고정 배치 작업이라 현재 구조로는 \"실시간\"이 아니라 \"끝나고 한 번에 몰아서\" — "
         "다음 슬라이드는 이걸 실시간 호환으로 만들 방법에 대한 제안.")

    # 슬라이드 18c — 실시간 호환 후처리 설계안
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§4 · PROPOSAL", "후처리를 실시간 호환으로 만들려면")
    bullet_block(s, MARGIN, Inches(1.5), BODY_W, Inches(1.0), [
        "핵심 문제: 26,000 iteration이 시퀀스 종료 후 한 번에 블로킹으로 몰림 — 온라인 시스템엔 \"종료 시점\"이 없음",
    ], size=13.5)
    ay = Inches(2.7)
    rows = [
        ("A · 상시 백그라운드화", BLUE, "26k를 세션 끝에 몰아서 돌리지 않고, 이미 지나간(프론티어 밖) keyframe에 대해 GPU 유휴 사이클마다 소량씩 지속 refinement — n_global_views(Phase7)가 이미 하던 걸 전용 백그라운드 루프로 확장"),
        ("B · 우선순위 스케줄링", BUDGET, "_gs_queue가 비어 mapper가 노는 시점(exp55 부록에서 이미 발견한 유휴 구간)을 감지해 그 틈에만 refinement 스텝을 끼워 넣음 — 신규 keyframe 처리를 항상 최우선으로 선점"),
        ("C · 반복수 자체 축소", GREEN, "iters=10→7이 온라인에서 그랬듯, 26,000도 고정 상수일 뿐 — 스윕(예: 5000/10000)해서 수확체감 지점을 찾으면 고정비용 자체를 줄일 수 있음(저위험, 미검증)"),
    ]
    for i, (tag, color, text) in enumerate(rows):
        action_row(s, MARGIN, ay + Inches(1.05) * i, BODY_W, tag, color, text)
    note(s, MARGIN, Inches(6.3), BODY_W, Inches(0.9),
         "권장: C(반복수 스윕)로 저위험 저비용 검증 먼저 → A/B(상시 백그라운드화)로 구조를 바꾸는 건 "
         "exp56 Phase7의 아이디어를 그대로 확장하는 것이라 다음 사이클 후보로 유력.")

    # 슬라이드 18d — post-processed급 online이 왜 어려운가 (현실적 한계)
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§4 · REALITY CHECK", "post-processed급 online — 쉽지 않은 이유")
    add_image_fit(s, IMG / "fig_realtime_gap.png", MARGIN, Inches(1.35), Inches(6.3), Inches(4.2))
    bx = MARGIN + Inches(6.6); bw = BODY_W - Inches(6.6)
    bullet_block(s, bx, Inches(1.5), bw, Inches(4.2), [
        "계산량 격차: 26,000 iter × ~6.4ms/view-op(Phase8 실측치 그대로 적용) ≈ 166~196초 필요 — 지금 파이프라인 여유는 65.1s 예산 대비 19.3초뿐, 약 8~10배 부족",
        "더 근본적 문제: color_refinement는 \"전체 궤적을 이미 다 본 상태\"에서 모든 keyframe을 반복 방문하는 global 작업 — 라이브 스트림엔 애초에 \"다 봤다\"는 시점이 없음",
        "즉 \"매 순간 완전히 캐치업된 post-processed급\"은 계산량 문제 이전에 정의상 불가능",
    ], size=12.5)
    note(s, bx, Inches(6.1), bw, Inches(1.0),
         "그래서 목표를 \"항상 최종 품질\"이 아니라 \"프론티어 + 뒤에서 점점 따라잡는 품질\"로 재정의하는 게 "
         "현실적 — 다음 슬라이드가 그 그림.")

    # 슬라이드 18e — bounded-lag 목표 재정의
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§4 · REALITY CHECK", "현실적 목표: bounded-lag 온라인 정제")
    add_image_fit(s, IMG / "fig_bounded_lag.png", MARGIN, Inches(1.4), BODY_W, Inches(3.9))
    note(s, MARGIN, Inches(5.6), BODY_W, Inches(1.5),
         "루프클로저 global BA를 백그라운드로 돌리는 성숙한 SLAM 시스템들과 같은 패턴 — 프론티어(방금 본 곳)는 "
         "항상 실시간으로 그려지고, 그 뒤를 정제(refinement)가 시간차를 두고 계속 따라오며 품질을 끌어올린다. "
         "\"실시간 = 항상 완성본\"이 아니라 \"실시간 = 프론티어가 안 밀린다 + 지나간 곳은 계속 좋아진다\"로 "
         "성공 기준 자체를 바꾸는 게 이번 발견의 결론.")

    # 슬라이드 19 — 결론
    s = blank(prs); set_bg(s)
    eyebrow_title(s, "§4 · CONCLUSION", "결론 한 줄")
    ay = Inches(1.7)
    rows = [
        ("EXP53+54", BLUE, "실시간 최초 돌파 (0.94배)"),
        ("EXP55", GREEN, "품질 레버 확보 (gaussian −35.9%, floater −7.5%, 비용 없음)"),
        ("EXP56", BUDGET, "\"왜 안 빨라지나\"를 끝까지 규명해 4단계 연속 개선 + 실패한 시도(batch화)까지 안전하게 검증"),
    ]
    for i, (tag, color, text) in enumerate(rows):
        action_row(s, MARGIN, ay + Inches(0.85) * i, BODY_W, tag, color, text)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(4.7), BODY_W, Inches(1.3))
    box.adjustments[0] = 0.06
    box.fill.solid(); box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = BUDGET; box.line.width = Pt(1.5)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "실시간을 달성하고, 그 위에서 시간을 추가로 1/4 더 줄이면서 품질까지 함께 올렸다"
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = TEXT; r.font.name = FONT
    note(s, MARGIN, Inches(6.4), BODY_W, Inches(0.7),
         "다음 미팅까지 목표(localization 통합 vs 커널 batch화 여부)에 대한 팀 의견 요청.")

    prs.save(OUT)
    print(f"[saved] {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
